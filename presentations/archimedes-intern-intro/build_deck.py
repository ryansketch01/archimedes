"""Archimedes — intern orientation deck. 10 slides, ~5-7 min.

Audience: college intern with no CTI background, no Claude background,
basic programming familiarity. The goal is to ORIENT them to what
we built, not train them to use or rebuild it.

Design intent:
  - Plain English. No jargon (or jargon defined on first use).
  - Analogies they'd know (news aggregator, ChatGPT, Discord bot).
  - No Admiralty / WEP / MITRE / Hard Rule N references.
  - One idea per slide. ~30-45 sec each.
  - Same stylistic neutrality as other decks (Calibri, gray/blue
    accents, no custom backgrounds) for template-merge.

Run:
  uv run python presentations/archimedes-intern-intro/build_deck.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------- style ----------

C_TITLE = RGBColor(0x1F, 0x3A, 0x5F)
C_ACCENT = RGBColor(0x4A, 0x90, 0xE2)
C_BODY = RGBColor(0x33, 0x33, 0x33)
C_SUBTLE = RGBColor(0x66, 0x66, 0x66)
C_BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
C_DISCORD_BG = RGBColor(0x31, 0x33, 0x38)
C_DISCORD_TEXT = RGBColor(0xDC, 0xDD, 0xDE)
C_KNOWN = RGBColor(0x2E, 0x7D, 0x32)


# ---------- helpers ----------


def add_title(slide, text: str, *, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = C_TITLE
    r.font.name = "Calibri"

    if subtitle:
        b2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.33), Inches(0.5))
        tf2 = b2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(16)
        r2.font.italic = True
        r2.font.color.rgb = C_SUBTLE
        r2.font.name = "Calibri"


def add_bullets(
    slide,
    bullets: list[str],
    *,
    left: float = 0.5,
    top: float = 1.9,
    width: float = 12.33,
    height: float = 5.0,
    body_size: int = 20,
    line_spacing: float = 1.3,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if idx > 0:
            p.space_before = Pt(14)
        r = p.add_run()
        r.text = f"•  {item}"
        r.font.size = Pt(body_size)
        r.font.color.rgb = C_BODY
        r.font.name = "Calibri"


def add_notes(slide, notes: str) -> None:
    slide.notes_slide.notes_text_frame.text = notes


def blank_slide(prs: Presentation):
    return prs.slide_layouts[6]
    # actually return the slide, not the layout
def blank_slide(prs: Presentation):  # noqa: F811
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_pagenum(slide, n: int, total: int) -> None:
    box = slide.shapes.add_textbox(Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{n} / {total}"
    r.font.size = Pt(9)
    r.font.color.rgb = C_SUBTLE
    r.font.name = "Calibri"


def add_footer(slide, text: str = "Archimedes · Quick intro") -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(11.0), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(9)
    r.font.color.rgb = C_SUBTLE
    r.font.italic = True
    r.font.name = "Calibri"


# ---------- build ----------


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    TOTAL = 10

    # =================================================================
    # Slide 1 — Title
    # =================================================================
    s = blank_slide(prs)

    box = s.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(12.33), Inches(1.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Meet Archimedes"
    r.font.size = Pt(64)
    r.font.bold = True
    r.font.color.rgb = C_TITLE
    r.font.name = "Calibri"

    box2 = s.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.33), Inches(0.6))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "An AI agent that reads cyber threat news so we don't have to"
    r2.font.size = Pt(22)
    r2.font.italic = True
    r2.font.color.rgb = C_SUBTLE
    r2.font.name = "Calibri"

    box3 = s.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.33), Inches(0.5))
    tf3 = box3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = "5-minute intro   ·   [Date]   ·   [Presenter name]"
    r3.font.size = Pt(14)
    r3.font.color.rgb = C_SUBTLE
    r3.font.name = "Calibri"

    add_notes(
        s,
        "Set the scene casually: 'I want to spend about five minutes walking you through what we built, why it exists, and what you'll see it doing day to day.' "
        "Tell them up front: no jargon test, ask any question that comes up, the goal is for you to understand the system end-to-end at a high level."
    )

    # =================================================================
    # Slide 2 — The problem
    # =================================================================
    s = blank_slide(prs)
    add_title(s, "The problem we're solving",
              subtitle="Every morning, security teams have a giant pile of news to read.")

    add_bullets(
        s,
        [
            "Hundreds of new articles every day — vendor blogs, news sites, X/Twitter, CISA alerts, RSS feeds",
            "Most of it isn't relevant. Some of it is critical and time-sensitive.",
            "Sorting the firehose by hand takes hours every morning before any actual analysis starts",
            "Leadership wants a 5-minute summary by 8 AM. Analysts have to deliver that on top of everything else.",
        ],
        body_size=18,
    )

    add_notes(
        s,
        "Frame this like a relatable workday problem. 'Imagine you're a security analyst — your inbox has 200 articles from overnight. Most are noise. Some are emergencies. Your boss needs a 5-minute summary in your hand by 8 AM. That's the job we're trying to make easier.' "
        "Don't dwell on the technical details of CTI — just get them to feel the volume problem."
    )
    add_footer(s)
    add_pagenum(s, 2, TOTAL)

    # =================================================================
    # Slide 3 — What we built
    # =================================================================
    s = blank_slide(prs)
    add_title(s, "What we built",
              subtitle="An AI agent that does the rote 80% so humans can focus on the judgment 20%.")

    # Big centered statement
    statement = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.0), Inches(11.33), Inches(2.8)
    )
    statement.line.color.rgb = C_ACCENT
    statement.line.width = Pt(2)
    statement.fill.solid()
    statement.fill.fore_color.rgb = C_BG_LIGHT
    tf = statement.text_frame
    tf.margin_top = Inches(0.4)
    tf.margin_left = Inches(0.4)
    tf.margin_right = Inches(0.4)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Archimedes reads ~45 cyber news sources every day,"
    r.font.size = Pt(26)
    r.font.color.rgb = C_BODY
    r.font.name = "Calibri"

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)
    r2 = p2.add_run()
    r2.text = "decides what's important,"
    r2.font.size = Pt(26)
    r2.font.color.rgb = C_BODY
    r2.font.name = "Calibri"

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(10)
    r3 = p3.add_run()
    r3.text = "and posts a 5-minute summary to Discord twice a day."
    r3.font.size = Pt(26)
    r3.font.bold = True
    r3.font.color.rgb = C_TITLE
    r3.font.name = "Calibri"

    # Sub-line
    sb = s.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12.33), Inches(0.5))
    stf = sb.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run()
    sr.text = "It also sends urgent alerts when something breaks — and you can ask it questions in chat."
    sr.font.size = Pt(16)
    sr.font.italic = True
    sr.font.color.rgb = C_SUBTLE
    sr.font.name = "Calibri"

    add_notes(
        s,
        "One sentence. Read it out loud, slowly. 'It reads 45 cyber news sources every day, decides what's important, and posts a summary to Discord twice a day.' "
        "If they remember nothing else from this talk, they should remember that sentence. "
        "Mention the urgent alerts (FLASH) and the chat interaction (slash commands) briefly — those come up later in the deck."
    )
    add_footer(s)
    add_pagenum(s, 3, TOTAL)

    # =================================================================
    # Slide 4 — What it looks like (Discord example)
    # =================================================================
    s = blank_slide(prs)
    add_title(s, "What it looks like",
              subtitle="This is the actual morning summary that posts to our Discord channel every day at 8 AM.")

    # Discord-style mockup
    discord_text = """Good morning. Here's your 0800 brief — Monday, May 11.

🚨 Active Threats

•  Checkmarx Jenkins AST plugin compromised in supply-chain attack
    Checkmarx warned Friday May 9 of a malicious Jenkins AST plugin on the
    Marketplace; weekend variants surfaced on GitHub. Remediation 2.0.13-848
    is out.

•  SailPoint discloses GitHub repository hack
    SailPoint disclosed an April 20 incident; a third-party app vulnerability
    exposed a subset of its GitHub repos.

🔓 Vulnerabilities

•  CVE-2026-6973 (Ivanti EPMM): federal patch deadline expired last night.
•  CVE-2026-42208 (BerriAI LiteLLM): KEV deadline closes today."""

    bg = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.9),
        Inches(12.33), Inches(5.0),
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_DISCORD_BG

    tb = s.shapes.add_textbox(
        Inches(0.8), Inches(2.1),
        Inches(11.7), Inches(4.6),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(discord_text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line if line else " "
        r.font.size = Pt(13)
        r.font.color.rgb = C_DISCORD_TEXT
        r.font.name = "Consolas"

    add_notes(
        s,
        "Show them a real artifact. This is what shipped to #intel-briefs on May 11. Three things to point out: "
        "(1) Plain language — not jargon-y. 'Patches expired last night.' Anyone can read it. "
        "(2) Source links — the headlines are clickable; click and you get the actual article. "
        "(3) Short — ~250 words; whole thing fits on a phone screen. Leadership reads this on the way to a meeting."
    )
    add_footer(s)
    add_pagenum(s, 4, TOTAL)

    # =================================================================
    # Slide 5 — How it works: Collecting
    # =================================================================
    s = blank_slide(prs)
    add_title(s, "How it works — Step 1: Collecting",
              subtitle="Think of it as a news aggregator on steroids.")

    add_bullets(
        s,
        [
            "~45 sources monitored — CISA alerts, Mandiant blog, BleepingComputer, RSS feeds, threat-intel APIs",
            "Pulls everything new in the last few hours",
            "Filters out the noise — only keeps items relevant to what we care about (defense industry, tracked threat groups, sector-specific incidents)",
            "Result: maybe 4-8 actually relevant items out of 200+ raw articles",
        ],
        body_size=18,
    )

    add_notes(
        s,
        "Analogy: this is like Google News, but specifically for cyber threats, and it filters way more aggressively. "
        "Out of hundreds of articles published overnight, maybe 4-8 are worth telling leadership about. The collector's job is to find those 4-8."
    )
    add_footer(s)
    add_pagenum(s, 5, TOTAL)

    # =================================================================
    # Slide 6 — How it works: Grading
    # =================================================================
    s = blank_slide(prs)
    add_title(s, "How it works — Step 2: Grading",
              subtitle="Not all news is equal. The agent scores how trustworthy each source is, every time.")

    add_bullets(
        s,
        [
            "Mandiant publishing on a Chinese threat group: very reliable. The agent trusts it.",
            "An anonymous Reddit post claiming the same thing: not reliable. The agent flags it as a rumor.",
            "If only ONE source says something, the agent caps how confident it sounds — even if the source is great",
            "If TWO independent sources agree, confidence goes up",
            "Same rules applied every time — no one analyst's mood changing the scoring",
        ],
        body_size=17,
    )

    add_notes(
        s,
        "This is the part most people are surprised by. The agent doesn't just summarize articles — it grades how much to trust them, using a framework that intelligence analysts have used for decades. "
        "Big point: consistency. Two human analysts grading the same finding might disagree based on experience or mood. The agent applies the same rule every time. "
        "Skip the words 'Admiralty' and 'WEP' — those are jargon for the intern."
    )
    add_footer(s)
    add_pagenum(s, 6, TOTAL)

    # =================================================================
    # Slide 7 — How it works: Briefing
    # =================================================================
    s = blank_slide(prs)
    add_title(s, "How it works — Step 3: Writing the brief",
              subtitle="The agent writes TWO versions of every brief — one for the deep-readers, one for everyone else.")

    # Two columns
    left_box = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(2.0),
        Inches(6.0), Inches(4.5),
    )
    left_box.line.color.rgb = C_TITLE
    left_box.line.width = Pt(2)
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    ltf = left_box.text_frame
    ltf.margin_top = Inches(0.3)
    ltf.margin_left = Inches(0.3)
    ltf.margin_right = Inches(0.3)
    ltf.word_wrap = True
    lp = ltf.paragraphs[0]
    lp.alignment = PP_ALIGN.LEFT
    lr = lp.add_run()
    lr.text = "Long version (~700 words)"
    lr.font.size = Pt(20)
    lr.font.bold = True
    lr.font.color.rgb = C_TITLE
    lr.font.name = "Calibri"

    lp2 = ltf.add_paragraph()
    lp2.space_before = Pt(14)
    lr2 = lp2.add_run()
    lr2.text = (
        "For the intel team to review carefully.\n\n"
        "Includes:\n"
        "•  All the source citations\n"
        "•  Confidence ratings per item\n"
        "•  Cross-references to other findings\n"
        "•  Caveats and uncertainty notes\n\n"
        "Lives in our records (git) forever."
    )
    lr2.font.size = Pt(14)
    lr2.font.color.rgb = C_BODY
    lr2.font.name = "Calibri"

    right_box = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.83), Inches(2.0),
        Inches(6.0), Inches(4.5),
    )
    right_box.line.color.rgb = C_ACCENT
    right_box.line.width = Pt(2)
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    rtf = right_box.text_frame
    rtf.margin_top = Inches(0.3)
    rtf.margin_left = Inches(0.3)
    rtf.margin_right = Inches(0.3)
    rtf.word_wrap = True
    rp = rtf.paragraphs[0]
    rp.alignment = PP_ALIGN.LEFT
    rr = rp.add_run()
    rr.text = "Short version (~250 words)"
    rr.font.size = Pt(20)
    rr.font.bold = True
    rr.font.color.rgb = C_ACCENT
    rr.font.name = "Calibri"

    rp2 = rtf.add_paragraph()
    rp2.space_before = Pt(14)
    rr2 = rp2.add_run()
    rr2.text = (
        "For everyone else.\n\n"
        "Includes:\n"
        "•  The headlines and what they mean\n"
        "•  Clickable links to the source\n"
        "•  Dates in normal English\n"
        "•  Action items where they apply\n\n"
        "This is what goes to Discord."
    )
    rr2.font.size = Pt(14)
    rr2.font.color.rgb = C_BODY
    rr2.font.name = "Calibri"

    add_notes(
        s,
        "This is a key design choice: same intel, two audiences, two versions written from the same source. "
        "The intel team needs the detailed version with all the citations. Leadership wants the 5-minute version on their phone. "
        "The agent writes both at the same time — one file, two sections. The Discord post you saw on the previous slide is the short version."
    )
    add_footer(s)
    add_pagenum(s, 7, TOTAL)

    # =================================================================
    # Slide 8 — Safety rails
    # =================================================================
    s = blank_slide(prs)
    add_title(s, "Safety rails",
              subtitle="The agent has rules it can't break — even if you ask it to.")

    add_bullets(
        s,
        [
            "Will not scan random websites or other companies (only systems we own and authorized)",
            "Will not make up where an attack came from — it only repeats what other sources have said, and cites them",
            "Will not write malicious code, even for 'educational' purposes",
            "Will not store passwords or credentials, ever",
            "For high-stakes calls (like upgrading a threat group to HIGH priority), it asks a human to approve first",
        ],
        body_size=17,
    )

    add_notes(
        s,
        "Important framing for an intern: AI agents can do harmful things if not constrained. Archimedes has rules baked into its code — not just instructions in a prompt, but actual code that refuses to do certain things. "
        "Example: if you tried to make it scan, say, microsoft.com to look for vulnerabilities, the agent would refuse — and log the refusal. The system is built to be safe by default, not safe-if-instructed."
    )
    add_footer(s)
    add_pagenum(s, 8, TOTAL)

    # =================================================================
    # Slide 9 — You can talk to it
    # =================================================================
    s = blank_slide(prs)
    add_title(s, "You can talk to it",
              subtitle="Type a slash command in Discord; the agent investigates and writes back.")

    # Show example slash commands
    commands_box = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.9),
        Inches(12.33), Inches(4.5),
    )
    commands_box.line.color.rgb = C_ACCENT
    commands_box.line.width = Pt(1.5)
    commands_box.fill.solid()
    commands_box.fill.fore_color.rgb = C_BG_LIGHT

    tf = commands_box.text_frame
    tf.margin_top = Inches(0.3)
    tf.margin_left = Inches(0.4)
    tf.margin_right = Inches(0.4)
    tf.word_wrap = True

    examples = [
        ("/cve CVE-2026-0300",        "→ Get a detailed report on that specific vulnerability"),
        ("/investigate APT37",        "→ Get a deep dive on a threat group: what they do, who they target, what to watch for"),
        ("/ioc-hunt 1.2.3.4",         "→ Check whether an IP address has shown up anywhere in our systems"),
        ("/help",                     "→ List everything you can ask"),
    ]
    first = True
    for cmd, desc in examples:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(14)
        r = p.add_run()
        r.text = cmd
        r.font.size = Pt(17)
        r.font.bold = True
        r.font.color.rgb = C_ACCENT
        r.font.name = "Consolas"

        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        r2 = p2.add_run()
        r2.text = "      " + desc
        r2.font.size = Pt(15)
        r2.font.color.rgb = C_BODY
        r2.font.name = "Calibri"

    # Note below
    note = s.shapes.add_textbox(Inches(0.5), Inches(6.55), Inches(12.33), Inches(0.5))
    ntf = note.text_frame
    ntf.word_wrap = True
    np_ = ntf.paragraphs[0]
    np_.alignment = PP_ALIGN.CENTER
    nr = np_.add_run()
    nr.text = "Same agent as the brief. Same sources. Same grading rules. Just on-demand."
    nr.font.size = Pt(13)
    nr.font.italic = True
    nr.font.color.rgb = C_SUBTLE
    nr.font.name = "Calibri"

    add_notes(
        s,
        "Think of it like ChatGPT, but specifically tuned for our threat intel. "
        "If you're curious about something during the day — 'what's the deal with this CVE?' — you don't have to wait for the next morning brief. You can ask. "
        "The agent pulls from the same sources it uses for the briefs, applies the same grading rules, and writes back in the channel."
    )
    add_footer(s)
    add_pagenum(s, 9, TOTAL)

    # =================================================================
    # Slide 10 — Wrap-up + Q&A
    # =================================================================
    s = blank_slide(prs)

    box = s.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.33), Inches(0.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "That's it"
    r.font.size = Pt(48)
    r.font.bold = True
    r.font.color.rgb = C_TITLE
    r.font.name = "Calibri"

    # Three takeaways
    takeaways = [
        ("Twice-daily summary",
         "Morning + afternoon briefs in #intel-briefs. Five minutes of reading covers the day's biggest threats."),
        ("Built to be safe",
         "Safety rails baked into the code. Humans approve the highest-stakes calls."),
        ("Ask it questions",
         "Slash commands in Discord. Same brain as the briefs, on-demand."),
    ]
    tb_w = Inches(4.0)
    tb_h = Inches(2.5)
    tb_top = Inches(2.2)
    tb_gap = Inches(0.15)
    tb_total_w = tb_w * 3 + tb_gap * 2
    tb_start_x = (Inches(13.333) - tb_total_w) / 2
    for i, (title, body) in enumerate(takeaways):
        x = tb_start_x + (tb_w + tb_gap) * i
        bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, tb_top, tb_w, tb_h)
        bx.line.color.rgb = C_ACCENT
        bx.line.width = Pt(1.5)
        bx.fill.solid()
        bx.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        ttf = bx.text_frame
        ttf.margin_top = Inches(0.3)
        ttf.margin_left = Inches(0.25)
        ttf.margin_right = Inches(0.25)
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.alignment = PP_ALIGN.CENTER
        tr = tp.add_run()
        tr.text = title
        tr.font.size = Pt(18)
        tr.font.bold = True
        tr.font.color.rgb = C_TITLE
        tr.font.name = "Calibri"

        tp2 = ttf.add_paragraph()
        tp2.alignment = PP_ALIGN.CENTER
        tp2.space_before = Pt(14)
        tr2 = tp2.add_run()
        tr2.text = body
        tr2.font.size = Pt(13)
        tr2.font.color.rgb = C_BODY
        tr2.font.name = "Calibri"

    qa_box = s.shapes.add_textbox(Inches(0.5), Inches(5.4), Inches(12.33), Inches(1.2))
    qtf = qa_box.text_frame
    qtf.word_wrap = True
    qp = qtf.paragraphs[0]
    qp.alignment = PP_ALIGN.CENTER
    qr = qp.add_run()
    qr.text = "Questions?"
    qr.font.size = Pt(36)
    qr.font.bold = True
    qr.font.color.rgb = C_TITLE
    qr.font.name = "Calibri"

    qp2 = qtf.add_paragraph()
    qp2.alignment = PP_ALIGN.CENTER
    qp2.space_before = Pt(10)
    qr2 = qp2.add_run()
    qr2.text = "Curious about how anything works under the hood? Happy to dig in."
    qr2.font.size = Pt(14)
    qr2.font.italic = True
    qr2.font.color.rgb = C_SUBTLE
    qr2.font.name = "Calibri"

    add_notes(
        s,
        "Three quick takeaways, then Q&A. The three boxes are the system's elevator pitch — twice-daily summary, safety rails, ask it questions. "
        "Likely intern questions: "
        "(1) 'Does it ever get things wrong?' — Yes, and we have a process for that — retractions are public, never silent. "
        "(2) 'Can I see the code?' — Yes, it's all in our git repo; show them where if appropriate. "
        "(3) 'Why Discord instead of Teams / Slack?' — Convenience and existing channel structure; not a technical reason. "
        "(4) 'What sources does it read?' — Show source-grades.yaml if helpful. "
        "(5) 'Did you build this in Python?' — Mostly markdown + bash + Python for the integrations. The 'brain' is Claude. "
        "End on: 'There's a more technical version of this talk if you ever want to learn how it's built — just ask.'"
    )
    add_footer(s)
    add_pagenum(s, 10, TOTAL)

    # ---------- save ----------
    out_path = Path(__file__).parent / "archimedes-intern-intro.pptx"
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    path = build()
    print(f"Wrote: {path}")
    print(f"Size:  {path.stat().st_size:,} bytes")
