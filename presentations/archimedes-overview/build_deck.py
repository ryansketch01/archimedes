"""Archimedes overview deck — 30-min, 24-slide presentation for a Claude-
familiar / CTI-new software dev team.

Output: archimedes-overview.pptx in the same directory as this script.

Design intent (per session-13 build conversation):
  - Stylistically neutral so it template-merges cleanly into the operator's
    company deck (Calibri, simple layouts, gray/dark-blue accents only)
  - Real content from the Archimedes corpus where it lands; mock-rendered
    Discord brief block where a literal screenshot would be hard to embed
  - Speaker notes 2-3 sentences per slide
  - 24 slides total; runs ~75s/slide = 30 minutes
  - Architecture diagram built native in pptx shapes (no external image)
  - Audience: Claude-familiar devs, new to CTI tradecraft — scaffold the
    domain before showing the Claude-shaped solution

Run:
  uv run python presentations/archimedes-overview/build_deck.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------- style constants (neutral; safe for template-merge) ----------

C_TITLE = RGBColor(0x1F, 0x3A, 0x5F)        # dark blue
C_ACCENT = RGBColor(0x4A, 0x90, 0xE2)       # mid blue
C_BODY = RGBColor(0x33, 0x33, 0x33)         # near-black
C_SUBTLE = RGBColor(0x66, 0x66, 0x66)       # mid-gray
C_BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)     # light gray for code/quote blocks
C_DISCORD_BG = RGBColor(0x31, 0x33, 0x38)   # discord dark bg
C_DISCORD_TEXT = RGBColor(0xDC, 0xDD, 0xDE) # discord text
C_DISCORD_LINK = RGBColor(0x00, 0xAF, 0xF4) # discord blue link
C_REFUSAL = RGBColor(0xC0, 0x39, 0x2B)      # muted red for "refused" callouts


# ---------- helpers ----------


def add_title(slide, text: str, *, subtitle: str | None = None) -> None:
    """Add a top-banded title to a slide using a textbox (avoids placeholder
    quirks across templates)."""
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(12.33)
    height = Inches(0.7)

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = C_TITLE
    run.font.name = "Calibri"

    if subtitle:
        box2 = slide.shapes.add_textbox(left, Inches(1.0), width, Inches(0.4))
        tf2 = box2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(14)
        run2.font.italic = True
        run2.font.color.rgb = C_SUBTLE
        run2.font.name = "Calibri"


def add_bullets(
    slide,
    bullets: list[str | tuple[str, str]],
    *,
    left: float = 0.5,
    top: float = 1.7,
    width: float = 12.33,
    height: float = 5.0,
    body_size: int = 18,
    sub_size: int = 14,
    line_spacing: float = 1.15,
) -> None:
    """Add a bullet list to a slide. Each entry can be a string (single line)
    or a (headline, sub) tuple (headline bold, sub regular underneath)."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True

    for idx, item in enumerate(bullets):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if idx > 0:
            p.space_before = Pt(8)

        if isinstance(item, str):
            run = p.add_run()
            run.text = f"• {item}"
            run.font.size = Pt(body_size)
            run.font.color.rgb = C_BODY
            run.font.name = "Calibri"
        else:
            headline, sub = item
            run = p.add_run()
            run.text = f"• {headline}"
            run.font.size = Pt(body_size)
            run.font.bold = True
            run.font.color.rgb = C_BODY
            run.font.name = "Calibri"

            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.LEFT
            p2.line_spacing = line_spacing
            run2 = p2.add_run()
            run2.text = f"   {sub}"
            run2.font.size = Pt(sub_size)
            run2.font.color.rgb = C_SUBTLE
            run2.font.italic = True
            run2.font.name = "Calibri"


def add_code_block(
    slide,
    code: str,
    *,
    left: float = 0.5,
    top: float = 1.7,
    width: float = 12.33,
    height: float = 5.0,
    font_size: int = 11,
    discord: bool = False,
) -> None:
    """Add a monospace block to mimic terminal/code/yaml/Discord output."""
    # Background rectangle
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    bg.line.fill.background()
    if discord:
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_DISCORD_BG
        text_color = C_DISCORD_TEXT
    else:
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_BG_LIGHT
        text_color = C_BODY

    # Text on top
    tb = slide.shapes.add_textbox(
        Inches(left + 0.15),
        Inches(top + 0.1),
        Inches(width - 0.3),
        Inches(height - 0.2),
    )
    tf = tb.text_frame
    tf.word_wrap = True

    lines = code.split("\n")
    for idx, line in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line if line else " "
        run.font.size = Pt(font_size)
        run.font.color.rgb = text_color
        run.font.name = "Consolas"


def add_notes(slide, notes: str) -> None:
    """Speaker notes."""
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = notes


def blank_slide(prs: Presentation):
    """Use the blank layout (index 6) — no template-side placeholders to
    fight with."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_footer(slide, text: str = "Archimedes overview · 2026-05-11") -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.color.rgb = C_SUBTLE
    run.font.italic = True
    run.font.name = "Calibri"


def add_pagenum(slide, n: int, total: int) -> None:
    box = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.3), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{n} / {total}"
    run.font.size = Pt(9)
    run.font.color.rgb = C_SUBTLE
    run.font.name = "Calibri"


# ---------- build ----------


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    TOTAL = 24

    # ===================================================================
    # Slide 1 — Title
    # ===================================================================
    s = blank_slide(prs)
    # Big centered title
    box = s.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(12.33), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Building a CTI Analyst with Claude"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = C_TITLE
    run.font.name = "Calibri"

    box2 = s.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12.33), Inches(0.6))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = (
        "Lessons from Archimedes — an autonomous threat-intel agent shipping "
        "graded briefs twice a day"
    )
    run2.font.size = Pt(18)
    run2.font.italic = True
    run2.font.color.rgb = C_SUBTLE
    run2.font.name = "Calibri"

    box3 = s.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.33), Inches(0.5))
    tf3 = box3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "[Presenter name]   ·   [Date]   ·   For: [Software dev team / OpenCTI build]"
    run3.font.size = Pt(14)
    run3.font.color.rgb = C_SUBTLE
    run3.font.name = "Calibri"

    add_notes(
        s,
        "Open by naming the audience — Claude-familiar devs building something similar on OpenCTI. "
        "Frame: 30 minutes, no live demo, lots of real artifacts from a working system. "
        "Goal of the talk is to make the case for WHY this kind of tool needs to exist, then show that Claude makes it feasible."
    )

    # ===================================================================
    # Slide 2 — The CTI analyst's day
    # ===================================================================
    s = blank_slide(prs)
    add_title(s, "The CTI analyst's day", subtitle="What an intel analyst is actually doing at 0700")
    add_bullets(
        s,
        [
            ("~45 OSINT sources to monitor",
             "CISA KEV, vendor blogs (Mandiant, Unit 42, MSTIC, CrowdStrike), security trades "
             "(BleepingComputer, The Record, Krebs), X/RSS firehose, threat-feed APIs"),
            ("Hundreds of items per day",
             "Most are noise. A handful are signal. Sorting one from the other is the job."),
            ("0800 deadline for the morning brief",
             "Leadership wants a 5-minute readable summary — graded, prioritized, sector-relevant — "
             "by 0800 EDT. Not raw signal. Not links to articles."),
            ("Then again at 1600. And FLASH if something breaks.",
             "Plus weekly synthesis, threat-actor deep dives, retraction handling when something "
             "ships wrong. The cadence never stops."),
            ("Burnout is the norm",
             "One human, one inbox, one 'too much signal, too little time' problem. The grading and "
             "delivery layers eat the time the analyst should spend on judgment."),
        ],
        body_size=18,
        sub_size=13,
    )
    add_notes(
        s,
        "Spend most of this slide on the volume problem. Devs in the audience can map this to alert fatigue in SOC work — same shape, broader source set. "
        "The handoff next slide is: leadership doesn't see this volume; they see the analyst's output. The gap between firehose and decision-maker is where the system has to live."
    )
    add_footer(s)
    add_pagenum(s, 2, TOTAL)

    # ===================================================================
    # Slide 3 — Where leadership sits
    # ===================================================================
    s = blank_slide(prs)
    add_title(s, "Where leadership sits", subtitle="Different audience, different document")
    add_bullets(
        s,
        [
            ("Reads on a phone, on the way to a meeting",
             "Not at a desk. Not with attachments. Scan-on-mobile is the constraint."),
            ("Wants the call, not the evidence",
             "\"Should we be worried about X?\" — they want yes/no/watch with one sentence of why."),
            ("Trusts the brief or doesn't read it",
             "If two briefs in a row land wrong, you've lost them. Trust is binary."),
            ("Doesn't speak Admiralty, WEP, or MITRE IDs",
             "The framework is for the analyst's rigor — the leadership view strips it back to plain English."),
            ("Wants source links — for the reader who wants to dig",
             "Trust-but-verify. Headlines that link to the source article are non-negotiable."),
        ],
        body_size=18,
        sub_size=13,
    )
    add_notes(
        s,
        "The audience here might assume 'leadership' means CISO — broaden it: program managers, IR leads, contract officers in regulated industries also consume these. "
        "Point being: one canonical analytic record, two audiences. The product has to serve both without compromising either. That's the design problem the rest of the talk addresses."
    )
    add_footer(s)
    add_pagenum(s, 3, TOTAL)

    # ===================================================================
    # Slide 4 — The three jobs
    # ===================================================================
    s = blank_slide(prs)
    add_title(s, "The three jobs", subtitle="Collect → Grade → Deliver. Each has a human bottleneck.")

    # Three column layout — boxes side by side
    col_width = Inches(4.0)
    col_height = Inches(4.5)
    col_top = Inches(1.8)
    col_lefts = [Inches(0.5), Inches(4.7), Inches(8.9)]
    col_titles = ["COLLECT", "GRADE", "DELIVER"]
    col_bodies = [
        "Read everything that hit your sources in the last 14 hours.\n\n"
        "Decide what's worth pulling forward to the brief.\n\n"
        "Pull IOCs, actor mentions, CVE numbers, sector tags. Don't lose context.\n\n"
        "Cross-reference against tracked actors and prior coverage.",

        "Apply the Admiralty 6×6 to each source.\n\n"
        "Apply Words of Estimative Probability to every forward claim.\n\n"
        "Single-source veto on 'very likely' or higher.\n\n"
        "Refuse to originate attribution.\n\n"
        "Test alternative hypotheses on load-bearing claims.",

        "Compose the analyst-grade markdown — full doctrine framings.\n\n"
        "Compose the Smart Brevity Discord summary — mobile-readable.\n\n"
        "Post to the channel, commit to git, log the audit event.\n\n"
        "Handle retractions additively. Never silently delete.",
    ]
    col_colors = [C_ACCENT, C_TITLE, C_REFUSAL]

    for left, ctitle, body, color in zip(col_lefts, col_titles, col_bodies, col_colors):
        # Header band
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, col_top, col_width, Inches(0.5))
        hdr.line.fill.background()
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr_tf = hdr.text_frame
        hdr_tf.margin_top = Inches(0.05)
        hdr_p = hdr_tf.paragraphs[0]
        hdr_p.alignment = PP_ALIGN.CENTER
        hdr_run = hdr_p.add_run()
        hdr_run.text = ctitle
        hdr_run.font.size = Pt(18)
        hdr_run.font.bold = True
        hdr_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        hdr_run.font.name = "Calibri"

        # Body
        body_box = s.shapes.add_textbox(
            left, Inches(col_top.inches + 0.6), col_width, col_height
        )
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        first = True
        for line in body.split("\n"):
            if first:
                p = body_tf.paragraphs[0]
                first = False
            else:
                p = body_tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.2
            run = p.add_run()
            run.text = line
            run.font.size = Pt(13)
            run.font.color.rgb = C_BODY
            run.font.name = "Calibri"

    add_notes(
        s,
        "Each column corresponds to a layer of the agent architecture you'll see on slide 14. "
        "Critical framing: humans CAN do all three of these. They just can't do them at the volume and cadence leadership demands without burning out. "
        "Claude doesn't replace the analyst — it does the rote 80% so the human can spend time on the 20% only humans should do (attribution calls, retraction sign-off, doctrine updates)."
    )
    add_footer(s)
    add_pagenum(s, 4, TOTAL)

    # ===================================================================
    # Slide 5 — Why Claude is the right tool
    # ===================================================================
    s = blank_slide(prs)
    add_title(s, "Why Claude is the right tool", subtitle="Three concrete asks Claude is good at")
    add_bullets(
        s,
        [
            ("Read 200 articles, surface 4 worth grading",
             "Long-context, summarization-at-volume. Cheap relative to an analyst-hour."),
            ("Apply structured frameworks consistently",
             "Admiralty 6×6, WEP probability bands, MITRE ATT&CK tagging. Frameworks are rules; "
             "rules are what LLMs follow well when the doctrine is written precisely."),
            ("Compose to a strict format",
             "Smart Brevity — banned phrases, active voice, lead with impact, bold the 'what.' "
             "Format compliance is checklist-shaped; the agent regenerates if it fails preflight."),
            ("Tool use for collection + delivery",
             "MCP wrappers expose Shodan / VirusTotal / theHarvester / SpiderFoot / Splunk as "
             "callable tools. Claude calls them; Claude reads the responses; doctrine bounds what's allowed."),
            ("Subagent context isolation",
             "Each role gets its own context — collector never sees the brief; briefer never sees raw "
             "articles; librarian is the only writer to git/Splunk/Discord. Lower error rate by construction."),
        ],
        body_size=17,
        sub_size=13,
    )
    add_notes(
        s,
        "This slide is the 'so what' for the dev audience. They already know Claude — frame it in terms they understand: long-context reads, tool use, structured outputs, multi-agent. "
        "The CTI tradecraft slides earlier set up WHAT to apply Claude to; this slide says WHY Claude is the right engine. "
        "Land hard on context isolation — it's the design insight most teams miss when they try to build something like this with one big prompt."
    )
    add_footer(s)
    add_pagenum(s, 5, TOTAL)

    # ===================================================================
    # Section 2 — What it does (slides 6-13)
    # ===================================================================

    # Slide 6 — Daily rhythm
    s = blank_slide(prs)
    add_title(s, "Daily rhythm", subtitle="The cadence is the product. Fully autonomous; humans approve only HIGH-impact calls.")

    schedule = [
        ("00:00 EDT", "FLASH sweep", "Check for triggers; queue if outside active hours"),
        ("06:00 EDT", "FLASH sweep", "Same"),
        ("07:30 EDT", "Pre-brief collection", "Pull last 14h of OSINT into raw signal"),
        ("08:00 EDT", "MORNING BRIEF", "Graded findings → analyst markdown + Discord summary"),
        ("12:00 EDT", "FLASH sweep", "Same"),
        ("15:30 EDT", "Pre-brief collection", "Same"),
        ("16:00 EDT", "AFTERNOON BRIEF", "Same"),
        ("18:00 EDT", "FLASH sweep", "Same"),
        ("Wed 10:30", "Threat Detection Weekly", "Detection-engineering focus"),
        ("Fri 12:00", "Threat Actor Summary", "Deep dive on 1-2 tracked actors"),
        ("Sun 10:00", "Weekly Synthesis", "Patterns across the week"),
    ]

    # Render as a tidy table
    tbl_left = Inches(0.8)
    tbl_top = Inches(1.7)
    row_h = 0.4
    col_widths = [Inches(1.6), Inches(3.2), Inches(7.0)]

    # Header row
    headers = ["TIME (EDT)", "EVENT", "OUTPUT"]
    for ci, h in enumerate(headers):
        x = tbl_left + sum(col_widths[:ci], Emu(0))
        hbox = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, tbl_top, col_widths[ci], Inches(row_h))
        hbox.line.fill.background()
        hbox.fill.solid()
        hbox.fill.fore_color.rgb = C_TITLE
        htf = hbox.text_frame
        htf.margin_top = Inches(0.05)
        hp = htf.paragraphs[0]
        hp.alignment = PP_ALIGN.LEFT if ci > 0 else PP_ALIGN.CENTER
        hr = hp.add_run()
        hr.text = "  " + h
        hr.font.size = Pt(11)
        hr.font.bold = True
        hr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        hr.font.name = "Calibri"

    for ri, (t, ev, out) in enumerate(schedule):
        y = tbl_top + Inches(row_h * (ri + 1))
        is_brief = "BRIEF" in ev.upper()
        bg_color = RGBColor(0xE8, 0xEF, 0xF5) if is_brief else (
            RGBColor(0xFF, 0xFF, 0xFF) if ri % 2 == 0 else RGBColor(0xF8, 0xF8, 0xF8)
        )
        for ci, val in enumerate([t, ev, out]):
            x = tbl_left + sum(col_widths[:ci], Emu(0))
            rbox = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_widths[ci], Inches(row_h))
            rbox.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            rbox.line.width = Pt(0.5)
            rbox.fill.solid()
            rbox.fill.fore_color.rgb = bg_color
            tf = rbox.text_frame
            tf.margin_top = Inches(0.05)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = "  " + val
            r.font.size = Pt(11)
            r.font.bold = is_brief and ci == 1
            r.font.color.rgb = C_BODY
            r.font.name = "Calibri"

    add_notes(
        s,
        "Point out: this is the same cadence as a real intel shop. The agent runs without human oversight between events — 30+ hours of unattended ops across the launch month with zero operator interventions. "
        "FLASH sweeps catch breaking news between scheduled briefs. Quiet hours (21:00-09:00 EDT) queue routine FLASHes unless the critical-override conditions all hit (CVSS 10 + active exploitation + tracked actor + sector watchlist) — then it bypasses and pings immediately."
    )
    add_footer(s)
    add_pagenum(s, 6, TOTAL)

    # ===================================================================
    # Slide 7 — A real brief (Layer 2 / Discord)
    # ===================================================================
    s = blank_slide(prs)
    add_title(
        s,
        "A real brief — Discord summary",
        subtitle="Layer 2 — what leadership reads on their phone at 0800. ~240 words, source-linked, natural-language dates.",
    )

    discord_text = """Good morning. Here's your 0800 brief — Monday, May 11.

🚨  Active Threats

•  Checkmarx Jenkins AST plugin compromised in supply-chain attack
    Checkmarx warned Friday May 9 of a malicious Jenkins AST plugin on the
    Marketplace; weekend variants surfaced on GitHub. Remediation 2.0.13-848
    is out. DIB CI/CD: pin to fix version now, capture plugin hashes for IOC
    backfill.

•  SailPoint discloses GitHub repository hack
    SailPoint disclosed an April 20 incident; a third-party app vulnerability
    exposed a subset of its GitHub repos. No production customer data
    accessed; some customer info was in the repos, extent undisclosed.
    SecurityWeek floats a possible TeamPCP link; Archimedes does not endorse.

🔓  Vulnerabilities

•  CVE-2026-6973 (Ivanti EPMM on-prem): federal patch deadline expired last
   night. Unpatched on-prem fleet now in non-compliance + standing
   exploitation risk.

•  CVE-2026-42208 (BerriAI LiteLLM): KEV deadline closes today (FCEB only).
   LiteLLM-proxying shops: inventory and patch by EOD."""
    add_code_block(s, discord_text, top=1.5, height=5.4, font_size=10, discord=True)
    add_notes(
        s,
        "This is the actual Discord post that shipped from #intel-briefs on 2026-05-11 — message ID is on the slide notes if anyone asks. "
        "Three design choices to call out: (1) headline IS the hyperlink to the source article; (2) dates are natural-language so freshness is scan-readable; (3) no themed-character voice — clean professional. "
        "Discord renders the bullets with bold-blue link headlines and embed previews under each link; show that visually if you can pull it up live."
    )
    add_footer(s)
    add_pagenum(s, 7, TOTAL)

    # ===================================================================
    # Slide 8 — Same brief, Layer 1 (analyst-grade)
    # ===================================================================
    s = blank_slide(prs)
    add_title(
        s,
        "Same brief — analyst-grade markdown",
        subtitle="Layer 1 — what the intel team reviews. ~730 words. Lives in git. Audit-grade.",
    )

    markdown_excerpt = """[Two enterprise security vendors — Checkmarx and SailPoint — disclosed
supply-chain compromises on SecurityWeek 18 minutes apart this morning;
both relays are vendor-self-disclosure restatements, not new third-party
research...](https://www.securityweek.com/)
Per Hard Rule 2, the TeamPCP attribution to Checkmarx Jenkins AST is a
restatement of prior reporting; SecurityWeek's framing of the SailPoint-
TeamPCP link is explicitly speculative. Archimedes does not endorse either.

Why it matters: Both vendors are widely deployed in DIB CI/CD and identity-
governance stacks — Jenkins AST in Checkmarx-using DevSecOps pipelines,
SailPoint IGA in CMMC and ITAR-program access governance — so the structural
exposure is real even where no A&D prime has been named as a victim...

## Active Threats

**[Checkmarx warned of malicious Jenkins AST plugin published to Jenkins
Marketplace](https://www.securityweek.com/...)** ... Digraph: B2 · WEP:
likely (procedural facts only — relational/attribution layer carries lower
WEP) · finding-2026-05-11-0001.

🔗 Connects to: Actor #001 TeamPCP (HIGH, roster) — Jenkins AST plugin
distribution channel is the first appearance of CI/CD-pipeline-poisoning
in the TeamPCP TTP register; dossier update queued via actor-profiler."""

    add_code_block(s, markdown_excerpt, top=1.5, height=5.4, font_size=10)
    add_notes(
        s,
        "Same morning, same findings — but this is the analyst-grade version. Note the doctrine framings throughout: Hard Rule 2 (no originating attribution), Digraph + WEP per item, finding-id citations, single-source-veto annotations. "
        "The intel team red-team-reviews this layer. Leadership never sees it — they get Layer 2. ONE canonical record, TWO renderings, one source of truth in git. "
        "Critical insight for OpenCTI builders: this maps to STIX-Report objects with extended properties for the doctrine framings."
    )
    add_footer(s)
    add_pagenum(s, 8, TOTAL)

    # ===================================================================
    # Slide 9 — Source grading (Admiralty)
    # ===================================================================
    s = blank_slide(prs)
    add_title(s, "Source grading — Admiralty 6×6", subtitle="NATO-standard reliability + credibility framework, applied per source per finding.")

    # Reliability column (A-F)
    add_bullets(
        s,
        [
            ("A — Completely reliable", "Sustained track record across many findings (Mandiant, Unit 42, MSTIC, CrowdStrike, MITRE, CISA)."),
            ("B — Usually reliable", "Established but with occasional gaps (BleepingComputer, The Record, Krebs)."),
            ("C — Fairly reliable", "Vendor blogs with mixed history; new entrants on first citation."),
            ("D / E / F", "Unreliable / unknown / questionable. Findings from D-F can monitor but never load-bear."),
        ],
        left=0.5, top=1.7, width=6.0, height=3.5, body_size=15, sub_size=12,
    )

    add_bullets(
        s,
        [
            ("1 — Confirmed", "Multiple independent A/B-grade sources agree."),
            ("2 — Probably true", "One A-grade source + plausibility test passes."),
            ("3 — Possibly true", "Single source, cannot be confirmed or denied."),
            ("4 / 5 / 6", "Doubtful / improbable / cannot be judged."),
        ],
        left=6.8, top=1.7, width=6.0, height=3.5, body_size=15, sub_size=12,
    )

    # Footer note
    footer_box = s.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.33), Inches(1.5))
    ftf = footer_box.text_frame
    ftf.word_wrap = True
    fp = ftf.paragraphs[0]
    fr = fp.add_run()
    fr.text = (
        "Every finding gets a digraph (e.g., A2, B3, C3). A1 = confirmed by completely reliable source. "
        "C3 = possibly true from a single fairly-reliable source. The single-source-veto rule caps any forward "
        "claim at \"likely\" when only one source backs it — even if both source and content are A-grade."
    )
    fr.font.size = Pt(13)
    fr.font.italic = True
    fr.font.color.rgb = C_SUBTLE
    fr.font.name = "Calibri"

    add_notes(
        s,
        "The CTI audience knows this; the dev audience may not. Frame it as: 'this is older than cyber — NATO has used Admiralty since the Cold War.' Why it matters: one bad source poisons a brief. "
        "The single-source-veto is the load-bearing safety primitive — it's what stops the agent from over-promoting unverified rumors into 'very likely' assessments. "
        "OpenCTI maps this cleanly to opinion / confidence fields on STIX SDOs."
    )
    add_footer(s)
    add_pagenum(s, 9, TOTAL)

    # ===================================================================
    # Slide 10 — WEP
    # ===================================================================
    s = blank_slide(prs)
    add_title(s, "Words of Estimative Probability", subtitle="CIA-standard probability vocabulary. The agent uses these — and refuses bare hedges like \"may\", \"might\", \"could.\"")

    wep_rows = [
        ("Almost certainly", "≥95%", C_TITLE),
        ("Very likely", "85–95%", C_ACCENT),
        ("Likely", "55–85%", C_ACCENT),
        ("Roughly even chance", "~50%", C_SUBTLE),
        ("Unlikely", "15–45%", C_SUBTLE),
        ("Very unlikely", "5–15%", C_SUBTLE),
        ("Remote", "<5%", C_SUBTLE),
    ]

    tbl_left = Inches(2.5)
    tbl_top = Inches(1.8)
    row_h = 0.55
    for ri, (term, band, color) in enumerate(wep_rows):
        y = tbl_top + Inches(row_h * ri)
        # Term cell
        tbox = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, tbl_left, y, Inches(4.5), Inches(row_h))
        tbox.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        tbox.line.width = Pt(0.5)
        tbox.fill.solid()
        tbox.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        ttf = tbox.text_frame
        ttf.margin_top = Inches(0.1)
        tp = ttf.paragraphs[0]
        tr = tp.add_run()
        tr.text = "  " + term
        tr.font.size = Pt(15)
        tr.font.bold = True
        tr.font.color.rgb = color
        tr.font.name = "Calibri"

        # Band cell
        bbox = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, tbl_left + Inches(4.5), y, Inches(2.5), Inches(row_h))
        bbox.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        bbox.line.width = Pt(0.5)
        bbox.fill.solid()
        bbox.fill.fore_color.rgb = C_BG_LIGHT
        btf = bbox.text_frame
        btf.margin_top = Inches(0.1)
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = band
        br.font.size = Pt(15)
        br.font.color.rgb = C_BODY
        br.font.name = "Consolas"

    # Why-it-matters footer
    footer_box = s.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12.33), Inches(0.8))
    ftf = footer_box.text_frame
    ftf.word_wrap = True
    fp = ftf.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    fr = fp.add_run()
    fr.text = (
        "Why: hedging language is the #1 way analysts get misread by leadership. "
        "\"Could be a problem\" gets ignored; \"likely escalates within 72 hours\" gets a meeting."
    )
    fr.font.size = Pt(13)
    fr.font.italic = True
    fr.font.color.rgb = C_SUBTLE
    fr.font.name = "Calibri"

    add_notes(
        s,
        "WEP comes from a 1964 CIA paper by Sherman Kent — frame it as 'this is how every IC analyst has been trained for 60 years.' "
        "Show that the preflight checklist rejects bare 'may/might/could' — every forward claim has to use the WEP vocabulary or regenerate. "
        "Dev audience analog: this is type-safety for probability statements."
    )
    add_footer(s)
    add_pagenum(s, 10, TOTAL)

    # ===================================================================
    # Slide 11 — FLASH alerts
    # ===================================================================
    s = blank_slide(prs)
    add_title(s, "FLASH alerts", subtitle="Async — for when something can't wait for the next scheduled brief.")
    add_bullets(
        s,
        [
            ("Sweeps every 6 hours", "00:00 / 06:00 / 12:00 / 18:00 EDT. Collector + grader on the fast path."),
            ("Triggers (any one fires a FLASH)",
             "(1) Critical CVE actively exploited · (2) New attribution to a tracked actor · "
             "(3) Major breach in scope · (4) Tracked actor TTP change · (5) Sector-specific zero-day · "
             "(6) Zero-day with no patch · (7) Supply-chain compromise impacting tracked vendors"),
            ("Quiet hours: 21:00 – 09:00 EDT", "Routine FLASHes queue to 09:00 catchup sweep — operators are not woken up casually."),
            ("Critical override bypasses quiet hours",
             "ALL of: CVSS 10 + confirmed active exploitation + tracked actor + sector watchlist hit. "
             "Anything less queues. This is the 'actually wake up' threshold."),
            ("72-hour auto-downgrade clock on single-source FLASHes",
             "If no independent A/B-grade source corroborates within 72h, attribution leg auto-drops to "
             "C3 'possibly true.' Tested and fired clean on MuddyWater 2026-05-09."),
        ],
        body_size=16, sub_size=12,
    )
    add_notes(
        s,
        "FLASH is the answer to 'what if something breaks between the morning and afternoon briefs?' "
        "Critical-override has only fired once in launch ops — the gate is intentionally tight. False FLASHes train operators to ignore them; a high-precision channel preserves the 'wake-up' signal. "
        "The 72h auto-downgrade clock is the doctrine pattern for handling unverified single-source attribution — ship the FLASH, but commit to retracting/downgrading if independent corroboration doesn't arrive."
    )
    add_footer(s)
    add_pagenum(s, 11, TOTAL)

    # ===================================================================
    # Slide 12 — Threat actor profiles
    # ===================================================================
    s = blank_slide(prs)
    add_title(s, "Threat actor profiles", subtitle="Per-actor dossier: profile, IOCs, threat-box scoring. Reviewed every 90 days.")

    # Two columns: left = roster overview, right = scoring example
    add_bullets(
        s,
        [
            ("23 actors tracked",
             "Iranian cluster (UNC1549, MuddyWater, APT34, Charming Kitten, Handala), DPRK (Lazarus, "
             "APT37, Stardust Chollima), Chinese (APT28, APT40), unattributed (TeamPCP, Cl0p, Lapsus$)."),
            ("Per-actor dossier",
             "profile.md (background), iocs.yaml (machine-readable), threat-box.yaml (scoring), "
             "threat-box.md (narrative)."),
            ("Threat-box methodology",
             "Per-category composite scores (Espionage 35%, Supply Chain 20%, Destructive 15%, "
             "Disruptive 15%, Cyber-Crime 15%) weighted to overall HIGH / MEDIUM / LOW."),
            ("Hard Rule 5: HIGH requires human sign-off",
             "Agent never auto-commits a HIGH score. Posts to #actor-review, waits for /approve-scoring."),
        ],
        left=0.5, top=1.7, width=6.2, height=4.5, body_size=15, sub_size=12,
    )

    # Right column — scoring sample
    sample_box_y = Inches(1.7)
    sample_box = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(7.0), sample_box_y, Inches(5.8), Inches(4.5)
    )
    sample_box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    sample_box.line.width = Pt(0.5)
    sample_box.fill.solid()
    sample_box.fill.fore_color.rgb = C_BG_LIGHT

    sample_tf = sample_box.text_frame
    sample_tf.margin_top = Inches(0.2)
    sample_tf.margin_left = Inches(0.2)
    sample_tf.margin_right = Inches(0.2)
    sample_tf.word_wrap = True

    sample_lines = [
        ("APT37 (#024) — /update-tracking 2026-05-10", True, 16, C_TITLE),
        ("", False, 12, C_BODY),
        ("Espionage          composite 8/10  →  HIGH", False, 12, C_BODY),
        ("Supply Chain       composite 6/10  →  MEDIUM", False, 12, C_BODY),
        ("Destructive        composite 2/10  →  LOW", False, 12, C_BODY),
        ("Disruptive         composite 2/10  →  LOW", False, 12, C_BODY),
        ("Cyber-Crime        composite 2/10  →  LOW", False, 12, C_BODY),
        ("", False, 12, C_BODY),
        ("Weighted overall: 4.9  →  MEDIUM", True, 14, C_TITLE),
        ("", False, 12, C_BODY),
        ("Why not HIGH:", True, 12, C_BODY),
        ("Intent capped at 3 (Sector Association) — public reporting", False, 11, C_SUBTLE),
        ("names think tanks / journalists / defectors, no A&D primes.", False, 11, C_SUBTLE),
        ("The methodology refused to extrapolate.", False, 11, C_SUBTLE),
    ]
    first = True
    for text, bold, size, color in sample_lines:
        if first:
            p = sample_tf.paragraphs[0]
            first = False
        else:
            p = sample_tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Consolas" if "composite" in text or "Weighted" in text else "Calibri"

    add_notes(
        s,
        "Critical empirical observation: 5 of 5 actors scored to date landed non-HIGH (UNC1549 MEDIUM, Charming Kitten LOW, APT34 MEDIUM, MuddyWater LOW, APT37 MEDIUM). "
        "The gate is tight by design — HIGH requires an A-grade source explicitly naming an A&D-prime victim, which most public reporting on these actors does not provide. "
        "When HIGH eventually lands, it'll mean something. Audience takeaway: the framework refuses to invent confidence the source material doesn't support."
    )
    add_footer(s)
    add_pagenum(s, 12, TOTAL)

    # ===================================================================
    # Slide 13 — Hard Rules
    # ===================================================================
    s = blank_slide(prs)
    add_title(s, "Hard Rules — the safety boundary", subtitle="Non-negotiable. Agent refuses even when prompted to bypass.")
    add_bullets(
        s,
        [
            ("1 · Legal policy is non-negotiable", "Active recon only against authorized targets. Refuses to query PII-heavy sources without justification."),
            ("2 · Never originate attribution", "Only reports what other sources have attributed, citing them. The agent does not invent attribution claims."),
            ("3 · No exploitation, ever", "Never generates PoC code, payloads, or exploit guides. Not for testing, not for research, not 'educational.'"),
            ("4 · Never scan third parties", "Active reconnaissance only against targets in authorized-targets.yaml. Passive-only for everything else. SpiderFoot + theHarvester MCPs enforce module-level allowlists."),
            ("5 · Human sign-off for HIGH threat levels", "When actor-profiler proposes HIGH, posts to #actor-review and waits for /approve-scoring. Does not auto-commit."),
            ("6 · 15-word quote limit, one quote per source", "Copyright compliance. Hard-enforced in the preflight checklist; refuses to ship a brief that violates it."),
            ("7 · Credentials are radioactive", "If a query surfaces credentials, never stores them. Counts and reports exposure; discards."),
            ("8 · Splunk first-party > external sources", "When first-party telemetry contradicts an external source, first-party wins and the external source gets graded down."),
        ],
        body_size=14, sub_size=12,
    )
    add_notes(
        s,
        "Hard Rules are the audit-defense layer. CTI is regulated-adjacent — ITAR, export-control, copyright, evidence chain. "
        "Frame for the dev audience: these are operational constraints baked into the doctrine files that load into every agent invocation. The agent refuses to violate them even when the prompt asks it to. "
        "Show what the refusal looks like — collector subagent will refuse a non-authorized port-scan request and log to policy-violations.yaml. This is what 'trust' looks like for a CTI tool."
    )
    add_footer(s)
    add_pagenum(s, 13, TOTAL)

    # ===================================================================
    # Section 3 — How (light architecture)
    # ===================================================================

    # ===================================================================
    # Slide 14 — Architecture diagram
    # ===================================================================
    s = blank_slide(prs)
    add_title(s, "Architecture", subtitle="Orchestrator + 9 subagents + ~8 MCP wrappers + doctrine. Everything flows through git, Splunk, Discord.")

    # Top: data sources (left), orchestrator (center), outputs (right)
    # Middle: subagent ring
    # Bottom: doctrine band

    # Orchestrator box (top center)
    orch = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.7), Inches(1.6),
        Inches(2.0), Inches(0.6),
    )
    orch.line.color.rgb = C_TITLE
    orch.line.width = Pt(2)
    orch.fill.solid()
    orch.fill.fore_color.rgb = C_TITLE
    otf = orch.text_frame
    otf.margin_top = Inches(0.08)
    op = otf.paragraphs[0]
    op.alignment = PP_ALIGN.CENTER
    orun = op.add_run()
    orun.text = "Claude (orchestrator)"
    orun.font.size = Pt(13)
    orun.font.bold = True
    orun.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    orun.font.name = "Calibri"

    # MCP block (top left)
    mcp = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.6), Inches(2.6), Inches(0.6)
    )
    mcp.line.color.rgb = C_ACCENT
    mcp.line.width = Pt(1.5)
    mcp.fill.solid()
    mcp.fill.fore_color.rgb = C_BG_LIGHT
    mtf = mcp.text_frame
    mtf.margin_top = Inches(0.08)
    mp = mtf.paragraphs[0]
    mp.alignment = PP_ALIGN.CENTER
    mrun = mp.add_run()
    mrun.text = "8 MCP wrappers"
    mrun.font.size = Pt(12)
    mrun.font.bold = True
    mrun.font.color.rgb = C_ACCENT
    mrun.font.name = "Calibri"

    mcp_sub = s.shapes.add_textbox(Inches(0.5), Inches(2.25), Inches(2.6), Inches(0.5))
    msub_tf = mcp_sub.text_frame
    msub_tf.word_wrap = True
    msub_p = msub_tf.paragraphs[0]
    msub_p.alignment = PP_ALIGN.CENTER
    msub_run = msub_p.add_run()
    msub_run.text = "Splunk · VT · Shodan · Censys\nurlscan · theHarvester · SpiderFoot · RSS"
    msub_run.font.size = Pt(9)
    msub_run.font.color.rgb = C_SUBTLE
    msub_run.font.name = "Calibri"

    # Outputs (top right)
    out = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(10.2), Inches(1.6), Inches(2.6), Inches(0.6)
    )
    out.line.color.rgb = C_REFUSAL
    out.line.width = Pt(1.5)
    out.fill.solid()
    out.fill.fore_color.rgb = C_BG_LIGHT
    otf2 = out.text_frame
    otf2.margin_top = Inches(0.08)
    op2 = otf2.paragraphs[0]
    op2.alignment = PP_ALIGN.CENTER
    orun2 = op2.add_run()
    orun2.text = "Outputs"
    orun2.font.size = Pt(12)
    orun2.font.bold = True
    orun2.font.color.rgb = C_REFUSAL
    orun2.font.name = "Calibri"

    out_sub = s.shapes.add_textbox(Inches(10.2), Inches(2.25), Inches(2.6), Inches(0.5))
    osub_tf = out_sub.text_frame
    osub_tf.word_wrap = True
    osub_p = osub_tf.paragraphs[0]
    osub_p.alignment = PP_ALIGN.CENTER
    osub_run = osub_p.add_run()
    osub_run.text = "Discord  ·  git  ·  Splunk"
    osub_run.font.size = Pt(9)
    osub_run.font.color.rgb = C_SUBTLE
    osub_run.font.name = "Calibri"

    # Arrows
    arrow_l = s.shapes.add_connector(1, Inches(3.1), Inches(1.9), Inches(5.7), Inches(1.9))
    arrow_l.line.color.rgb = C_BODY
    arrow_l.line.width = Pt(1.5)
    arrow_r = s.shapes.add_connector(1, Inches(7.7), Inches(1.9), Inches(10.2), Inches(1.9))
    arrow_r.line.color.rgb = C_BODY
    arrow_r.line.width = Pt(1.5)

    # Subagent ring (middle)
    subagents = [
        ("collector", "raw signal"),
        ("grader", "Admiralty + WEP"),
        ("analyst", "SATs / ACH / KAC"),
        ("red-team-analyst", "challenge HIGH"),
        ("actor-profiler", "dossiers"),
        ("vuln-tracker", "CVE corpus"),
        ("briefer", "Layer 1 + 2"),
        ("librarian", "git / Splunk / Discord"),
        ("(human approval)", "HIGH scores"),
    ]
    ring_top = Inches(3.4)
    box_w = Inches(1.35)
    box_h = Inches(0.95)
    n = len(subagents)
    spacing = Inches(0.05)
    total_w = box_w * n + spacing * (n - 1)
    ring_left = (Inches(13.333) - total_w) / 2
    for i, (name, role) in enumerate(subagents):
        x = ring_left + (box_w + spacing) * i
        is_human = "human" in name.lower()
        sb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, ring_top, box_w, box_h)
        sb.line.color.rgb = C_REFUSAL if is_human else C_ACCENT
        sb.line.width = Pt(1.2)
        sb.fill.solid()
        sb.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if not is_human else RGBColor(0xFD, 0xEF, 0xED)
        stf = sb.text_frame
        stf.margin_top = Inches(0.08)
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        srun = sp.add_run()
        srun.text = name
        srun.font.size = Pt(10)
        srun.font.bold = True
        srun.font.color.rgb = C_TITLE if not is_human else C_REFUSAL
        srun.font.name = "Calibri"

        sp2 = stf.add_paragraph()
        sp2.alignment = PP_ALIGN.CENTER
        srun2 = sp2.add_run()
        srun2.text = role
        srun2.font.size = Pt(8)
        srun2.font.italic = True
        srun2.font.color.rgb = C_SUBTLE
        srun2.font.name = "Calibri"

    # Doctrine band (bottom)
    doc_top = Inches(5.0)
    doc_band = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), doc_top, Inches(12.33), Inches(0.6)
    )
    doc_band.line.color.rgb = C_SUBTLE
    doc_band.line.width = Pt(1)
    doc_band.fill.solid()
    doc_band.fill.fore_color.rgb = C_BG_LIGHT
    dtf = doc_band.text_frame
    dtf.margin_top = Inches(0.1)
    dp = dtf.paragraphs[0]
    dp.alignment = PP_ALIGN.CENTER
    drun = dp.add_run()
    drun.text = (
        "Doctrine (.md files, versioned in git):  LEGAL-POLICY  ·  INTEL-GRADING  ·  "
        "INTEL-BRIEF-STANDARDS  ·  THREAT-BOX-METHODOLOGY  ·  RETRACTION-POLICY  ·  FLASH-POLICY"
    )
    drun.font.size = Pt(11)
    drun.font.bold = True
    drun.font.color.rgb = C_BODY
    drun.font.name = "Calibri"

    # Caption
    cap_box = s.shapes.add_textbox(Inches(0.5), Inches(5.9), Inches(12.33), Inches(1.0))
    cap_tf = cap_box.text_frame
    cap_tf.word_wrap = True
    cap_p = cap_tf.paragraphs[0]
    cap_p.alignment = PP_ALIGN.CENTER
    cap_run = cap_p.add_run()
    cap_run.text = (
        "Each subagent has its own context window, its own write scope, its own doctrine subset. "
        "Orchestrator never reads raw articles. Librarian is the only writer to git / Splunk / Discord. "
        "Context isolation is the design that makes the multi-agent setup robust."
    )
    cap_run.font.size = Pt(12)
    cap_run.font.italic = True
    cap_run.font.color.rgb = C_SUBTLE
    cap_run.font.name = "Calibri"

    add_notes(
        s,
        "This is the slide where devs will start drawing parallels to their own stack. Land the three big ideas: "
        "(1) subagents = role separation, (2) doctrine files = persistent prompts that load on every run, (3) one writer to externalize side effects (the librarian). "
        "If they remember nothing else from the architecture section, this slide is the thing to remember."
    )
    add_footer(s)
    add_pagenum(s, 14, TOTAL)

    # ===================================================================
    # Slide 15 — Subagents = role separation
    # ===================================================================
    s = blank_slide(prs)
    add_title(s, "Subagents = role separation", subtitle="Each agent's job description, write scope, and doctrine are minimal.")

    rows = [
        ("collector", "raw OSINT collection", "threats/raw-signal/", "source-grades, watchlists, FLASH triggers"),
        ("grader", "Admiralty + WEP promotion", "threats/findings/", "INTEL-GRADING"),
        ("analyst", "SAT/ACH/KAC structured analysis", "findings (analysis sections)", "ACH/KAC method"),
        ("red-team-analyst", "challenge HIGH-WEP findings", "findings (red-team section)", "contrarian SAT"),
        ("actor-profiler", "dossier maintenance + scoring", "threats/threat-actors/*", "THREAT-BOX-METHODOLOGY, ACTOR-PROFILE-STANDARD"),
        ("vuln-tracker", "CVE tracking + KEV monitoring", "threats/vulnerabilities/*", "NVD + KEV"),
        ("briefer", "compose all brief types", "threats/briefs/", "INTEL-BRIEF-STANDARDS, smart-brevity skill"),
        ("librarian", "ship + commit + log + index", "git, Splunk, Discord, indices", "RETRACTION-POLICY"),
    ]

    tbl_left = Inches(0.5)
    tbl_top = Inches(1.7)
    row_h = 0.5
    col_widths = [Inches(2.0), Inches(3.5), Inches(3.3), Inches(3.5)]
    headers = ["Subagent", "Role", "Write scope", "Doctrine read"]

    for ci, h in enumerate(headers):
        x = tbl_left + sum(col_widths[:ci], Emu(0))
        hb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, tbl_top, col_widths[ci], Inches(row_h))
        hb.line.fill.background()
        hb.fill.solid()
        hb.fill.fore_color.rgb = C_TITLE
        htf = hb.text_frame
        htf.margin_top = Inches(0.1)
        hp = htf.paragraphs[0]
        hp.alignment = PP_ALIGN.LEFT
        hr = hp.add_run()
        hr.text = "  " + h
        hr.font.size = Pt(11)
        hr.font.bold = True
        hr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        hr.font.name = "Calibri"

    for ri, row in enumerate(rows):
        y = tbl_top + Inches(row_h * (ri + 1))
        bg = RGBColor(0xFF, 0xFF, 0xFF) if ri % 2 == 0 else RGBColor(0xF8, 0xF8, 0xF8)
        for ci, val in enumerate(row):
            x = tbl_left + sum(col_widths[:ci], Emu(0))
            cb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_widths[ci], Inches(row_h))
            cb.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            cb.line.width = Pt(0.5)
            cb.fill.solid()
            cb.fill.fore_color.rgb = bg
            tf = cb.text_frame
            tf.margin_top = Inches(0.1)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = "  " + val
            r.font.size = Pt(10)
            r.font.bold = (ci == 0)
            r.font.color.rgb = C_BODY
            r.font.name = "Calibri"

    add_notes(
        s,
        "The audience should walk away noticing: write scopes never overlap. Only the librarian can git push. Only the actor-profiler can edit a dossier. The orchestrator never writes anything. "
        "That's the design that prevents context bleed and lets red-team-analyst be genuinely contrarian — it has no investment in the primary finding because it never wrote one. "
        "Counterintuitive but real: less context per agent = lower error rate. Bigger prompts are not always better."
    )
    add_footer(s)
    add_pagenum(s, 15, TOTAL)

    # ===================================================================
    # Slide 16 — Doctrine as code
    # ===================================================================
    s = blank_slide(prs)
    add_title(s, "Doctrine as code", subtitle="The agent's working memory. Versioned in git. Reviewed like code.")

    add_bullets(
        s,
        [
            ("doctrine/LEGAL-POLICY.md",
             "Prohibited query patterns, authorized targets, ITAR boundary, copyright. Read before every tool call."),
            ("doctrine/INTEL-GRADING.md",
             "Admiralty 6×6 + WEP + single-source-veto + corroboration rules. Read by grader on every promotion."),
            ("doctrine/INTEL-BRIEF-STANDARDS.md",
             "Layered format (analyst-grade + Discord), word/char budgets, Smart Brevity, preflight checklist. Read by briefer."),
            ("doctrine/THREAT-BOX-METHODOLOGY.md",
             "Per-category composite scoring with evidence-minimum tables. Read by actor-profiler on every /update-tracking."),
            ("doctrine/RETRACTION-POLICY.md",
             "When to retract vs. correct, additive (never silent), 72h auto-downgrade clock. Read by grader + librarian."),
            ("doctrine/FLASH-POLICY.md",
             "7 trigger conditions, quiet hours, critical-override threshold. Read by collector on every sweep."),
        ],
        body_size=14, sub_size=12,
    )

    add_notes(
        s,
        "This is the slide where the OpenCTI angle matters: doctrine files map to versioned policy docs in your repo, loaded into the agent's context as system prompts or via skill files. "
        "Frame: doctrine is what the agent's job description USED to be locked inside a single prompt. Splitting it into versioned files means you can review, test, and iterate doctrine like code. "
        "Show one example briefly — INTEL-GRADING.md has the Admiralty 6×6 lookup table + the single-source-veto rule + how to handle source corroboration. Plain English. Reviewable by an analyst lead, not just an ML engineer."
    )
    add_footer(s)
    add_pagenum(s, 16, TOTAL)

    # ===================================================================
    # Slide 17 — Where OpenCTI fits
    # ===================================================================
    s = blank_slide(prs)
    add_title(s, "Where OpenCTI fits in your stack", subtitle="Same shape, different substrate. Map Archimedes patterns onto STIX entities + Teams delivery.")

    # Two-column compare
    left_x = Inches(0.5)
    right_x = Inches(6.8)
    col_w = Inches(6.0)
    col_h = Inches(4.8)
    top = Inches(1.8)

    # Left column — Archimedes
    lb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, top, col_w, Inches(0.45))
    lb.line.fill.background()
    lb.fill.solid()
    lb.fill.fore_color.rgb = C_ACCENT
    ltf = lb.text_frame
    ltf.margin_top = Inches(0.07)
    lp = ltf.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    lr = lp.add_run()
    lr.text = "ARCHIMEDES (what we built)"
    lr.font.size = Pt(13)
    lr.font.bold = True
    lr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    lr.font.name = "Calibri"

    # Right column — Your build
    rb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, top, col_w, Inches(0.45))
    rb.line.fill.background()
    rb.fill.solid()
    rb.fill.fore_color.rgb = C_TITLE
    rtf = rb.text_frame
    rtf.margin_top = Inches(0.07)
    rp = rtf.paragraphs[0]
    rp.alignment = PP_ALIGN.CENTER
    rr = rp.add_run()
    rr.text = "YOUR BUILD (OpenCTI + Claude + Teams)"
    rr.font.size = Pt(13)
    rr.font.bold = True
    rr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    rr.font.name = "Calibri"

    rows_l = [
        ("Corpus", "Markdown + YAML in git"),
        ("Entities", "Free-form actor / vuln / finding files"),
        ("Collection", "MCP wrappers around OSINT APIs"),
        ("Analyst layer", "Claude subagents + doctrine"),
        ("Delivery", "Discord (#intel-briefs, #flash-alerts)"),
        ("Audit", "git history + Splunk events + Discord log"),
        ("Source of truth", ".md files; one record, two renderings"),
    ]
    rows_r = [
        ("Corpus", "OpenCTI STIX 2.x entity store"),
        ("Entities", "STIX-native (Threat-Actor, Indicator, Report, Campaign, Vulnerability)"),
        ("Collection", "OpenCTI connectors + optional MCP wrappers"),
        ("Analyst layer", "Same — Claude subagents + doctrine"),
        ("Delivery", "Teams (webhook + Adaptive Cards)"),
        ("Audit", "OpenCTI provenance + git for doctrine + Splunk events"),
        ("Source of truth", "OpenCTI; analyst layer publishes Reports + Opinions"),
    ]

    row_h = 0.55
    for ri, ((kl, vl), (kr, vr)) in enumerate(zip(rows_l, rows_r)):
        y = top + Inches(0.55 + row_h * ri)
        for x_off, key, val, bg in [(left_x, kl, vl, RGBColor(0xF1, 0xF7, 0xFD)),
                                     (right_x, kr, vr, RGBColor(0xEF, 0xF4, 0xFA))]:
            cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_off, y, col_w, Inches(row_h))
            cell.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            cell.line.width = Pt(0.5)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg if ri % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
            tf = cell.text_frame
            tf.margin_top = Inches(0.07)
            tf.margin_left = Inches(0.15)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"{key}  —  {val}"
            run.font.size = Pt(11)
            run.font.color.rgb = C_BODY
            run.font.name = "Calibri"
            # Bold the key
            run.font.bold = False

    add_notes(
        s,
        "This is the big takeaway slide for the audience. The shape transfers: same analyst layer, same doctrine, same subagent decomposition. "
        "What changes: the corpus (OpenCTI STIX entities replace markdown files), connectors replace some of our MCPs, Teams replaces Discord. "
        "Critical: STIX 2.x has Report, Opinion, and Analyst-Note objects that map cleanly to our finding + analyst section + red-team section. The mapping is straightforward."
    )
    add_footer(s)
    add_pagenum(s, 17, TOTAL)

    # ===================================================================
    # Slide 18 — MCP wrappers
    # ===================================================================
    s = blank_slide(prs)
    add_title(s, "MCP wrappers — the integration layer", subtitle="One per external data source. Doctrine enforces what's callable.")

    code = """# theharvester MCP — passive-only enumeration
# (.claude/mcps/theharvester/src/theharvester_mcp/models.py)

PASSIVE_SOURCE_ALLOWLIST: tuple[str, ...] = (
    "crtsh", "otx", "hackertarget", "rapiddns", "certspotter",
    "duckduckgo", "bingsearch", "virustotal", "shodan", "censys",
    # ... 45 entries total — all passive, all OSINT, no PII-heavy
)

# spiderfoot MCP — same pattern, module-level allowlist
PASSIVE_MODULE_ALLOWLIST: tuple[str, ...] = (
    "sfp_dnsresolve", "sfp_whois", "sfp_crt", "sfp_certspotter",
    "sfp_archiveorg", "sfp_threatfox", "sfp_virustotal",
    # ... 45 entries
)

# Active modules are HARD-REJECTED at input validation:
PROHIBITED_MODULES: tuple[str, ...] = (
    "sfp_tool_nmap",        # port scanning
    "sfp_tool_nuclei",      # vulnerability probing
    "sfp_spider",           # web crawling target sites
    "sfp_dnsbrute",         # DNS brute force
    "sfp_screenshot",       # target screenshots
)

# Caller asking for an active module → SpiderFootPolicyError raised
# BEFORE any HTTP call to SpiderFoot. Cannot be bypassed by prompt."""
    add_code_block(s, code, top=1.5, height=5.4, font_size=11)
    add_notes(
        s,
        "Show what 'doctrine enforced at the integration layer' looks like in code. The allowlist is the load-bearing safety primitive — Hard Rule 4 (no scanning third parties) is enforced before the HTTP call, regardless of prompt. "
        "Pattern is portable: wrap any data source as an MCP, encode the policy in the wrapper, the agent literally cannot violate it. "
        "Each of the 8 wrappers has a similar policy primitive — VirusTotal/Censys don't have allowlists but they have rate-limit + auth + structured-output discipline."
    )
    add_footer(s)
    add_pagenum(s, 18, TOTAL)

    # ===================================================================
    # Slide 19 — Audit trail
    # ===================================================================
    s = blank_slide(prs)
    add_title(s, "Audit trail", subtitle="Every action logged. Every change versioned. Every retraction additive.")
    add_bullets(
        s,
        [
            ("Git — full corpus history",
             "Every brief, every finding, every actor dossier update, every doctrine change. "
             "60+ commits across the last week alone. Provenance is queryable: git log + git blame."),
            ("Splunk — operational telemetry",
             "Every brief composed, every Discord post, every refusal, every retraction. "
             "Indexed under 'archimedes' for dashboarding and alerting on the agent's behavior."),
            ("Discord — delivery log",
             "Every channel post archived with message IDs. Retractions appended to original briefs, never replacing."),
            ("Hard Rule logging",
             "Every legal-policy refusal logged to infrastructure/policy-violations.yaml. Surfaces patterns in attempted misuse."),
            ("Why this matters for CTI",
             "CTI is regulated-adjacent — ITAR, export-control, copyright. The audit trail isn't documentation; it's "
             "evidence. When an intel team gets asked 'how did you arrive at this attribution,' the answer is git log + Splunk + Discord."),
        ],
        body_size=15, sub_size=12,
    )
    add_notes(
        s,
        "Counterintuitive insight: agents are MORE auditable than human analysts, not less. Every Claude action is logged with timestamp, tool, doctrine version, and outcome. "
        "Frame this as a feature, not just compliance — it's how you debug a bad brief, how you defend a retraction in court, how you train the next iteration of the doctrine. "
        "OpenCTI's provenance model adds another layer here — STIX has built-in 'created_by_ref' and 'modified' fields that you get for free."
    )
    add_footer(s)
    add_pagenum(s, 19, TOTAL)

    # ===================================================================
    # Section 4 — Lessons (slides 20-22)
    # ===================================================================

    # Slide 20 — Build trust before automation
    s = blank_slide(prs)
    add_title(s, "Lesson 1 — Build trust before automation", subtitle="The first bad brief breaks intel-team trust permanently.")
    add_bullets(
        s,
        [
            ("Human approval gates on high-stakes calls",
             "HIGH actor threat-box scoring → /approve-scoring. New actor scaffolding → /new-actor. "
             "Source-grade ratification → human review. The agent proposes; the human approves."),
            ("Agent never originates attribution",
             "Hard Rule 2. The agent only reports what cited sources have attributed. If Mandiant says X, "
             "Archimedes says 'Mandiant says X.' If no one has said X yet, Archimedes refuses to be the first."),
            ("Single-source veto on 'very likely' or higher",
             "Even if a source is A-grade, a single citation caps the forward claim at 'likely.' The 72h auto-downgrade "
             "clock then forces re-grading if no independent corroboration arrives."),
            ("Retraction is additive, never silent",
             "When something ships wrong (and it will), the original brief stays in the record with a retraction "
             "appended. The record of being wrong IS part of the record."),
            ("Five non-HIGH scorings empirically validate the gate",
             "UNC1549, Charming Kitten, MuddyWater, APT34, APT37 — zero auto-HIGH. The gate is honest; the "
             "methodology refuses to invent confidence the evidence doesn't support."),
        ],
        body_size=14, sub_size=12,
    )
    add_notes(
        s,
        "This is the most important slide of the talk. If the audience walks away with one thing it should be: build the gates BEFORE you build the automation. "
        "Pattern: agent proposes, human approves on anything with reputational stakes. The agent doesn't make the call where being wrong has consequences — it surfaces the call to the human. "
        "This is also why the agent never originates attribution: attribution is a load-bearing claim with policy + legal + reputational consequences."
    )
    add_footer(s)
    add_pagenum(s, 20, TOTAL)

    # ===================================================================
    # Slide 21 — Live validation > mocks
    # ===================================================================
    s = blank_slide(prs)
    add_title(s, "Lesson 2 — Live validation > mocks for tool wrappers", subtitle="6 real bugs caught by live tests that mocks couldn't have surfaced.")

    add_bullets(
        s,
        [
            ("theHarvester 4.10.1 — 3 bugs the mocks missed",
             "(1) PYTHONHOME env poisoning under `uv run` → SRE module mismatch. "
             "(2) Modern JSON shape emits only cmd/hosts/shodan — older docs claimed ips/vhosts/asns. "
             "(3) securityTrails source name is case-sensitive (camelCase) — argparse rejects lowercase."),
            ("SpiderFoot 4.0.0 — 3 more",
             "(1) /ping returns [\"SUCCESS\", \"4.0.0\"] (JSON list), not the documented \"pong\". "
             "(2) /scanstatus is a 7-element list, not 6 — status at index 5 not 4. "
             "(3) /scaneventresultexport is CSV-only — the real JSON endpoint is /scaneventresults."),
            ("All 6 were 'docs match mocks; docs don't match reality'",
             "Every wrapper had unit tests modeled on the documentation. Documentation was wrong/stale. "
             "Mocks happily passed; live runs failed in ways that took an hour each to debug."),
            ("Fix: bind your test suite to a running service",
             "Spin up theHarvester (or SpiderFoot, or OpenCTI) in CI. Hit real endpoints. Catch shape "
             "drift before it ships."),
        ],
        body_size=14, sub_size=12,
    )
    add_notes(
        s,
        "Concrete, painful, valuable. Six bugs is enough to make the point. The pattern: third-party tool APIs drift, documentation lags, and mocks based on documentation are confidence-without-verification. "
        "Recommendation for OpenCTI builders: run an OpenCTI instance in your CI environment and exercise your Claude-integration tests against it. Painful upfront, pays back constantly."
    )
    add_footer(s)
    add_pagenum(s, 21, TOTAL)

    # ===================================================================
    # Slide 22 — Context isolation
    # ===================================================================
    s = blank_slide(prs)
    add_title(s, "Lesson 3 — Context isolation reduces error rate", subtitle="Counterintuitive: less context per agent = lower error rate.")
    add_bullets(
        s,
        [
            ("Orchestrator never reads raw articles",
             "The orchestrator schedules and delegates. The collector reads articles; the grader reads "
             "raw signal; the briefer reads graded findings. Each step's context is minimal."),
            ("Briefer never sees the coverage log",
             "Anti-repetition logic runs in the orchestrator and gets passed as 'these topics are stale.' "
             "Briefer composes against a clean slate per its scope."),
            ("Red-team-analyst never wrote the primary finding",
             "It's contrarian by construction because it has no investment in the leading hypothesis. "
             "The architecture creates the conditions for genuine disagreement."),
            ("Librarian is the only writer to side effects",
             "git, Splunk, Discord — all funnel through one agent. Easy to audit; easy to add rate-limiting; "
             "easy to enforce LEGAL-POLICY content-scan before any external publication."),
            ("Practical implication for your OpenCTI build",
             "Each Claude-Code-style subagent corresponds to a service or worker in your architecture. "
             "Don't try to do all of CTI in one big prompt — split by role, give each minimal context, "
             "and have one writer to the OpenCTI database."),
        ],
        body_size=14, sub_size=12,
    )
    add_notes(
        s,
        "Surprising result that most teams miss. The instinct is 'give the model everything it might need.' The opposite works better — minimal scope per agent, with the orchestrator handling cross-agent state. "
        "Land hard on the practical implication for OpenCTI: their architecture should map subagents to services or worker queues, not to one big prompt."
    )
    add_footer(s)
    add_pagenum(s, 22, TOTAL)

    # ===================================================================
    # Slide 23 — What you can build
    # ===================================================================
    s = blank_slide(prs)
    add_title(s, "What you can build", subtitle="Archimedes shape on OpenCTI + Teams. ~6-week MVP if OpenCTI is already standing.")
    add_bullets(
        s,
        [
            ("Week 1 — Doctrine + scaffolding",
             "Write the 5-6 doctrine files (LEGAL-POLICY, INTEL-GRADING, INTEL-BRIEF-STANDARDS, "
             "THREAT-BOX-METHODOLOGY, RETRACTION-POLICY, FLASH-POLICY). Set up the OpenCTI connector "
             "patterns for your top 5 OSINT sources."),
            ("Week 2-3 — Subagents",
             "Collector, grader, briefer, librarian to start. Each as a Claude subagent or as a service "
             "calling Claude. Test against a small set of canned findings first."),
            ("Week 4 — Delivery + audit",
             "Teams webhook + Adaptive Cards. Splunk (or equivalent) for telemetry. Git for doctrine "
             "versioning. Get the daily-cadence rhythm working end-to-end."),
            ("Week 5 — Threat actor + vuln tracking",
             "actor-profiler + vuln-tracker subagents. Threat-box scoring methodology. /approve-scoring "
             "gate posting to Teams for human review."),
            ("Week 6 — Hardening",
             "Red-team-analyst subagent. Retraction handling. Hard Rule refusals at every tool wrapper. "
             "Live-validate every external integration before production cutover."),
            ("Then iterate",
             "Real-world findings will surface gaps. Doctrine evolves. Add MCPs / connectors as needed. "
             "The first quarter is when the system earns its place in the intel team's workflow."),
        ],
        body_size=14, sub_size=11,
    )
    add_notes(
        s,
        "Realistic timeline assuming OpenCTI is already running and you have Claude API access. The doctrine work is the underestimated piece — it's not just markdown, it's the agent's job description. "
        "Don't skip the live validation in week 6. It's where the 6 bugs from lesson 2 surface in your own stack."
    )
    add_footer(s)
    add_pagenum(s, 23, TOTAL)

    # ===================================================================
    # Slide 24 — Q&A / Contact
    # ===================================================================
    s = blank_slide(prs)

    box = s.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(12.33), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Questions?"
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = C_TITLE
    run.font.name = "Calibri"

    box2 = s.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.33), Inches(0.5))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Happy to dig into any layer — architecture, doctrine, specific subagent, OpenCTI mapping."
    run2.font.size = Pt(16)
    run2.font.italic = True
    run2.font.color.rgb = C_SUBTLE
    run2.font.name = "Calibri"

    box3 = s.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.33), Inches(0.5))
    tf3 = box3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "[Presenter name]   ·   [Email / Slack / Teams handle]"
    run3.font.size = Pt(14)
    run3.font.color.rgb = C_SUBTLE
    run3.font.name = "Calibri"

    add_notes(
        s,
        "Likely questions to be ready for: "
        "(1) cost — Claude API token usage at this volume, ballpark $X/month; "
        "(2) what doesn't work yet — the dashboard (deferred indefinitely), the threat-box gate hasn't fired in five tries (feature not bug); "
        "(3) team size — Archimedes is one operator + Claude; their build will scale with team size and OpenCTI complexity; "
        "(4) what would you do differently — probably build OpenCTI-on-day-1 instead of markdown corpus; we got the analytic patterns right but the data model is constrained by file-based storage."
    )
    add_footer(s)
    add_pagenum(s, 24, TOTAL)

    # ---------- save ----------
    out_path = Path(__file__).parent / "archimedes-overview.pptx"
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    path = build()
    print(f"Wrote: {path}")
    print(f"Size:  {path.stat().st_size:,} bytes")
