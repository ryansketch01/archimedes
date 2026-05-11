"""Athena Discovery Proposal — 22-min, 19-slide presentation.

For a mixed audience (software dev team + leadership/sponsors) in the
same room. Discovery framing — proposed approach with explicit
"validation needed" callouts, not a committed architecture. Output:
athena-discovery.pptx in the same directory.

Design intent:
  - Discovery proposal tone, NOT confident commitment
  - Every requirement bullet has "we know / we need to validate"
  - Names Mandiant Advantage TI by name (GTI)
  - CMDB stays generic ("your CMDB") with tooling footnote
  - Project name "Athena" used as placeholder — find-replaceable
  - Same stylistic neutrality as archimedes-overview for template-merge
  - Speaker notes 2-3 sentences per slide

Run:
  uv run python presentations/athena-discovery/build_deck.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------- style constants ----------

C_TITLE = RGBColor(0x1F, 0x3A, 0x5F)
C_ACCENT = RGBColor(0x4A, 0x90, 0xE2)
C_BODY = RGBColor(0x33, 0x33, 0x33)
C_SUBTLE = RGBColor(0x66, 0x66, 0x66)
C_BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
C_VALIDATE = RGBColor(0xC8, 0x7F, 0x0A)  # amber for "needs validation"
C_KNOWN = RGBColor(0x2E, 0x7D, 0x32)     # green for "we know"
C_DISCORD_BG = RGBColor(0x31, 0x33, 0x38)
C_DISCORD_TEXT = RGBColor(0xDC, 0xDD, 0xDE)


# ---------- helpers ----------


def add_title(slide, text: str, *, subtitle: str | None = None) -> None:
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(12.33)

    box = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = C_TITLE
    run.font.name = "Calibri"

    if subtitle:
        box2 = slide.shapes.add_textbox(left, Inches(1.0), width, Inches(0.5))
        tf2 = box2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(14)
        run2.font.italic = True
        run2.font.color.rgb = C_SUBTLE
        run2.font.name = "Calibri"


def add_bullets(
    slide,
    bullets: list[str | tuple[str, str]],
    *,
    left: float = 0.5,
    top: float = 1.7,
    width: float = 12.33,
    height: float = 5.0,
    body_size: int = 18,
    sub_size: int = 14,
    line_spacing: float = 1.15,
) -> None:
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True

    for idx, item in enumerate(bullets):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if idx > 0:
            p.space_before = Pt(8)

        if isinstance(item, str):
            run = p.add_run()
            run.text = f"• {item}"
            run.font.size = Pt(body_size)
            run.font.color.rgb = C_BODY
            run.font.name = "Calibri"
        else:
            headline, sub = item
            run = p.add_run()
            run.text = f"• {headline}"
            run.font.size = Pt(body_size)
            run.font.bold = True
            run.font.color.rgb = C_BODY
            run.font.name = "Calibri"

            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.LEFT
            p2.line_spacing = line_spacing
            run2 = p2.add_run()
            run2.text = f"   {sub}"
            run2.font.size = Pt(sub_size)
            run2.font.color.rgb = C_SUBTLE
            run2.font.italic = True
            run2.font.name = "Calibri"


def add_kvtable(
    slide,
    rows: list[tuple[str, str, str]],
    *,
    left: float = 0.5,
    top: float = 1.7,
    col_widths: tuple[float, float, float] = (4.0, 4.7, 3.6),
    row_h: float = 0.55,
    header: tuple[str, str, str] | None = None,
) -> None:
    """Three-column reference table.

    Used for requirement→feature→know-vs-validate and similar.
    Each row's 3rd cell color-codes green (we know) vs amber (validate).
    """
    if header:
        for ci, h in enumerate(header):
            x = Inches(left + sum(col_widths[:ci]))
            hb = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, Inches(top), Inches(col_widths[ci]), Inches(row_h)
            )
            hb.line.fill.background()
            hb.fill.solid()
            hb.fill.fore_color.rgb = C_TITLE
            tf = hb.text_frame
            tf.margin_top = Inches(0.1)
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = "  " + h
            r.font.size = Pt(11)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            r.font.name = "Calibri"

    for ri, row in enumerate(rows):
        y = Inches(top + (row_h if header else 0) + row_h * ri)
        bg = RGBColor(0xFF, 0xFF, 0xFF) if ri % 2 == 0 else RGBColor(0xF8, 0xF8, 0xF8)
        for ci, val in enumerate(row):
            x = Inches(left + sum(col_widths[:ci]))
            cb = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, Inches(col_widths[ci]), Inches(row_h)
            )
            cb.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            cb.line.width = Pt(0.5)
            cb.fill.solid()
            cb.fill.fore_color.rgb = bg
            tf = cb.text_frame
            tf.margin_top = Inches(0.08)
            tf.margin_left = Inches(0.1)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = "  " + val
            r.font.size = Pt(9)
            # Color the third column by content prefix
            if ci == 2:
                if val.lower().startswith(("we know", "known")):
                    r.font.color.rgb = C_KNOWN
                elif "validate" in val.lower() or "tbd" in val.lower():
                    r.font.color.rgb = C_VALIDATE
                else:
                    r.font.color.rgb = C_BODY
            else:
                r.font.color.rgb = C_BODY
            r.font.name = "Calibri"
            if ci == 0:
                r.font.bold = True


def add_code_block(
    slide,
    code: str,
    *,
    left: float = 0.5,
    top: float = 1.7,
    width: float = 12.33,
    height: float = 5.0,
    font_size: int = 11,
    discord: bool = False,
) -> None:
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_DISCORD_BG if discord else C_BG_LIGHT
    text_color = C_DISCORD_TEXT if discord else C_BODY

    tb = slide.shapes.add_textbox(
        Inches(left + 0.15),
        Inches(top + 0.1),
        Inches(width - 0.3),
        Inches(height - 0.2),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(code.split("\n")):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line if line else " "
        run.font.size = Pt(font_size)
        run.font.color.rgb = text_color
        run.font.name = "Consolas"


def add_notes(slide, notes: str) -> None:
    slide.notes_slide.notes_text_frame.text = notes


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_footer(slide, text: str = "Athena · Discovery Proposal · 2026") -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.color.rgb = C_SUBTLE
    run.font.italic = True
    run.font.name = "Calibri"


def add_pagenum(slide, n: int, total: int) -> None:
    box = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.3), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{n} / {total}"
    run.font.size = Pt(9)
    run.font.color.rgb = C_SUBTLE
    run.font.name = "Calibri"


def add_discovery_badge(slide) -> None:
    """Small amber badge in the top right marking the slide as Discovery (not commitment)."""
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.0), Inches(0.3), Inches(1.8), Inches(0.4)
    )
    badge.line.fill.background()
    badge.fill.solid()
    badge.fill.fore_color.rgb = C_VALIDATE
    tf = badge.text_frame
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "DISCOVERY"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    r.font.name = "Calibri"


# ---------- build ----------


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    TOTAL = 19

    # ==================================================================
    # Slide 1 — Title
    # ==================================================================
    s = blank_slide(prs)

    # Discovery banner
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.4))
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = C_VALIDATE
    btf = band.text_frame
    btf.margin_top = Inches(0.05)
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = "DISCOVERY PROPOSAL  —  for discussion, not commitment"
    br.font.size = Pt(12)
    br.font.bold = True
    br.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    br.font.name = "Calibri"

    box = s.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(12.33), Inches(1.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Athena"
    run.font.size = Pt(72)
    run.font.bold = True
    run.font.color.rgb = C_TITLE
    run.font.name = "Calibri"

    box2 = s.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(12.33), Inches(0.6))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "AI-enabled CTI enrichment + decision support"
    run2.font.size = Pt(22)
    run2.font.italic = True
    run2.font.color.rgb = C_SUBTLE
    run2.font.name = "Calibri"

    box3 = s.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.33), Inches(0.5))
    tf3 = box3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "[Presenter name]   ·   [Date]   ·   Audience: dev team + program sponsors"
    run3.font.size = Pt(14)
    run3.font.color.rgb = C_SUBTLE
    run3.font.name = "Calibri"

    add_notes(
        s,
        "Open by naming the audience explicitly — dev team here to evaluate buildability, sponsors here to evaluate whether to fund a Discovery Phase. "
        "Frame: this is a Discovery proposal, not a build commitment. We have a working architecture pattern (Archimedes) we'd anchor on, but we're asking for 2-3 weeks of structured working sessions to validate the requirements before sizing a build. "
        "The deck shows the proposed approach plus the open questions that need stakeholder input."
    )

    # ==================================================================
    # Slide 2 — The problem
    # ==================================================================
    s = blank_slide(prs)
    add_title(s, "The problem", subtitle="From the intake — verbatim distilled")

    # Quote-style block
    qbox = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.6), Inches(12.33), Inches(2.7))
    qbox.line.color.rgb = C_ACCENT
    qbox.line.width = Pt(1.5)
    qbox.fill.solid()
    qbox.fill.fore_color.rgb = C_BG_LIGHT
    qtf = qbox.text_frame
    qtf.margin_top = Inches(0.2)
    qtf.margin_left = Inches(0.3)
    qtf.margin_right = Inches(0.3)
    qtf.word_wrap = True
    qp = qtf.paragraphs[0]
    qp.alignment = PP_ALIGN.LEFT
    qr = qp.add_run()
    qr.text = (
        "“Threat intelligence analysts spend significant time manually reviewing "
        "and judging the trustworthiness and relevance of intelligence (especially from "
        "OpenCTI), then separately correlating it to our environment (assets in CMDB) "
        "and external risk context (e.g., Mandiant Advantage / GTI).”"
    )
    qr.font.size = Pt(20)
    qr.font.italic = True
    qr.font.color.rgb = C_BODY
    qr.font.name = "Calibri"

    qp2 = qtf.add_paragraph()
    qp2.space_before = Pt(14)
    qr2 = qp2.add_run()
    qr2.text = (
        "“This manual work slows triage, creates inconsistent confidence scoring "
        "and prioritization across analysts, and delays response actions for high-impact "
        "vulnerabilities and threats.”"
    )
    qr2.font.size = Pt(20)
    qr2.font.italic = True
    qr2.font.color.rgb = C_BODY
    qr2.font.name = "Calibri"

    # Three takeaway tags
    takeaways = [
        ("Inconsistent", "scoring varies analyst-to-analyst"),
        ("Slow", "manual correlation across OpenCTI + GTI + CMDB"),
        ("Reactive", "high-impact threats wait on analyst capacity"),
    ]
    for i, (k, v) in enumerate(takeaways):
        x = Inches(0.5 + i * 4.27)
        tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.8), Inches(4.0), Inches(1.5))
        tag.line.color.rgb = C_TITLE
        tag.line.width = Pt(1.5)
        tag.fill.solid()
        tag.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        ttf = tag.text_frame
        ttf.margin_top = Inches(0.2)
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.alignment = PP_ALIGN.CENTER
        tr = tp.add_run()
        tr.text = k
        tr.font.size = Pt(24)
        tr.font.bold = True
        tr.font.color.rgb = C_TITLE
        tr.font.name = "Calibri"

        tp2 = ttf.add_paragraph()
        tp2.alignment = PP_ALIGN.CENTER
        tr2 = tp2.add_run()
        tr2.text = v
        tr2.font.size = Pt(12)
        tr2.font.italic = True
        tr2.font.color.rgb = C_SUBTLE
        tr2.font.name = "Calibri"

    add_notes(
        s,
        "Read the problem statement out loud — most of the room hasn't seen this slide before. "
        "The three tags at the bottom are what we'll come back to throughout — Athena's job is to make these three things go away. "
        "Don't get into the solution yet; this slide is anchoring why we're here."
    )
    add_footer(s)
    add_pagenum(s, 2, TOTAL)

    # ==================================================================
    # Slide 3 — The cost
    # ==================================================================
    s = blank_slide(prs)
    add_title(s, "The cost of status quo", subtitle="What the manual workflow costs today — placeholders for your numbers")
    add_bullets(
        s,
        [
            ("Analyst-hours on triage, not analysis",
             "Most of an analyst's day is the rote 80% — re-reading what Mandiant + Unit 42 + MSTIC published, "
             "looking it up in OpenCTI, cross-checking CMDB by hand. The 20% only humans can do gets squeezed."),
            ("Inconsistent confidence across analysts",
             "Two analysts grading the same finding can land at 'likely' vs. 'almost certainly' depending on "
             "experience, day, mood. Leadership sees the inconsistency and discounts the output."),
            ("Delayed response on high-impact threats",
             "By the time a Mandiant FLASH gets correlated to your CMDB exposure + your GTI risk overlay, the "
             "patch window is often half gone. Time-to-action on a CVSS-10 should be hours, not days."),
            ("Prioritization drift across the org",
             "SOC, Vuln Mgmt, IR, and Leadership often have different lists of 'top threats this week' — "
             "because there's no shared, transparent prioritization model. Decisions don't align."),
            ("Stakeholder reporting is its own job",
             "Weekly summaries to each team are hand-written. Tailoring the same intel to four audiences "
             "costs hours that don't produce new analysis."),
        ],
        body_size=15,
        sub_size=12,
    )
    add_notes(
        s,
        "If you have concrete numbers (analyst-hours per week, time-to-triage averages, MTTR for high-impact threats), drop them in here BEFORE this meeting — leadership wants the specific cost. "
        "If you don't have those numbers yet, the Discovery Phase produces them as part of validation. Frame that explicitly during the talk."
    )
    add_footer(s)
    add_pagenum(s, 3, TOTAL)

    # ==================================================================
    # Slide 4 — Why this is solvable now
    # ==================================================================
    s = blank_slide(prs)
    add_title(s, "Why this is solvable now", subtitle="The pattern is proven; the substrate exists; the missing piece is the analyst layer.")
    add_bullets(
        s,
        [
            ("The pattern is proven elsewhere",
             "Archimedes — internal CTI agent — has been running this shape autonomously for several months. "
             "Twice-daily graded briefs, FLASH alerts, threat-actor tracking, full audit trail. Same problem "
             "class, different stack."),
            ("Substrate already exists in your environment",
             "OpenCTI for the entity store. Mandiant Advantage / GTI for the external risk overlay. CMDB for "
             "the asset inventory. SOAR for the playbook engine. The pieces are in place — they don't talk to "
             "each other yet."),
            ("Claude is the missing analyst layer",
             "Long-context reads (200 articles → 4 graded findings). Consistent framework application "
             "(Admiralty, WEP, MITRE). Strict-format output. Tool use across all your existing systems. "
             "Context-isolated subagents — one per role — reduce error rate by construction."),
            ("Discovery proves it for YOUR environment",
             "A 2-3 week structured discovery phase validates the integration paths, locks the requirements, "
             "and produces a sized build plan with confidence intervals. Then you decide whether to fund the build."),
        ],
        body_size=15,
        sub_size=12,
    )
    add_notes(
        s,
        "This is the 'we know what we're doing' slide for leadership. Three points: pattern proven, substrate in place, Claude is the right engine. "
        "If asked to demo Archimedes, you can show #intel-briefs Discord post or screen-share a brief markdown file. "
        "The final bullet pivots to the Discovery ask without committing to a build yet."
    )
    add_footer(s)
    add_pagenum(s, 4, TOTAL)

    # ==================================================================
    # Slide 5 — Requirements → working hypotheses
    # ==================================================================
    s = blank_slide(prs)
    add_title(s, "Requirements → working hypotheses",
              subtitle="Your spec, mapped to a proposed feature, with explicit \"we know / we need to validate\" sub-line per row")
    add_discovery_badge(s)

    rows = [
        ("Confidence grading",
         "Admiralty 6×6 + WEP, applied per source per finding",
         "We know — same methodology in Archimedes today"),
        ("Threat prioritization",
         "Weighted threat-box scoring; per-category composites",
         "Validate — weights tuned to YOUR stakeholder priorities"),
        ("CMDB exposure scoring",
         "Finding IOCs/CVEs → CMDB asset match → criticality overlay",
         "Validate — your CMDB tool + asset-criticality field"),
        ("GTI risk cross-reference",
         "Mandiant Advantage risk score overlay on every graded finding",
         "Validate — your GTI license tier + API access"),
        ("Media attention signal",
         "Tracked-news source set + signal-volume metric per topic",
         "We know — pattern exists; needs source weighting"),
        ("Interactive Q&A",
         "Slash commands in Teams: /cve /investigate /ioc-hunt",
         "Validate — Teams as delivery platform; webhook plumbing"),
        ("“What matters now” dashboard",
         "Real-time view: top threats, exposure heatmap, FLASH queue, gates",
         "Validate — hosting model + auth/SSO + per-team filters"),
        ("Automated response playbooks",
         "Agent → SOAR webhook → playbook fires → result back to corpus",
         "Validate — which SOAR (Tines / Torq / ServiceNow / other)"),
        ("Weekly trend summaries",
         "Per-team renderings of the same canonical intel; “so what” framing",
         "Validate — stakeholder team list + contact owners"),
        ("On-demand threat reports",
         "/report <topic> against OpenCTI + GTI corpus",
         "We know — Archimedes /investigate already does this shape"),
    ]
    add_kvtable(
        s, rows,
        left=0.5, top=1.8,
        col_widths=(3.2, 5.3, 4.3),
        row_h=0.42,
        header=("Requirement", "Proposed feature", "We know / need to validate"),
    )
    add_notes(
        s,
        "This is the load-bearing slide of the deck. Every requirement from the intake is mapped to a concrete proposed feature and labeled with what we already know vs. what needs Discovery validation. "
        "Green text = pattern proven in Archimedes, low risk. Amber text = needs your input to lock. "
        "This is also the place to invite questions from the room — anything they want to flag as 'we need more on this' is fair game."
    )
    add_footer(s)
    add_pagenum(s, 5, TOTAL)

    # ==================================================================
    # Slide 6 — Confidence grading methodology
    # ==================================================================
    s = blank_slide(prs)
    add_title(s, "Confidence grading — proposed methodology",
              subtitle="Admiralty 6×6 + Words of Estimative Probability. Transparent, doctrine-versioned, applied uniformly.")

    # Two-column compare: Reliability + Credibility
    add_bullets(
        s,
        [
            ("A — Completely reliable", "Mandiant, Unit 42, MSTIC, CrowdStrike, MITRE, CISA"),
            ("B — Usually reliable", "BleepingComputer, The Record, Krebs"),
            ("C — Fairly reliable", "Vendor blogs with mixed history; first-citation entrants"),
            ("D / E / F", "Unreliable / unknown / questionable — monitoring only"),
        ],
        left=0.5, top=1.7, width=6.0, height=3.2, body_size=14, sub_size=11,
    )
    add_bullets(
        s,
        [
            ("1 — Confirmed", "Multiple independent A/B-grade sources"),
            ("2 — Probably true", "One A-grade source + plausibility test passes"),
            ("3 — Possibly true", "Single source, cannot be confirmed or denied"),
            ("4 / 5 / 6", "Doubtful / improbable / cannot be judged"),
        ],
        left=6.8, top=1.7, width=6.0, height=3.2, body_size=14, sub_size=11,
    )

    # WEP strip
    wep_box = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.0), Inches(12.33), Inches(1.2)
    )
    wep_box.line.color.rgb = C_TITLE
    wep_box.line.width = Pt(1)
    wep_box.fill.solid()
    wep_box.fill.fore_color.rgb = RGBColor(0xF1, 0xF7, 0xFD)
    wtf = wep_box.text_frame
    wtf.margin_top = Inches(0.15)
    wtf.margin_left = Inches(0.2)
    wtf.word_wrap = True
    wp = wtf.paragraphs[0]
    wr = wp.add_run()
    wr.text = (
        "Words of Estimative Probability bound every forward claim:  "
        "almost certainly (≥95%)  ·  very likely (85–95%)  ·  likely (55–85%)  ·  "
        "roughly even chance (~50%)  ·  unlikely (15–45%)  ·  very unlikely (5–15%)  ·  remote (<5%)"
    )
    wr.font.size = Pt(12)
    wr.font.color.rgb = C_BODY
    wr.font.name = "Calibri"
    wp2 = wtf.add_paragraph()
    wp2.space_before = Pt(4)
    wr2 = wp2.add_run()
    wr2.text = "Hedging language banned by doctrine. “May / might / could” → fails preflight, regenerated until WEP-compliant."
    wr2.font.size = Pt(11)
    wr2.font.italic = True
    wr2.font.color.rgb = C_SUBTLE
    wr2.font.name = "Calibri"

    # Footer note
    fbox = s.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.5))
    ftf = fbox.text_frame
    ftf.word_wrap = True
    fp = ftf.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    fr = fp.add_run()
    fr.text = (
        "Single-source-veto rule: even an A-grade source caps the forward claim at “likely” on its own. "
        "Solves the analyst-to-analyst inconsistency problem by construction."
    )
    fr.font.size = Pt(12)
    fr.font.italic = True
    fr.font.color.rgb = C_SUBTLE
    fr.font.name = "Calibri"
    add_notes(
        s,
        "Frame: Admiralty is older than cyber — NATO has used it since the Cold War. WEP comes from a 1964 CIA paper. We're not inventing methodology; we're encoding existing methodology in doctrine the agent applies uniformly. "
        "The single-source-veto rule is the load-bearing answer to 'inconsistent confidence across analysts' — the agent literally cannot promote a single-source claim above 'likely' regardless of how strong it sounds."
    )
    add_footer(s)
    add_pagenum(s, 6, TOTAL)

    # ==================================================================
    # Slide 7 — Transparent prioritization
    # ==================================================================
    s = blank_slide(prs)
    add_title(s, "Transparent threat prioritization — proposed",
              subtitle="Threat-box scoring: per-category composites, weighted to overall HIGH / MEDIUM / LOW. Auditable. Doctrine-versioned.")
    add_discovery_badge(s)

    # Left: how it works
    add_bullets(
        s,
        [
            ("Per-category scoring", "Espionage · Supply Chain · Destructive · Disruptive · Cyber-Crime · Influence. Each scored 0-10 against your environment."),
            ("Evidence-minimum tables", "Hard rules per score band. Agent refuses to manufacture HIGH from thin evidence. “Intent=5” requires named victim — not extrapolation."),
            ("Weighted to overall", "Default weights tunable per your stakeholder priorities. Espionage 35% / Supply 20% / Destructive 15% / Disruptive 15% / Cyber-Crime 15% as a starting point."),
            ("Human approval on HIGH", "Doctrine gate. Agent never auto-commits a HIGH score. Posts to review channel; waits for sign-off."),
            ("Empirically tight gate", "Archimedes ran 5 scorings; zero auto-HIGH. Gate refused to invent confidence. Same discipline carries over."),
        ],
        left=0.5, top=1.8, width=6.2, height=4.5, body_size=14, sub_size=11,
    )

    # Right: sample scoring
    sb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(1.8), Inches(5.8), Inches(4.5))
    sb.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    sb.line.width = Pt(0.5)
    sb.fill.solid()
    sb.fill.fore_color.rgb = C_BG_LIGHT
    stf = sb.text_frame
    stf.margin_top = Inches(0.2)
    stf.margin_left = Inches(0.2)
    stf.margin_right = Inches(0.2)
    stf.word_wrap = True

    lines = [
        ("Example: APT37 threat-box scoring", True, 15, C_TITLE, "Calibri"),
        ("", False, 12, C_BODY, "Calibri"),
        ("Espionage          composite 8/10  →  HIGH", False, 11, C_BODY, "Consolas"),
        ("Supply Chain       composite 6/10  →  MEDIUM", False, 11, C_BODY, "Consolas"),
        ("Destructive        composite 2/10  →  LOW", False, 11, C_BODY, "Consolas"),
        ("Disruptive         composite 2/10  →  LOW", False, 11, C_BODY, "Consolas"),
        ("Cyber-Crime        composite 2/10  →  LOW", False, 11, C_BODY, "Consolas"),
        ("", False, 12, C_BODY, "Calibri"),
        ("Weighted overall: 4.9  →  MEDIUM", True, 13, C_TITLE, "Calibri"),
        ("", False, 12, C_BODY, "Calibri"),
        ("Rationale — why not HIGH:", True, 11, C_BODY, "Calibri"),
        ("Intent capped at 3 (Sector Association) — public", False, 10, C_SUBTLE, "Calibri"),
        ("reporting names think tanks / journalists, no", False, 10, C_SUBTLE, "Calibri"),
        ("A&D primes. Methodology refused to extrapolate.", False, 10, C_SUBTLE, "Calibri"),
    ]
    first = True
    for text, bold, size, color, font in lines:
        if first:
            p = stf.paragraphs[0]
            first = False
        else:
            p = stf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font

    add_notes(
        s,
        "Critical for the inconsistency-across-analysts problem: the methodology is mechanical. Every analyst hitting this with the same evidence gets the same answer. "
        "The example on the right shows the discipline working — APT37 has espionage capability that would justify HIGH if you only looked at one category, but the four floor scores dilute to overall MEDIUM. "
        "Sponsors will ask 'what if we disagree with the weights?' — answer: the weights are doctrine, versioned in git, tunable in Discovery Phase per your priorities."
    )
    add_footer(s)
    add_pagenum(s, 7, TOTAL)

    # ==================================================================
    # Slide 8 — CMDB + GTI integration
    # ==================================================================
    s = blank_slide(prs)
    add_title(s, "CMDB exposure + GTI cross-reference — proposed shape",
              subtitle="The bridge between external intel and your environment. The piece OpenCTI alone doesn't give you.")
    add_discovery_badge(s)

    # Pipeline diagram
    boxes = [
        ("Graded finding", "OpenCTI report\n+ Admiralty digraph\n+ WEP", C_ACCENT, RGBColor(0xE3, 0xEF, 0xFA)),
        ("IOC / CVE\nextraction", "Pull indicators\n+ CVE refs from\nthe finding body", C_TITLE, RGBColor(0xE8, 0xEF, 0xF5)),
        ("CMDB asset match", "Match against\nyour asset inventory\n+ criticality field", C_TITLE, RGBColor(0xE8, 0xEF, 0xF5)),
        ("GTI risk overlay", "Mandiant Advantage\nrisk score for the\nCVE / actor / TTP", C_TITLE, RGBColor(0xE8, 0xEF, 0xF5)),
        ("Prioritized output", "Per-team brief\n+ dashboard tile\n+ playbook trigger", C_KNOWN, RGBColor(0xE6, 0xF4, 0xEA)),
    ]
    box_w = Inches(2.3)
    box_h = Inches(1.6)
    gap = Inches(0.2)
    total_w = box_w * 5 + gap * 4
    start_x = (Inches(13.333) - total_w) / 2
    y = Inches(2.2)

    for i, (title, body, color, fill) in enumerate(boxes):
        x = start_x + (box_w + gap) * i
        sb_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        sb_box.line.color.rgb = color
        sb_box.line.width = Pt(1.5)
        sb_box.fill.solid()
        sb_box.fill.fore_color.rgb = fill
        tf = sb_box.text_frame
        tf.margin_top = Inches(0.15)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = color
        r.font.name = "Calibri"

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(4)
        r2 = p2.add_run()
        r2.text = body
        r2.font.size = Pt(9)
        r2.font.color.rgb = C_BODY
        r2.font.name = "Calibri"

        # Arrow to next box
        if i < 4:
            ax = x + box_w
            ay = y + box_h / 2
            arrow = s.shapes.add_connector(1, ax, ay, ax + gap, ay)
            arrow.line.color.rgb = C_BODY
            arrow.line.width = Pt(1.5)

    # Validation callouts beneath
    callouts = [
        "Validate: CMDB tool + API access\n(ServiceNow / Tanium / Axonius / other)",
        "Validate: asset-criticality field schema\n(score vs. tag vs. owner-tier)",
        "Validate: GTI license tier — read API\nfor risk scores at finding-level granularity",
    ]
    co_y = Inches(4.6)
    co_w = Inches(4.0)
    co_h = Inches(1.0)
    co_gap = Inches(0.15)
    total_co_w = co_w * 3 + co_gap * 2
    co_start_x = (Inches(13.333) - total_co_w) / 2
    for i, txt in enumerate(callouts):
        cx = co_start_x + (co_w + co_gap) * i
        cb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, co_y, co_w, co_h)
        cb.line.color.rgb = C_VALIDATE
        cb.line.width = Pt(1.5)
        cb.fill.solid()
        cb.fill.fore_color.rgb = RGBColor(0xFE, 0xF6, 0xE7)
        tf = cb.text_frame
        tf.margin_top = Inches(0.12)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(10)
        r.font.color.rgb = C_BODY
        r.font.name = "Calibri"

    add_notes(
        s,
        "This is the slide where 'AI in CTI' becomes specific to YOUR environment. OpenCTI alone gives you the external intel; CMDB alone gives you the asset inventory; GTI alone gives you a risk overlay. "
        "Athena connects all three so 'CVE-2026-XXXX is hot' becomes '23 of YOUR assets are exposed; 4 are crown-jewel-tier; GTI rates this an 80; here's your prioritized list.' "
        "The amber callouts at the bottom are the explicit Discovery questions — three integration paths we need to validate before sizing the build."
    )
    add_footer(s)
    add_pagenum(s, 8, TOTAL)

    # ==================================================================
    # Slide 9 — Stakeholder-tailored output
    # ==================================================================
    s = blank_slide(prs)
    add_title(s, "Stakeholder-tailored output — proposed matrix",
              subtitle="One canonical record. Many renderings. Each team gets what they need, no more.")
    add_discovery_badge(s)

    matrix = [
        ("SOC", "IOCs, Splunk hunt queries, detection-tier signal", "Daily Layer-2 + ad-hoc FLASH"),
        ("Vuln Mgmt", "CVE deadlines, KEV status, CMDB exposure list, patch posture", "Weekly summary + FLASH on KEV adds"),
        ("Incident Response", "Actor TTPs, MITRE mapping, playbook handoffs, lessons", "Weekly + on-demand /investigate"),
        ("Threat Intel team", "Full analyst-grade markdown, source grading detail, doctrine audit", "Daily Layer-1 + weekly synthesis"),
        ("Leadership", "Executive Layer-2: top threats, exposure heatmap, posture summary", "Weekly executive brief"),
    ]
    add_kvtable(
        s, matrix,
        left=0.5, top=1.8,
        col_widths=(2.5, 7.5, 2.8),
        row_h=0.65,
        header=("Stakeholder team", "What they get", "Cadence"),
    )

    # Footer note
    fbox = s.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12.33), Inches(0.8))
    ftf = fbox.text_frame
    ftf.word_wrap = True
    fp = ftf.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    fr = fp.add_run()
    fr.text = (
        "Validation needed: confirm the team list above matches your org. Confirm contact owner per team "
        "(who signs off on the per-team format). Confirm cadence expectations with each team lead before build."
    )
    fr.font.size = Pt(12)
    fr.font.italic = True
    fr.font.color.rgb = C_VALIDATE
    fr.font.name = "Calibri"

    add_notes(
        s,
        "'Prioritization drift across the org' is solved here. Every team sees the same canonical findings, rendered for their workflow. Vuln Mgmt sees CVE deadlines; SOC sees Splunk hunt queries; Leadership sees the executive heatmap. "
        "Same source intel; same prioritization model; different surface. The discipline is in the doctrine — INTEL-BRIEF-STANDARDS specifies the per-team rendering rules."
    )
    add_footer(s)
    add_pagenum(s, 9, TOTAL)

    # ==================================================================
    # Slide 10 — Dashboard sketch
    # ==================================================================
    s = blank_slide(prs)
    add_title(s, "“What matters now” dashboard — wireframe sketch",
              subtitle="For discussion. The dashboard rendering surface; refined in Discovery.")
    add_discovery_badge(s)

    # Wireframe: 4-quadrant dashboard mock
    db_left = Inches(0.5)
    db_top = Inches(1.8)
    db_w = Inches(12.33)
    db_h = Inches(4.5)

    # Outer frame
    frame = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, db_left, db_top, db_w, db_h)
    frame.line.color.rgb = C_TITLE
    frame.line.width = Pt(2)
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFC)

    # Header bar
    hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, db_left, db_top, db_w, Inches(0.4))
    hdr.line.fill.background()
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = C_TITLE
    htf = hdr.text_frame
    htf.margin_top = Inches(0.06)
    htf.margin_left = Inches(0.3)
    hp = htf.paragraphs[0]
    hr = hp.add_run()
    hr.text = "Athena  ·  What matters now  ·  Monday, May 11 — 08:14 EDT"
    hr.font.size = Pt(11)
    hr.font.bold = True
    hr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    hr.font.name = "Calibri"

    # 4 quadrants
    quads = [
        ("TOP THREATS TODAY",
         "1.  CVE-2026-0300 (PAN-OS) · A1 · 23 exposed assets\n"
         "2.  Checkmarx Jenkins AST · B2 · DIB CI/CD impact\n"
         "3.  Ivanti EPMM CVE-2026-6973 · KEV expired\n"
         "4.  SailPoint GitHub repos · B2 · IGA stack\n"
         "5.  Operation HookedWing · C3 · monitoring"),
        ("YOUR EXPOSURE HEATMAP",
         "Critical:    4 assets   ████ \n"
         "High:        18 assets  ████████ \n"
         "Medium:      37 assets  ██████ \n"
         "Low:         142 assets ███ \n\n"
         "[hover: drilldown to CMDB CI list]"),
        ("FLASH QUEUE + PENDING GATES",
         "FLASH queued:    2 (quiet-hours)\n"
         "/approve-scoring waiting:  1 (UNC1549 → HIGH proposed)\n"
         "Retraction proposed:  0\n"
         "Source-grade ratify: 2 (Rapid7, SOCRadar)\n\n"
         "→ click any item to action"),
        ("RECENT ACTIVITY",
         "08:00  Morning brief published  ·  5 findings\n"
         "07:30  Pre-brief collection  ·  47 raw items\n"
         "06:00  FLASH sweep  ·  clean\n"
         "00:00  FLASH sweep  ·  clean\n"
         "Yesterday 16:00  Afternoon brief  ·  3 findings"),
    ]
    qw = Inches(6.0)
    qh = Inches(1.95)
    positions = [
        (Inches(0.65), Inches(2.3)),
        (Inches(6.8), Inches(2.3)),
        (Inches(0.65), Inches(4.3)),
        (Inches(6.8), Inches(4.3)),
    ]
    for (qtitle, qbody), (qx, qy) in zip(quads, positions):
        qb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, qx, qy, qw, qh)
        qb.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        qb.line.width = Pt(0.75)
        qb.fill.solid()
        qb.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf = qb.text_frame
        tf.margin_top = Inches(0.1)
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = qtitle
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = C_TITLE
        r.font.name = "Calibri"

        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        r2 = p2.add_run()
        r2.text = qbody
        r2.font.size = Pt(8)
        r2.font.color.rgb = C_BODY
        r2.font.name = "Consolas"

    # Footer callouts
    fbox = s.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.5))
    ftf = fbox.text_frame
    ftf.word_wrap = True
    fp = ftf.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    fr = fp.add_run()
    fr.text = (
        "Validation needed: hosting model (cloud / on-prem)  ·  auth/SSO integration  ·  per-team filters  ·  drilldown depth into CMDB"
    )
    fr.font.size = Pt(11)
    fr.font.italic = True
    fr.font.color.rgb = C_VALIDATE
    fr.font.name = "Calibri"

    add_notes(
        s,
        "Wireframe — for discussion, NOT a committed design. Four quadrants address four different consumption modes: 'what's hot' (top threats), 'where am I exposed' (heatmap), 'what needs my attention' (FLASH+gates), 'what just happened' (activity log). "
        "Sponsors will want to see this — leadership consumes the dashboard primarily. Dev team will ask about hosting / auth / data refresh; that's all Discovery work. "
        "Note this is also where 'inconsistent prioritization across the org' gets visually resolved — everybody sees the same top-5 list."
    )
    add_footer(s)
    add_pagenum(s, 10, TOTAL)

    # ==================================================================
    # Slide 11 — Interactive Q&A
    # ==================================================================
    s = blank_slide(prs)
    add_title(s, "Interactive Q&A — proposed slash commands",
              subtitle="Analyst asks a question in Teams; agent investigates against OpenCTI + GTI + the corpus; returns analyst-grade response.")
    add_discovery_badge(s)

    sample = """[Analyst]  /investigate UNC1549

[Athena]
UNC1549 — Iranian-cluster espionage actor (a.k.a. Smoke Sandstorm).
Tracking since 2024. Last 90 days: 3 named A&D victims (Mandiant
2026-04-12, CrowdStrike 2026-04-28, MSTIC 2026-05-02).

Threat-box scoring: HIGH (5.4 weighted) — Espionage HIGH (10),
Supply Chain LOW, Destructive LOW, Disruptive LOW, Cyber-Crime LOW.
Last /update-tracking: 2026-05-09.

Current exposure in your environment:
  • 0 confirmed IOC hits in Splunk -30d
  • 4 CVEs tied to UNC1549 tradecraft; 2 unpatched on 11 assets
  • GTI risk score: 84 (threat-actor entity)

Action items proposed:
  • Patch CVE-2026-XXXX on the 11 exposed assets (Vuln Mgmt)
  • Confirm detection coverage for the named TTPs (SOC)
  • Review UNC1549 dossier for any sector-specific update (Intel team)

Sources cited: Mandiant 2026-04-12, CrowdStrike 2026-04-28, MSTIC
2026-05-02, internal /update-tracking 2026-05-09.

[Analyst]  /report UNC1549 last 30 days

[Athena]
Generating threat report from OpenCTI + GTI corpus...
[returns Layer-1 markdown brief with full source citations]
"""
    add_code_block(s, sample, top=1.8, height=5.0, font_size=11)
    add_notes(
        s,
        "Show what 'ask the agent a question' actually looks like. The shape is borrowed from Archimedes' Discord listener — slash commands, structured response, source citations. "
        "Three things to emphasize: (1) the agent investigates LIVE against OpenCTI + GTI, not just the cached corpus; (2) it returns the same analyst-grade format the briefs use, not a chatbot ramble; (3) every claim cites a source. "
        "If the audience asks 'what if the agent gets it wrong' — the answer is single-source-veto + WEP discipline + every refusal logged for review."
    )
    add_footer(s)
    add_pagenum(s, 11, TOTAL)

    # ==================================================================
    # Slide 12 — Automated playbook integration
    # ==================================================================
    s = blank_slide(prs)
    add_title(s, "Automated response playbooks — proposed integration shape",
              subtitle="Agent posts; your SOAR orchestrates; result feeds back. Two webhooks each direction.")
    add_discovery_badge(s)

    # Sequence diagram
    actors = [("Athena", C_TITLE), ("SOAR (your tool)", C_ACCENT), ("Splunk + CMDB", C_KNOWN), ("Athena corpus", C_TITLE)]
    actor_y = Inches(1.9)
    actor_w = Inches(2.8)
    actor_h = Inches(0.5)
    actor_gap = Inches(0.4)
    actor_total_w = actor_w * 4 + actor_gap * 3
    actor_start = (Inches(13.333) - actor_total_w) / 2
    for i, (name, color) in enumerate(actors):
        x = actor_start + (actor_w + actor_gap) * i
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, actor_y, actor_w, actor_h)
        box.line.fill.background()
        box.fill.solid()
        box.fill.fore_color.rgb = color
        tf = box.text_frame
        tf.margin_top = Inches(0.1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = name
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        r.font.name = "Calibri"

    # Steps below
    steps = [
        "1. FLASH triggers — Athena graded CVE-2026-XXXX as A1 + active exploitation + 4 exposed assets",
        "2. Athena POSTs webhook to SOAR with finding-id, IOCs, CMDB asset list, suggested playbook",
        "3. SOAR runs playbook — quarantine VLAN, force patch, page on-call, open ServiceNow ticket",
        "4. SOAR POSTs result back to Athena — playbook id, action results, ticket ref",
        "5. Athena writes outcome to corpus — finding now carries response-action provenance",
        "6. Next brief surfaces the closed-loop event under 'Recent activity / closed actions'",
    ]
    step_box = s.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.33), Inches(3.5))
    stf = step_box.text_frame
    stf.word_wrap = True
    for i, step in enumerate(steps):
        if i == 0:
            p = stf.paragraphs[0]
        else:
            p = stf.add_paragraph()
        p.space_before = Pt(6)
        p.line_spacing = 1.3
        r = p.add_run()
        r.text = step
        r.font.size = Pt(13)
        r.font.color.rgb = C_BODY
        r.font.name = "Calibri"

    # Validation footer
    fb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.3), Inches(12.33), Inches(0.6))
    fb.line.color.rgb = C_VALIDATE
    fb.line.width = Pt(1.5)
    fb.fill.solid()
    fb.fill.fore_color.rgb = RGBColor(0xFE, 0xF6, 0xE7)
    ftf = fb.text_frame
    ftf.margin_top = Inches(0.12)
    fp = ftf.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    fr = fp.add_run()
    fr.text = (
        "Validation needed: which SOAR (Tines / Torq / ServiceNow SecOps / Splunk SOAR / other), playbook catalog, "
        "auth/auth model for the webhooks, change-management approval path for automated actions."
    )
    fr.font.size = Pt(11)
    fr.font.italic = True
    fr.font.color.rgb = C_BODY
    fr.font.name = "Calibri"

    add_notes(
        s,
        "The bidirectional webhook pattern is what closes the loop on 'delays response actions for high-impact threats.' Without it, the agent just produces graded intel; with it, the intel triggers actual response. "
        "Critical caveat for sponsors: automated response on production assets is high-stakes. Discovery Phase needs to lock down which playbooks are agent-triggerable vs. analyst-confirmed-then-triggered. "
        "Pattern: low-risk actions auto-fire (quarantine, ticket open); destructive actions require human approval (block, isolate, take down)."
    )
    add_footer(s)
    add_pagenum(s, 12, TOTAL)

    # ==================================================================
    # Slide 13 — Architecture sketch
    # ==================================================================
    s = blank_slide(prs)
    add_title(s, "Architecture sketch — proposed",
              subtitle="Starting point. Component layout. Discovery validates the integration paths.")
    add_discovery_badge(s)

    # Reusing the architecture pattern from the overview deck
    orch = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.7), Inches(1.7), Inches(2.0), Inches(0.6)
    )
    orch.line.color.rgb = C_TITLE
    orch.line.width = Pt(2)
    orch.fill.solid()
    orch.fill.fore_color.rgb = C_TITLE
    otf = orch.text_frame
    otf.margin_top = Inches(0.08)
    op = otf.paragraphs[0]
    op.alignment = PP_ALIGN.CENTER
    orun = op.add_run()
    orun.text = "Claude (Athena)"
    orun.font.size = Pt(13)
    orun.font.bold = True
    orun.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    orun.font.name = "Calibri"

    # Left — Inputs
    inp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.7), Inches(2.6), Inches(0.6))
    inp.line.color.rgb = C_ACCENT
    inp.line.width = Pt(1.5)
    inp.fill.solid()
    inp.fill.fore_color.rgb = C_BG_LIGHT
    itf = inp.text_frame
    itf.margin_top = Inches(0.08)
    ip = itf.paragraphs[0]
    ip.alignment = PP_ALIGN.CENTER
    ir = ip.add_run()
    ir.text = "Intel inputs"
    ir.font.size = Pt(12)
    ir.font.bold = True
    ir.font.color.rgb = C_ACCENT
    ir.font.name = "Calibri"

    isub = s.shapes.add_textbox(Inches(0.5), Inches(2.35), Inches(2.6), Inches(0.5))
    istf = isub.text_frame
    istf.word_wrap = True
    isp = istf.paragraphs[0]
    isp.alignment = PP_ALIGN.CENTER
    isr = isp.add_run()
    isr.text = "OpenCTI · Mandiant GTI\nRSS · vendor blogs · Splunk"
    isr.font.size = Pt(9)
    isr.font.color.rgb = C_SUBTLE
    isr.font.name = "Calibri"

    # Right — Outputs
    outb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.2), Inches(1.7), Inches(2.6), Inches(0.6))
    outb.line.color.rgb = C_KNOWN
    outb.line.width = Pt(1.5)
    outb.fill.solid()
    outb.fill.fore_color.rgb = C_BG_LIGHT
    otf2 = outb.text_frame
    otf2.margin_top = Inches(0.08)
    op2 = otf2.paragraphs[0]
    op2.alignment = PP_ALIGN.CENTER
    orun2 = op2.add_run()
    orun2.text = "Outputs"
    orun2.font.size = Pt(12)
    orun2.font.bold = True
    orun2.font.color.rgb = C_KNOWN
    orun2.font.name = "Calibri"

    osub = s.shapes.add_textbox(Inches(10.2), Inches(2.35), Inches(2.6), Inches(0.5))
    ostf = osub.text_frame
    ostf.word_wrap = True
    osp = ostf.paragraphs[0]
    osp.alignment = PP_ALIGN.CENTER
    osr = osp.add_run()
    osr.text = "Teams · Dashboard\nSOAR webhooks · audit log"
    osr.font.size = Pt(9)
    osr.font.color.rgb = C_SUBTLE
    osr.font.name = "Calibri"

    # Arrows
    al = s.shapes.add_connector(1, Inches(3.1), Inches(2.0), Inches(5.7), Inches(2.0))
    al.line.color.rgb = C_BODY
    al.line.width = Pt(1.5)
    ar = s.shapes.add_connector(1, Inches(7.7), Inches(2.0), Inches(10.2), Inches(2.0))
    ar.line.color.rgb = C_BODY
    ar.line.width = Pt(1.5)

    # Subagent ring
    subagents = [
        ("collector", "ingest + filter"),
        ("grader", "Admiralty + WEP"),
        ("analyst", "SAT / ACH / KAC"),
        ("red-team", "challenge HIGH"),
        ("actor-profiler", "dossiers + scoring"),
        ("vuln-tracker", "CVE corpus + KEV"),
        ("briefer", "per-team rendering"),
        ("librarian", "ship + audit + log"),
    ]
    ring_top = Inches(3.5)
    box_w = Inches(1.45)
    box_h = Inches(0.95)
    spacing = Inches(0.1)
    total_w = box_w * 8 + spacing * 7
    ring_left = (Inches(13.333) - total_w) / 2
    for i, (name, role) in enumerate(subagents):
        x = ring_left + (box_w + spacing) * i
        sb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, ring_top, box_w, box_h)
        sb.line.color.rgb = C_ACCENT
        sb.line.width = Pt(1.2)
        sb.fill.solid()
        sb.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        stf = sb.text_frame
        stf.margin_top = Inches(0.08)
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        sr = sp.add_run()
        sr.text = name
        sr.font.size = Pt(11)
        sr.font.bold = True
        sr.font.color.rgb = C_TITLE
        sr.font.name = "Calibri"

        sp2 = stf.add_paragraph()
        sp2.alignment = PP_ALIGN.CENTER
        sr2 = sp2.add_run()
        sr2.text = role
        sr2.font.size = Pt(8)
        sr2.font.italic = True
        sr2.font.color.rgb = C_SUBTLE
        sr2.font.name = "Calibri"

    # Integration band
    int_top = Inches(5.0)
    int_band = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), int_top, Inches(12.33), Inches(0.5)
    )
    int_band.line.color.rgb = C_VALIDATE
    int_band.line.width = Pt(1)
    int_band.fill.solid()
    int_band.fill.fore_color.rgb = RGBColor(0xFE, 0xF6, 0xE7)
    itf2 = int_band.text_frame
    itf2.margin_top = Inches(0.08)
    ip2 = itf2.paragraphs[0]
    ip2.alignment = PP_ALIGN.CENTER
    ir2 = ip2.add_run()
    ir2.text = (
        "Validation needed:  OpenCTI deployment status  ·  CMDB API path  ·  GTI license tier  ·  SOAR webhook capability  ·  Teams app/channel access"
    )
    ir2.font.size = Pt(11)
    ir2.font.bold = True
    ir2.font.color.rgb = C_BODY
    ir2.font.name = "Calibri"

    # Doctrine band
    doc_top = Inches(5.7)
    doc_band = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), doc_top, Inches(12.33), Inches(0.5)
    )
    doc_band.line.color.rgb = C_SUBTLE
    doc_band.line.width = Pt(1)
    doc_band.fill.solid()
    doc_band.fill.fore_color.rgb = C_BG_LIGHT
    dtf = doc_band.text_frame
    dtf.margin_top = Inches(0.08)
    dp = dtf.paragraphs[0]
    dp.alignment = PP_ALIGN.CENTER
    dr = dp.add_run()
    dr.text = (
        "Doctrine (versioned in git):  LEGAL-POLICY  ·  INTEL-GRADING  ·  INTEL-BRIEF-STANDARDS  ·  "
        "THREAT-BOX-METHODOLOGY  ·  RETRACTION-POLICY  ·  FLASH-POLICY"
    )
    dr.font.size = Pt(11)
    dr.font.bold = True
    dr.font.color.rgb = C_BODY
    dr.font.name = "Calibri"

    cap = s.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.5))
    ctf = cap.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = (
        "Each subagent: own context, own write scope, own doctrine subset. Orchestrator never reads raw articles. "
        "Librarian is the only writer to side effects (corpus, Teams, SOAR, audit log)."
    )
    cr.font.size = Pt(11)
    cr.font.italic = True
    cr.font.color.rgb = C_SUBTLE
    cr.font.name = "Calibri"

    add_notes(
        s,
        "This is the architecture for sponsors. Frame it as 'this is the starting point we'd anchor on; Discovery validates the integration paths.' "
        "The amber band lists the four integration unknowns explicitly — those are the questions we're asking sponsors to commit working-session time to resolve. "
        "If a dev in the room asks 'why subagents instead of one prompt' — answer: context isolation. Proven pattern; lower error rate; auditable per-role."
    )
    add_footer(s)
    add_pagenum(s, 13, TOTAL)

    # ==================================================================
    # Slide 14 — Doctrine as code
    # ==================================================================
    s = blank_slide(prs)
    add_title(s, "Doctrine as code — the working memory layer",
              subtitle="Versioned in git. Reviewable by analysts AND engineers. The agent's job description.")

    add_bullets(
        s,
        [
            ("LEGAL-POLICY.md", "Prohibited queries, authorized targets, ITAR boundary, copyright. Read before every tool call."),
            ("INTEL-GRADING.md", "Admiralty 6×6 + WEP + single-source-veto + corroboration rules. Read by grader."),
            ("INTEL-BRIEF-STANDARDS.md", "Per-stakeholder rendering rules, word/char budgets, Smart Brevity, preflight checklist. Read by briefer."),
            ("THREAT-BOX-METHODOLOGY.md", "Per-category composite scoring + evidence-minimum tables. Read by actor-profiler on every scoring run."),
            ("RETRACTION-POLICY.md", "When to retract vs. correct, additive (never silent), 72h auto-downgrade clock."),
            ("FLASH-POLICY.md", "Trigger conditions, quiet hours, critical-override threshold."),
        ],
        body_size=14, sub_size=11,
    )

    # Why it matters footer
    fb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.0), Inches(12.33), Inches(0.9))
    fb.line.color.rgb = C_KNOWN
    fb.line.width = Pt(1.5)
    fb.fill.solid()
    fb.fill.fore_color.rgb = RGBColor(0xE6, 0xF4, 0xEA)
    ftf = fb.text_frame
    ftf.margin_top = Inches(0.15)
    ftf.margin_left = Inches(0.2)
    ftf.margin_right = Inches(0.2)
    ftf.word_wrap = True
    fp = ftf.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    fr = fp.add_run()
    fr.text = (
        "Why this matters: doctrine is what the agent's job description used to be locked inside a single prompt. "
        "Splitting it into versioned files means analysts, engineers, and compliance can review and iterate it like code. "
        "It also makes “why did the agent do that” answerable — every behavior traces back to a doctrine line."
    )
    fr.font.size = Pt(12)
    fr.font.color.rgb = C_BODY
    fr.font.name = "Calibri"

    add_notes(
        s,
        "Doctrine-as-code is the architectural answer to 'how do we govern agent behavior in a regulated environment?' Every rule the agent follows is in a .md file. Every change is in git history. Every refusal is traceable. "
        "Engineers can review the structure; analysts can review the methodology. Compliance can audit both. This is the layer that makes the system defensible — not just performant. "
        "Discovery Phase includes a doctrine review with your intel team lead — they sign off on the methodology before we commit to building against it."
    )
    add_footer(s)
    add_pagenum(s, 14, TOTAL)

    # ==================================================================
    # Slide 15 — Discovery Phase scope
    # ==================================================================
    s = blank_slide(prs)
    add_title(s, "Discovery Phase — proposed scope",
              subtitle="2-3 weeks of structured work. Validated requirements doc, target architecture, sized build plan.")
    add_discovery_badge(s)

    # Week-by-week
    weeks = [
        ("Week 1",
         "Requirement validation + stakeholder interviews",
         "Working sessions with each stakeholder team (SOC, IR, Vuln Mgmt, Intel, Leadership). Confirm requirements, "
         "weight priorities, lock the per-team rendering format. Capture analyst-day shadow observations to ground the "
         "'cost of status quo' numbers."),
        ("Week 2",
         "Integration POCs + doctrine review",
         "Hands-on POCs for OpenCTI read, GTI API call, CMDB asset lookup, SOAR webhook. Each is a 1-day spike. "
         "Doctrine review with intel team lead — adopt-as-is, modify, or reject. Identify any deal-breakers early."),
        ("Week 3",
         "Synthesis + sized build plan",
         "Produce validated requirements doc, target architecture (no longer 'proposed'), build plan with phases + "
         "confidence intervals, integration spec for each external system. Decision gate: green-light build, "
         "request additional scoping, or stop."),
    ]
    add_kvtable(
        s, weeks,
        left=0.5, top=1.8,
        col_widths=(1.5, 4.0, 7.8),
        row_h=1.4,
        header=("When", "Focus", "Activities + outcomes"),
    )

    # Deliverables footer
    fb = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.5)
    )
    fb.line.color.rgb = C_KNOWN
    fb.line.width = Pt(1.5)
    fb.fill.solid()
    fb.fill.fore_color.rgb = RGBColor(0xE6, 0xF4, 0xEA)
    ftf = fb.text_frame
    ftf.margin_top = Inches(0.1)
    fp = ftf.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    fr = fp.add_run()
    fr.text = (
        "Deliverables at end of Discovery:  validated requirements doc  ·  target architecture  ·  sized build plan  ·  "
        "integration POC results  ·  stakeholder sign-off"
    )
    fr.font.size = Pt(11)
    fr.font.bold = True
    fr.font.color.rgb = C_BODY
    fr.font.name = "Calibri"

    add_notes(
        s,
        "Discovery is a contained scope with a hard decision gate at the end. Three weeks; named deliverables; either we have a sized build plan or we have a documented reason not to proceed. Either outcome is valuable. "
        "Push for sponsor commitment to attend working sessions personally — without their input on weights and per-team rendering, we'll build the wrong thing efficiently."
    )
    add_footer(s)
    add_pagenum(s, 15, TOTAL)

    # ==================================================================
    # Slide 16 — Open questions
    # ==================================================================
    s = blank_slide(prs)
    add_title(s, "Open questions to resolve in Discovery",
              subtitle="Every one of these blocks a build commitment. Discovery is how we close them.")
    add_discovery_badge(s)

    questions = [
        ("CMDB integration",
         "Which CMDB tool? API access available? Asset-criticality schema (score / tag / owner-tier)? "
         "Who owns the integration on your side?"),
        ("GTI / Mandiant Advantage license",
         "Which license tier? API access for risk-score-at-finding-level granularity? Rate-limit headroom for daily volume?"),
        ("SOAR / playbook engine",
         "Which tool (Tines / Torq / ServiceNow SecOps / Splunk SOAR / other)? Existing playbook catalog? "
         "Change-management path for automated actions on production?"),
        ("Stakeholder team list",
         "Confirm the team list (SOC, IR, Vuln Mgmt, Intel, Leadership). Contact owner per team. "
         "Cadence + format expectations per team."),
        ("OpenCTI deployment status",
         "Already running in production? Connector inventory? Reasonable read-load on current instance? "
         "STIX 2.x or older?"),
        ("Doctrine review",
         "Intel team lead available to review the 6 doctrine files? Modifications needed before the agent runs against them?"),
        ("Hosting + auth",
         "Where does the dashboard live (cloud / on-prem)? SSO integration path? Per-team access controls?"),
        ("Build team + timeline expectations",
         "FTE availability for the build phase? Target go-live date? Capacity for ongoing doctrine + corpus maintenance?"),
    ]
    add_bullets(s, questions, body_size=13, sub_size=11)
    add_notes(
        s,
        "Read these out loud. The point is to make explicit what's NOT yet known — sponsors should see that we've thought about what could go wrong and have a structured plan to surface answers. "
        "If a sponsor in the room can already answer one of these, capture it immediately — that's pre-Discovery progress. "
        "The deck deliberately doesn't pretend to have answers we don't. Honest scope wins."
    )
    add_footer(s)
    add_pagenum(s, 16, TOTAL)

    # ==================================================================
    # Slide 17 — The ask
    # ==================================================================
    s = blank_slide(prs)
    add_title(s, "The ask", subtitle="What we need from sponsors today to start Discovery next week.")

    asks = [
        ("Sponsor a 2-3 week Discovery Phase",
         "Named budget owner. Time commitment for working sessions. Clear scope as outlined."),
        ("Name working-session owners",
         "One per open dimension: CMDB, GTI, SOAR, Stakeholder list, OpenCTI, Doctrine. ~4-6 hours each across the phase."),
        ("Commit to the Discovery-end decision gate",
         "End of week 3: review validated requirements + target architecture + sized build plan. Decide: build / additional scoping / stop."),
        ("Confirm Athena as the working project name",
         "(or assign a different one). Used in working-session artifacts and the final build proposal."),
        ("Approve initial access requests for the POC integrations",
         "Read-only API access to OpenCTI, GTI, CMDB. Sandbox SOAR access. Teams app registration for the dashboard wireframe."),
    ]
    add_bullets(s, asks, body_size=15, sub_size=12)

    add_notes(
        s,
        "Five concrete asks. None of them commit to a build; all of them get us to a sized build plan. "
        "If sponsors push back on any specific ask — capture the resistance and route around it; don't argue. The point is to get to Discovery, not to win every line item. "
        "Closing line for this slide: 'with these five things, we can start Discovery next Monday and have a sized build plan in three weeks.'"
    )
    add_footer(s)
    add_pagenum(s, 17, TOTAL)

    # ==================================================================
    # Slide 18 — What Discovery delivers
    # ==================================================================
    s = blank_slide(prs)
    add_title(s, "What Discovery delivers",
              subtitle="End-of-phase deliverables. Then you decide whether to fund the build.")

    deliverables = [
        ("Validated requirements document",
         "Every requirement bullet from the intake, validated through working sessions, with explicit acceptance criteria and stakeholder sign-off."),
        ("Target architecture",
         "Component diagram with integration paths locked. No more 'proposed' — this is what we'd build. Reviewable, defensible, versionable."),
        ("Sized build plan with confidence intervals",
         "Phased delivery, FTE estimates, dependency map, risk register. “8 weeks ± 2” not “some amount of time.”"),
        ("Integration POC results",
         "OpenCTI read working. GTI risk-score lookup working. CMDB asset match working. SOAR webhook firing. Each demonstrated in the working sessions, not just claimed."),
        ("Doctrine sign-off",
         "Intel team lead has reviewed the 6 doctrine files and approved (or marked the modifications needed)."),
        ("Stakeholder sign-off",
         "Each stakeholder team owner has confirmed the proposed per-team rendering format and cadence."),
        ("A clear go/no-go decision",
         "Either: “we have what we need to build, here's the plan” — or — “here's the specific reason this isn't ready to build yet.” Both outcomes are valuable."),
    ]
    add_bullets(s, deliverables, body_size=13, sub_size=11)

    add_notes(
        s,
        "Close on the value of Discovery whether it leads to a build or not. A 'no, not yet' answer that surfaces three specific blockers is worth as much as a 'yes' — because we don't waste 8 weeks of build time on a project that wasn't ready. "
        "Both outcomes produce a defensible artifact. That's the discovery framing's strongest argument."
    )
    add_footer(s)
    add_pagenum(s, 18, TOTAL)

    # ==================================================================
    # Slide 19 — Q&A
    # ==================================================================
    s = blank_slide(prs)
    box = s.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(12.33), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Questions?"
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = C_TITLE
    run.font.name = "Calibri"

    box2 = s.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.33), Inches(0.5))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "And what you'd want to validate first in Discovery."
    run2.font.size = Pt(16)
    run2.font.italic = True
    run2.font.color.rgb = C_SUBTLE
    run2.font.name = "Calibri"

    box3 = s.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.33), Inches(0.5))
    tf3 = box3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "[Presenter name]   ·   [Email / Slack / Teams handle]"
    run3.font.size = Pt(14)
    run3.font.color.rgb = C_SUBTLE
    run3.font.name = "Calibri"

    add_notes(
        s,
        "Likely questions to be ready for: "
        "(1) Cost — Discovery is mostly your time + working-session hours; the build itself is the unknown until Discovery sizes it. "
        "(2) What if Discovery says 'no' — both outcomes produce a defensible artifact (see slide 18). "
        "(3) Why Claude vs. an in-house ML model — long-context + tool-use + framework-application discipline at this maturity is hard to build from scratch; Claude solves the analyst-layer problem cheaper and faster than a custom model would. "
        "(4) Security review — every external integration is auth-bounded; every action is logged; every doctrine refusal is auditable. Discovery includes a security review touchpoint."
    )
    add_footer(s)
    add_pagenum(s, 19, TOTAL)

    # ---------- save ----------
    out_path = Path(__file__).parent / "athena-discovery.pptx"
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    path = build()
    print(f"Wrote: {path}")
    print(f"Size:  {path.stat().st_size:,} bytes")
