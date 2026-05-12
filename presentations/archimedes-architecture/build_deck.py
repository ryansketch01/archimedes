"""Archimedes Architecture Walkthrough — 45-min, 32-slide deck.

For a senior engineer / architect audience that already knows Claude
(MCP, subagents, hooks) and needs to internalize Archimedes well
enough to rebuild it. Artifact-heavy: real file paths, real code
excerpts, real directory trees, real subagent inventory.

Design intent:
  - Engineering-handoff tone. Not marketing, not discovery.
  - Every slide has at least one real artifact from the codebase.
  - Diagrams at section openers; details in tables + code blocks.
  - Same stylistic neutrality as prior decks (Calibri, gray/blue
    accents, no custom backgrounds) so company template merges clean.
  - 32 slides @ ~85s = 45 min target.

Run:
  uv run python presentations/archimedes-architecture/build_deck.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------- style ----------

C_TITLE = RGBColor(0x1F, 0x3A, 0x5F)
C_ACCENT = RGBColor(0x4A, 0x90, 0xE2)
C_BODY = RGBColor(0x33, 0x33, 0x33)
C_SUBTLE = RGBColor(0x66, 0x66, 0x66)
C_BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
C_CODE_BG = RGBColor(0x2D, 0x30, 0x3D)
C_CODE_TEXT = RGBColor(0xE6, 0xE6, 0xE6)
C_CODE_KEYWORD = RGBColor(0x9C, 0xDC, 0xFE)
C_CODE_STRING = RGBColor(0xCE, 0x91, 0x78)
C_REFUSAL = RGBColor(0xC0, 0x39, 0x2B)
C_KNOWN = RGBColor(0x2E, 0x7D, 0x32)
C_HIGHLIGHT = RGBColor(0xC8, 0x7F, 0x0A)


# ---------- helpers ----------


def add_title(slide, text: str, *, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = C_TITLE
    run.font.name = "Calibri"

    if subtitle:
        b2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.33), Inches(0.4))
        tf2 = b2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(13)
        r2.font.italic = True
        r2.font.color.rgb = C_SUBTLE
        r2.font.name = "Calibri"


def add_bullets(
    slide,
    bullets: list[str | tuple[str, str]],
    *,
    left: float = 0.5,
    top: float = 1.6,
    width: float = 12.33,
    height: float = 5.2,
    body_size: int = 16,
    sub_size: int = 12,
    line_spacing: float = 1.15,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if idx > 0:
            p.space_before = Pt(6)
        if isinstance(item, str):
            r = p.add_run()
            r.text = f"• {item}"
            r.font.size = Pt(body_size)
            r.font.color.rgb = C_BODY
            r.font.name = "Calibri"
        else:
            headline, sub = item
            r = p.add_run()
            r.text = f"• {headline}"
            r.font.size = Pt(body_size)
            r.font.bold = True
            r.font.color.rgb = C_BODY
            r.font.name = "Calibri"
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.LEFT
            p2.line_spacing = line_spacing
            r2 = p2.add_run()
            r2.text = f"   {sub}"
            r2.font.size = Pt(sub_size)
            r2.font.color.rgb = C_SUBTLE
            r2.font.italic = True
            r2.font.name = "Calibri"


def add_code(
    slide,
    code: str,
    *,
    left: float = 0.5,
    top: float = 1.6,
    width: float = 12.33,
    height: float = 5.2,
    font_size: int = 11,
    dark: bool = True,
    label: str | None = None,
) -> None:
    """Code block with optional caption above (e.g., a filename)."""
    if label:
        lb = slide.shapes.add_textbox(Inches(left), Inches(top - 0.35), Inches(width), Inches(0.3))
        lp = lb.text_frame.paragraphs[0]
        lr = lp.add_run()
        lr.text = label
        lr.font.size = Pt(11)
        lr.font.bold = True
        lr.font.italic = True
        lr.font.color.rgb = C_HIGHLIGHT
        lr.font.name = "Consolas"

    bg_color = C_CODE_BG if dark else C_BG_LIGHT
    text_color = C_CODE_TEXT if dark else C_BODY

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color

    tb = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.1),
        Inches(width - 0.3), Inches(height - 0.2)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line if line else " "
        r.font.size = Pt(font_size)
        r.font.color.rgb = text_color
        r.font.name = "Consolas"


def add_table(
    slide,
    rows: list[tuple],
    *,
    left: float,
    top: float,
    col_widths: tuple,
    row_h: float,
    header: tuple | None = None,
    header_color=None,
    body_size: int = 10,
    header_size: int = 11,
) -> None:
    if header_color is None:
        header_color = C_TITLE
    if header:
        for ci, h in enumerate(header):
            x = Inches(left + sum(col_widths[:ci]))
            hb = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, Inches(top), Inches(col_widths[ci]), Inches(row_h)
            )
            hb.line.fill.background()
            hb.fill.solid()
            hb.fill.fore_color.rgb = header_color
            tf = hb.text_frame
            tf.margin_top = Inches(0.08)
            tf.margin_left = Inches(0.1)
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = "  " + h
            r.font.size = Pt(header_size)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            r.font.name = "Calibri"

    for ri, row in enumerate(rows):
        y = Inches(top + (row_h if header else 0) + row_h * ri)
        bg = RGBColor(0xFF, 0xFF, 0xFF) if ri % 2 == 0 else RGBColor(0xF8, 0xF8, 0xF8)
        for ci, val in enumerate(row):
            x = Inches(left + sum(col_widths[:ci]))
            cb = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, Inches(col_widths[ci]), Inches(row_h)
            )
            cb.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            cb.line.width = Pt(0.5)
            cb.fill.solid()
            cb.fill.fore_color.rgb = bg
            tf = cb.text_frame
            tf.margin_top = Inches(0.06)
            tf.margin_left = Inches(0.1)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = "  " + val
            r.font.size = Pt(body_size)
            r.font.color.rgb = C_BODY
            r.font.name = "Calibri" if ci > 0 else "Calibri"
            if ci == 0:
                r.font.bold = True


def add_notes(slide, notes: str) -> None:
    slide.notes_slide.notes_text_frame.text = notes


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_footer(slide, text: str = "Archimedes · Architecture Walkthrough") -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.33), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(8)
    r.font.color.rgb = C_SUBTLE
    r.font.italic = True
    r.font.name = "Calibri"


def add_pagenum(slide, n: int, total: int) -> None:
    box = slide.shapes.add_textbox(Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{n} / {total}"
    r.font.size = Pt(8)
    r.font.color.rgb = C_SUBTLE
    r.font.name = "Calibri"


def add_section_band(slide, section_no: int, section_name: str) -> None:
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.25))
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = C_ACCENT
    tf = band.text_frame
    tf.margin_top = Inches(0.02)
    tf.margin_left = Inches(0.5)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"§{section_no}  ·  {section_name}"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    r.font.name = "Calibri"


# ---------- build ----------


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    TOTAL = 32

    # =================================================================
    # Slide 1 — Title
    # =================================================================
    s = blank_slide(prs)
    box = s.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.33), Inches(1.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Archimedes"
    r.font.size = Pt(72)
    r.font.bold = True
    r.font.color.rgb = C_TITLE
    r.font.name = "Calibri"

    box2 = s.shapes.add_textbox(Inches(0.5), Inches(3.9), Inches(12.33), Inches(0.6))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Architecture Walkthrough — How to Rebuild It"
    r2.font.size = Pt(24)
    r2.font.italic = True
    r2.font.color.rgb = C_SUBTLE
    r2.font.name = "Calibri"

    box3 = s.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.33), Inches(0.4))
    tf3 = box3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = "[Presenter name]   ·   [Date]   ·   For: senior engineers / architects familiar with Claude + MCP"
    r3.font.size = Pt(13)
    r3.font.color.rgb = C_SUBTLE
    r3.font.name = "Calibri"

    add_notes(
        s,
        "Engineering-handoff framing — not marketing, not discovery. This deck is structured so that after 45 minutes the audience could start building their own Archimedes-shaped system. "
        "Open with: 'I'll spend roughly 45 minutes walking through the architecture. The intent is that you leave knowing what to build, in what order, with what invariants must hold. Stop me with questions; I'll dwell on whatever isn't landing.' "
        "Reference the other two decks if context helps — overview was problem-driven, Athena was forward-looking proposal. This is the engineering-deep-dive."
    )

    # =================================================================
    # Slide 2 — Mental model in one picture
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 0, "Mental model")
    add_title(s, "The whole system in one picture", subtitle="Four layers. Doctrine cuts across all four.")

    # Layer boxes — top to bottom
    layers = [
        ("HUMAN OPERATOR", "Approves HIGH scoring · names new actors · ratifies sources · reviews retractions", C_REFUSAL, RGBColor(0xFD, 0xEF, 0xED)),
        ("ORCHESTRATOR (Claude Code)", "Schedules · delegates · holds the run state across subagent invocations", C_TITLE, RGBColor(0xE8, 0xEF, 0xF5)),
        ("9 SUBAGENTS — each with its own context, write scope, doctrine", "collector · grader · analyst · red-team-analyst · actor-profiler · vuln-tracker · briefer · librarian + delegate-only orchestrator", C_ACCENT, RGBColor(0xE3, 0xEF, 0xFA)),
        ("MCP WRAPPERS + HOOKS — the tool interface", "Splunk · VirusTotal · Shodan · Censys · urlscan · theHarvester · SpiderFoot · RSS / discord-post.sh · splunk-log.sh", C_ACCENT, RGBColor(0xE3, 0xEF, 0xFA)),
        ("EXTERNAL — sources of truth and side effects", "OSINT feeds · Splunk indexes · git · Discord · Task Scheduler", C_KNOWN, RGBColor(0xE6, 0xF4, 0xEA)),
    ]
    box_h = Inches(0.85)
    box_top_start = Inches(1.6)
    box_w = Inches(11.0)
    box_left = Inches(1.5)
    gap = Inches(0.1)

    for i, (title, body, color, fill) in enumerate(layers):
        y = box_top_start + (box_h + gap) * i
        bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_left, y, box_w, box_h)
        bx.line.color.rgb = color
        bx.line.width = Pt(1.5)
        bx.fill.solid()
        bx.fill.fore_color.rgb = fill
        tf = bx.text_frame
        tf.margin_top = Inches(0.1)
        tf.margin_left = Inches(0.2)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = color
        r.font.name = "Calibri"

        p2 = tf.add_paragraph()
        p2.space_before = Pt(3)
        r2 = p2.add_run()
        r2.text = body
        r2.font.size = Pt(11)
        r2.font.color.rgb = C_BODY
        r2.font.name = "Calibri"

    # Doctrine ribbon down the left
    doc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.6), Inches(1.0), Inches(5.2))
    doc.line.fill.background()
    doc.fill.solid()
    doc.fill.fore_color.rgb = C_HIGHLIGHT
    tf = doc.text_frame
    tf.margin_top = Inches(2.2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "DOCTRINE\n(6 .md files\nversioned in git)"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    r.font.name = "Calibri"

    add_notes(
        s,
        "This is the map. Spend a full minute on it. Five layers, top to bottom: human operator (approves the high-stakes calls), orchestrator (Claude Code instance that delegates), subagents (9 specialized roles with isolated contexts), tools (MCPs + hooks), external (where the data lives and where side effects land). "
        "Doctrine is the orange ribbon down the left — it cuts across all five layers because every layer reads from it. "
        "Reference this slide throughout the rest of the deck when introducing each layer's details."
    )
    add_footer(s)
    add_pagenum(s, 2, TOTAL)

    # =================================================================
    # Slide 3 — Five invariants
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 0, "Mental model")
    add_title(s, "Five load-bearing invariants",
              subtitle="If your rebuild preserves these five, the rest is implementation detail.")

    invariants = [
        ("1.  Doctrine as code",
         "Every behavior the agent exhibits traces back to a versioned .md file. No prompts hidden in subagent code; "
         "no logic that contradicts the doctrine files. If you can't grep for the rule in doctrine/, the rule doesn't exist."),
        ("2.  Subagent context isolation",
         "Each subagent gets the minimum context to do its job. Orchestrator never reads raw articles. Briefer never sees "
         "the coverage log. Red-team never wrote the primary finding. Lower error rate by construction."),
        ("3.  Librarian-only side effects",
         "Only the librarian writes to git, Splunk, or Discord. Other subagents can write to disk inside their scope but "
         "the externalization happens in one place — easy to audit, easy to LEGAL-POLICY-scan, easy to rate-limit."),
        ("4.  Human approval gates on high-stakes calls",
         "HIGH threat-box scoring → /approve-scoring. New actor → /new-actor. Source-grade ratification → human review. "
         "The agent proposes; the human approves on anything with reputational consequences."),
        ("5.  Audit trail is additive, never silent",
         "Retractions append to the original, never overwrite. Source grade downgrades log to source-grade-log.md. Every "
         "policy refusal logs to policy-violations.yaml. Git history + Splunk events + Discord posts together reconstruct any action."),
    ]
    add_bullets(s, invariants, body_size=15, sub_size=11)
    add_notes(
        s,
        "These are the five things that, if you preserve them, the rebuild will still be Archimedes-shaped — even if you use Python where I used Bash, or OpenCTI where I used markdown, or Teams where I used Discord. "
        "If you break any of these in the rebuild, you've built something different. Worth showing the audience this slide twice — once now (so they have the frame) and once at the rebuild-guide section (so they remember which lines not to cross)."
    )
    add_footer(s)
    add_pagenum(s, 3, TOTAL)

    # =================================================================
    # Slide 4 — On-disk layout
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 0, "Mental model")
    add_title(s, "On-disk layout", subtitle="What lives where. Real directory tree, abbreviated.")

    tree = """archimedes/
├── doctrine/                       6 .md files — agent's working memory
│   ├── LEGAL-POLICY.md             prohibited queries, ITAR boundary, copyright
│   ├── INTEL-GRADING.md            Admiralty 6×6, WEP, single-source-veto
│   ├── INTEL-BRIEF-STANDARDS.md    Layer 1/Layer 2 spec, preflight checklist
│   ├── THREAT-BOX-METHODOLOGY.md   per-category scoring, evidence-min tables
│   ├── RETRACTION-POLICY.md        additive retractions, 72h auto-downgrade
│   └── FLASH-POLICY.md             7 triggers, quiet hours, critical override
│
├── .claude/
│   ├── agents/                     9 subagent definition .md files
│   │   ├── collector.md  grader.md  analyst.md  red-team-analyst.md
│   │   ├── actor-profiler.md  vuln-tracker.md  briefer.md  librarian.md
│   ├── skills/                     reusable skill bundles
│   │   ├── smart-brevity/   admiralty-grading/   sat-ach/   sat-kac/
│   ├── commands/                   slash-command definitions (/investigate, etc.)
│   └── hooks/                      bash scripts for deterministic side effects
│       ├── discord-post.sh         POST to Discord via webhook
│       └── splunk-log.sh           HEC ingest to Splunk
│
├── mcps/                           8 MCP wrappers — uv workspace members
│   ├── splunk-query/  virustotal/  shodan-mcp/  rss-bridge/
│   ├── urlscan/  censys/  theharvester/  spiderfoot/
│
├── threats/                        the corpus — agent writes; humans review
│   ├── raw-signal/                 collector output (un-promoted)
│   ├── findings/                   grader output (promoted, graded)
│   ├── briefs/                     briefer output (Layer 1 + Layer 2)
│   ├── threat-actors/<actor>/      per-actor dossier (profile + IOCs + threat-box)
│   ├── vulnerabilities/<cve>/      per-CVE tracking
│   └── iocs/_master-index.yaml     regenerated indexes
│
├── infrastructure/                 agent-readable YAML config
│   ├── source-grades.yaml          which sources, which grades, which categories
│   ├── source-health.yaml          runtime health state (gitignored — collector writes)
│   ├── watch-config.yaml           standing brief sections + watchlists
│   ├── authorized-targets.yaml     Hard Rule 4 — only these may be actively scanned
│   └── scheduler/                  Windows Task Scheduler XML templates
│
├── scripts/                        Python helpers (not subagents)
│   ├── run_phase.ps1               wrapper for scheduled brief invocation
│   ├── splunk_log.py               HEC log helper
│   ├── regenerate_ioc_index.py     index regen
│   └── discord_listener.py         inbound /command bridge
│
├── docs/handoffs/                  session retrospectives (session-N.md)
├── pyproject.toml                  root project + uv workspace declaration
├── uv.lock                         locked dependency graph
└── CLAUDE.md                       agent charter + operational notes log"""
    add_code(s, tree, top=1.55, height=5.55, font_size=10, dark=False)
    add_notes(
        s,
        "Walk through top-to-bottom. Three things to emphasize: "
        "(1) doctrine/ at the top — these six files are the agent's working memory; the rebuild starts here. "
        "(2) .claude/ contains everything the agent runtime needs — agents, skills, commands, hooks; the equivalent in your stack might be different (Claude Code uses this layout, other Claude runtimes have their own). "
        "(3) threats/ is the corpus — the agent's output lives here, versioned in git. This is what makes the system auditable: every brief, every finding, every dossier, every retraction is a markdown file in git history."
    )
    add_footer(s)
    add_pagenum(s, 4, TOTAL)

    # =================================================================
    # Slide 5 — Doctrine catalog
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 1, "Doctrine layer")
    add_title(s, "The six doctrine files",
              subtitle="The agent's working memory. Versioned in git. Reviewable by analysts and engineers.")

    rows = [
        ("LEGAL-POLICY.md",
         "prohibited queries · authorized targets · ITAR · copyright",
         "every subagent (before any tool call)",
         "session start + on update"),
        ("INTEL-GRADING.md",
         "Admiralty 6×6 · WEP · single-source-veto · corroboration rules",
         "grader",
         "on every finding promotion"),
        ("INTEL-BRIEF-STANDARDS.md",
         "Layer 1 / Layer 2 spec · preflight checklist · word/char budgets",
         "briefer",
         "on every brief compose"),
        ("THREAT-BOX-METHODOLOGY.md",
         "per-category composite scoring · evidence-minimum tables",
         "actor-profiler",
         "on every /update-tracking"),
        ("RETRACTION-POLICY.md",
         "when to retract vs correct · additive · 72h auto-downgrade",
         "grader + librarian",
         "on retraction trigger"),
        ("FLASH-POLICY.md",
         "7 triggers · quiet hours · critical override · 72h clock",
         "collector + grader",
         "on every FLASH sweep"),
    ]
    add_table(
        s, rows,
        left=0.5, top=1.5,
        col_widths=(2.9, 5.0, 2.6, 1.83),
        row_h=0.55,
        header=("Doctrine file", "Governs", "Read by", "When"),
        body_size=10,
    )

    add_notes(
        s,
        "Six files. Each one governs a slice of the agent's behavior. The 'Read by' column maps doctrine to subagent — the librarian doesn't read THREAT-BOX-METHODOLOGY because it doesn't score actors; the grader doesn't read INTEL-BRIEF-STANDARDS because it doesn't compose briefs. "
        "Minimal-required-context, applied even to doctrine. This pattern is load-bearing — without it, every subagent loads every doctrine file and the context bloat costs accuracy and latency."
    )
    add_footer(s)
    add_pagenum(s, 5, TOTAL)

    # =================================================================
    # Slide 6 — Anatomy of a doctrine file
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 1, "Doctrine layer")
    add_title(s, "Anatomy of a doctrine file",
              subtitle="Real excerpt from doctrine/INTEL-GRADING.md")

    code = """# INTEL-GRADING.md — Intelligence Source Grading System

> **Archimedes doctrine — grading.**
> This file is authoritative. Every finding's digraph must be
> justifiable against this document.

---

## Source Reliability (A–F)

| Grade | Label                | Description                                  |
|-------|----------------------|----------------------------------------------|
| A     | Completely Reliable  | Consistent, verified track record.           |
| B     | Usually Reliable     | Strong track record, minor inaccuracies.    |
| C     | Fairly Reliable      | Right before, but enough errors for caution. |
| D     | Not Usually Reliable | More misses than hits. Lead, not finding.   |
| E     | Unreliable           | History of false/misleading information.    |
| F     | Cannot Be Judged     | New, unknown, no track record yet.          |

## Credibility Assessment Checklist (authoritative)

### 1 — Confirmed — ALL of:
  - Multiple independent A/B-grade sources
  - No public contradictions in 24h corroboration window
  - First-party Splunk telemetry if applicable

### 2 — Probably True — ALL of:
  - One A-grade source
  - Plausibility test passes (consistent with known TTPs)
  - No active contradictions

### 3 — Possibly True — single source, cannot confirm or deny.

[...continues for 4, 5, 6...]"""
    add_code(s, code, top=1.55, height=5.5, font_size=10, dark=False)
    add_notes(
        s,
        "Three structural patterns to call out: "
        "(1) Authority statement at the top — 'this file is authoritative.' The grader subagent is instructed to defer to this file over its own judgment. "
        "(2) Lookup tables — Admiralty grades, WEP probability bands, evidence-minimum tables for threat-box. The agent uses these like a dev uses a switch statement. "
        "(3) Checklist-style rules — 'ALL of these conditions must hold' or 'ANY of these triggers' — written as if for human auditors, applied by the agent as decision logic. "
        "Doctrine files are plain English, but they read like code. That's intentional."
    )
    add_footer(s)
    add_pagenum(s, 6, TOTAL)

    # =================================================================
    # Slide 7 — Hard Rules
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 1, "Doctrine layer")
    add_title(s, "Hard Rules — the immutable refusal boundary",
              subtitle="From CLAUDE.md. Agent refuses even when prompted to bypass. Logged to policy-violations.yaml.")
    add_bullets(
        s,
        [
            ("1 · Legal policy is non-negotiable", "If LEGAL-POLICY prohibits it, refuse even if user insists. Log the attempt."),
            ("2 · Never originate attribution", "Agent only reports what other sources have attributed, citing them."),
            ("3 · No exploitation, ever", "No PoC code, payloads, exploit guides — not for testing, research, or 'educational' purposes."),
            ("4 · Never scan third parties", "Active recon only against targets in authorized-targets.yaml. SpiderFoot + theHarvester MCPs enforce module-level allowlists."),
            ("5 · Human sign-off for HIGH threat levels", "When actor-profiler proposes HIGH, posts to #actor-review and waits for /approve-scoring. Does not auto-commit."),
            ("6 · 15-word quote limit, one quote per source", "Copyright compliance. Hard-enforced in preflight checklist."),
            ("7 · Credentials are radioactive", "If a query surfaces credentials, never store them. Count, report exposure, discard."),
            ("8 · Splunk first-party > external", "When first-party telemetry contradicts external sources, first-party wins. The external source gets graded down."),
        ],
        body_size=13, sub_size=11,
    )
    add_notes(
        s,
        "Hard Rules are the audit-defense layer. They're enforced at three levels: in doctrine prose (agent reads and respects), in subagent definitions (every subagent has 'before any action — consult LEGAL-POLICY' as step 0), and in code at the MCP wrappers (allowlist + refusal before HTTP call). "
        "Three levels of defense in depth. Removing any one weakens the system. "
        "If a rebuild changes the Hard Rules list — e.g. adds rule 9, removes rule 7 — that's fine, but make sure the new list is enforced at all three levels."
    )
    add_footer(s)
    add_pagenum(s, 7, TOTAL)

    # =================================================================
    # Slide 8 — Skills
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 1, "Doctrine layer")
    add_title(s, "Skills — reusable analytic procedures",
              subtitle=".claude/skills/<name>/  — invoked by subagents, share procedures across them.")
    add_bullets(
        s,
        [
            ("smart-brevity/",
             "Brief composition rules. Banned phrases, active voice, lead-with-impact, preflight checklist. "
             "Invoked by briefer on every brief. Has references/ subdirectory: banned-phrases.md, brief-templates.md, preflight-checklist.md."),
            ("admiralty-grading/",
             "Source reliability + credibility evaluation procedure. Invoked by grader on every finding promotion. "
             "Generates the digraph the grader writes to the finding's frontmatter."),
            ("sat-ach/",
             "Analysis of Competing Hypotheses — structured analytic technique. Builds an N×M matrix of hypotheses × evidence, "
             "ranks by inconsistency count. Invoked by analyst (primary) and red-team-analyst (contrarian)."),
            ("sat-kac/",
             "Key Assumptions Check — surface load-bearing premises that the finding rests on. "
             "Invoked by analyst when a finding's WEP is 'likely' or higher."),
            ("threat-box-scoring/",
             "Per-category composite scoring procedure with evidence-minimum table enforcement. Invoked by actor-profiler. "
             "Outputs the proposed threat-box.yaml; routes HIGH proposals to the /approve-scoring gate."),
            ("ioc-extraction/",
             "Pull indicators from finding body text — IPs, domains, hashes, CVEs, actor mentions, MITRE IDs. Invoked by collector."),
        ],
        body_size=13, sub_size=11,
    )
    add_notes(
        s,
        "Skills are like libraries — reusable procedures the subagents call. They live under .claude/skills/<name>/. Each has a SKILL.md (the main procedure) and usually a references/ subdirectory with lookup tables or templates. "
        "Pattern: when two subagents need the same procedure, factor it into a skill. The grader and the analyst both need ACH; sat-ach is the skill they both call. "
        "For a rebuild: skills aren't strictly necessary — you could inline the procedures into each subagent — but they pay off as soon as you have any procedure shared across multiple roles."
    )
    add_footer(s)
    add_pagenum(s, 8, TOTAL)

    # =================================================================
    # Slide 9 — Subagent inventory
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 2, "Subagents")
    add_title(s, "Subagent inventory", subtitle="9 subagents, each with isolated context, write scope, doctrine read, tool set.")

    rows = [
        ("collector", "raw OSINT collection + watchlist filter", "threats/raw-signal/", "FLASH-POLICY, LEGAL-POLICY", "WebFetch, WebSearch, MCP tools"),
        ("grader", "promote raw → finding · apply Admiralty + WEP", "threats/findings/", "INTEL-GRADING, RETRACTION", "Read, Write, Edit, Splunk"),
        ("analyst", "SAT / ACH / KAC structured analysis", "findings (analysis sections)", "INTEL-GRADING, SAT skills", "Read, Write, Edit"),
        ("red-team-analyst", "challenge HIGH-confidence findings", "findings (red_team section)", "sat-ach (contrarian)", "Read, Edit"),
        ("actor-profiler", "dossier maintenance + threat-box scoring", "threats/threat-actors/*/", "THREAT-BOX, ACTOR-STANDARD", "Read, Write, Edit, Splunk"),
        ("vuln-tracker", "CVE tracking · KEV monitoring", "threats/vulnerabilities/*/", "NVD + KEV refs", "Read, Write, WebFetch"),
        ("briefer", "compose all brief types · Layer 1 + Layer 2", "threats/briefs/", "INTEL-BRIEF-STANDARDS", "Read, Write, Edit"),
        ("librarian", "commit · post Discord · log Splunk · regen indices", "git + Splunk + Discord + indices", "RETRACTION-POLICY", "Read, Bash, Edit"),
    ]
    add_table(
        s, rows,
        left=0.5, top=1.5,
        col_widths=(1.8, 3.8, 2.6, 2.3, 1.83),
        row_h=0.55,
        header=("Subagent", "Role", "Write scope", "Doctrine read", "Tools"),
        body_size=9,
        header_size=10,
    )

    add_notes(
        s,
        "Eight defined subagents + the orchestrator (which is a 9th, with no write scope of its own — it only delegates). Walk through them top to bottom; the table is dense but every column matters: "
        "Role is the elevator pitch. Write scope is what files/systems they're allowed to modify. Doctrine read is which .md files load into their context. Tools is the tool subset they're given (e.g. red-team-analyst only gets Read + Edit — can't write new files, only modify findings). "
        "The minimal tool surface per subagent is part of context isolation — fewer tools = simpler decision space = lower error rate."
    )
    add_footer(s)
    add_pagenum(s, 9, TOTAL)

    # =================================================================
    # Slide 10 — Subagent definition anatomy
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 2, "Subagents")
    add_title(s, "Anatomy of a subagent definition",
              subtitle="Real excerpt — .claude/agents/briefer.md (frontmatter + structure)")

    code = """---
name: briefer
description: Use to compose every Archimedes brief from graded findings.
  Invoke for scheduled morning brief (08:00 EDT), scheduled afternoon brief
  (16:00 EDT), async FLASH alerts ... Reads approved findings, the coverage
  log for anti-repetition, watch-config for standing sections, and doctrine.
  Invokes the smart-brevity skill always — drafts prose, runs the 13-item
  pre-flight checklist, regenerates failing sections until all pass.
  Writes only to threats/briefs/. Does not post to Discord (librarian's job),
  does not grade (grader's job), does not run SATs (analyst's job).
tools: Read, Write, Edit, Glob, Grep
model: opus
---

# Briefer Subagent

## Role
You are the briefer. You take the graded, analyzed, red-teamed findings
the other subagents produced and compose them into the documents Ryan
actually reads: six brief types plus retractions.

## Before any action — consult LEGAL-POLICY
- Hard Rule 6 (quote discipline) is your operational constraint
- ITAR/export-control check on every brief

## Procedure — scheduled brief (Morning or Afternoon)
1. Load inputs (findings, coverage log, watch-config)
2. Filter by digraph ≥ B2, by freshness window
3. Categorize into sections per INTEL-BRIEF-STANDARDS
4. Apply correlation callouts
5. Compose Layer 1 + Layer 2
6. Run 13-item preflight checklist (smart-brevity skill)
7. Write file, update _coverage-log.yaml, back-write findings"""
    add_code(s, code, top=1.55, height=5.5, font_size=10, dark=False)
    add_notes(
        s,
        "Three load-bearing parts of every subagent definition: "
        "(1) Frontmatter — name, description (this is the elevator pitch the orchestrator reads when deciding who to delegate to), tools (the minimum set), model (haiku for cheap, sonnet for medium, opus for complex). "
        "(2) Role + scope — written in second person ('you are the briefer') — Claude reads this as identity. Clear delineation of what this subagent does NOT do is as important as what it does. "
        "(3) Procedure — numbered steps. This is the agent's runbook. When the orchestrator delegates, the subagent reads its own definition and follows the procedure."
    )
    add_footer(s)
    add_pagenum(s, 10, TOTAL)

    # =================================================================
    # Slide 11 — Context isolation
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 2, "Subagents")
    add_title(s, "Context isolation — what each subagent sees vs. doesn't",
              subtitle="The counterintuitive insight: less context per agent = lower error rate.")

    # Two columns — sees vs. doesn't see
    add_bullets(
        s,
        [
            ("orchestrator", "schedule, delegation directives, pipeline state. Never reads raw articles. Never sees graded findings in detail."),
            ("collector", "source list, watchlists, time window. Never reads coverage log. Never sees prior briefs."),
            ("grader", "raw signal from current sweep, INTEL-GRADING, historical-source health. Never sees brief text."),
            ("analyst", "the finding it's analyzing + analytic framework references. Never sees other findings."),
            ("red-team-analyst", "ONE finding + the analyst's own analysis. Cannot have written the primary — contrarian by construction."),
            ("briefer", "approved findings + coverage log + watch-config + doctrine. Never reads raw signal."),
            ("librarian", "the file or files it's about to ship + the change summary. Never composes content."),
        ],
        body_size=13, sub_size=11,
    )

    # Bottom callout
    cb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.0), Inches(12.33), Inches(0.85))
    cb.line.color.rgb = C_KNOWN
    cb.line.width = Pt(1.5)
    cb.fill.solid()
    cb.fill.fore_color.rgb = RGBColor(0xE6, 0xF4, 0xEA)
    ctf = cb.text_frame
    ctf.margin_top = Inches(0.12)
    ctf.margin_left = Inches(0.2)
    ctf.margin_right = Inches(0.2)
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = (
        "Why this matters: agents hallucinate against context they don't need. "
        "Giving the briefer access to raw articles makes it tempted to re-grade. "
        "Giving the red-team the primary author's reasoning makes it harder to disagree. "
        "Minimum context = minimum surface for the wrong inference."
    )
    cr.font.size = Pt(12)
    cr.font.italic = True
    cr.font.color.rgb = C_BODY
    cr.font.name = "Calibri"

    add_notes(
        s,
        "This is the most counterintuitive design choice in the system. The instinct is 'give the model everything it might need.' The opposite works better — minimal scope per agent, with the orchestrator handling cross-agent state. "
        "Show one concrete example: red-team-analyst. It's contrarian by construction because it didn't write the primary finding. If you gave it the same context the analyst had, it would naturally converge on the same conclusion (anchoring bias). By keeping the contexts separate, you create the conditions for genuine disagreement."
    )
    add_footer(s)
    add_pagenum(s, 11, TOTAL)

    # =================================================================
    # Slide 12 — The grading chain
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 2, "Subagents")
    add_title(s, "The grading chain — case study",
              subtitle="One CVE-attribution claim traced through collector → grader → analyst → red-team")

    # Vertical flow
    steps = [
        ("collector",
         "Picks up Rapid7 blog post via RSS. Filters against watchlist (MuddyWater = tracked actor → KEEP). "
         "Writes raw-signal/2026-05-06-1218-rapid7-muddywater.md with minimal frontmatter. No grading.",
         C_ACCENT),
        ("grader",
         "Reads the raw signal. Applies INTEL-GRADING: Rapid7 provisional A · single source → credibility 2 max. "
         "Single-source-veto caps forward claim at 'likely.' Promotes to findings/finding-2026-05-06-FLASH-0002.md with "
         "digraph A2, wep likely, corroboration.independent=false, 72h auto-downgrade clock set.",
         C_TITLE),
        ("analyst",
         "Reads ONLY this finding + sat-ach / sat-kac skills. Builds ACH matrix of attribution hypotheses (TeamPCP / "
         "MuddyWater / unattributed). Records analysis_sections.ach + .kac on the finding. Flags red_team_review_required "
         "because WEP=likely on single source — not 'very likely' but close.",
         C_TITLE),
        ("red-team-analyst",
         "Receives just this finding + analyst's analysis. Argues FOR rejected hypotheses. Tests for confirmation bias. "
         "Either confirms the assessment (signs off) or flags weaknesses. Writes findings.red_team_review section.",
         C_REFUSAL),
        ("→ next: briefer reads this finding (along with others promoted same window) and composes Layer 1 + Layer 2",
         "Briefer never re-grades. Librarian never re-composes. Each subagent stays in its lane.",
         C_KNOWN),
    ]

    step_top = Inches(1.55)
    step_h = Inches(1.0)
    gap = Inches(0.1)
    box_w = Inches(12.33)
    for i, (title, body, color) in enumerate(steps):
        y = step_top + (step_h + gap) * i
        bx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, box_w, step_h)
        bx.line.color.rgb = color
        bx.line.width = Pt(1.5)
        bx.fill.solid()
        bx.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf = bx.text_frame
        tf.margin_top = Inches(0.1)
        tf.margin_left = Inches(0.2)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = color
        r.font.name = "Calibri"

        p2 = tf.add_paragraph()
        p2.space_before = Pt(3)
        r2 = p2.add_run()
        r2.text = body
        r2.font.size = Pt(10)
        r2.font.color.rgb = C_BODY
        r2.font.name = "Calibri"

    add_notes(
        s,
        "Trace one real finding through the chain. The MuddyWater / Rapid7 example actually happened — this is the FLASH-0002 finding referenced in CLAUDE.md's operational notes. "
        "Five stages, each with a different subagent, each with isolated context. The grader doesn't speculate beyond what INTEL-GRADING permits. The analyst doesn't grade. The red-team doesn't analyze — it challenges. Strict separation of concerns. "
        "Outcome: that finding's attribution leg eventually auto-downgraded to C3 (possibly true) when the 72h corroboration window closed unbroken. The discipline pays off."
    )
    add_footer(s)
    add_pagenum(s, 12, TOTAL)

    # =================================================================
    # Slide 13 — The librarian
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 2, "Subagents")
    add_title(s, "The librarian — the only writer to side effects",
              subtitle="LEGAL-POLICY content scan + git + Splunk + Discord, in that order, every time.")

    add_bullets(
        s,
        [
            ("Why one writer",
             "Every external publication funnels through one agent — easy to audit, easy to rate-limit, easy to add a "
             "LEGAL-POLICY content-scan gate before anything leaves the building."),
            ("Mode-based procedure catalog (8 modes)",
             "Mode 1: post scheduled brief · Mode 2: post FLASH (with quiet-hours/critical-override logic) · "
             "Mode 3: process FLASH queue at 09:00 · Mode 4: actor-profiler commit · Mode 5: vuln-tracker commit · "
             "Mode 6: retraction · Mode 7: source-grade update · Mode 8: index regen"),
            ("Tools = Bash, Read, Edit, Glob, Grep",
             "Bash for hook invocation (.claude/hooks/discord-post.sh, .claude/hooks/splunk-log.sh, git commands). "
             "Read + Edit for the file(s) being shipped. NO Write — librarian doesn't author content."),
            ("Hook layer for deterministic side effects",
             "Discord posts go through discord-post.sh; Splunk events through splunk-log.sh. Bash scripts, not subagents — "
             "no LLM cost on every post, deterministic behavior, easy to debug."),
            ("LEGAL-POLICY content scan is the gate",
             "Every Discord post preceded by content scan: no credentials, no TLP:RED, no ITAR-questionable detail. "
             "If anything trips, halt and flag. The librarian is the last gate before content becomes public."),
            ("Wrapper catchup-push",
             "scripts/run_phase.ps1 runs `git push origin main` at end of every phase. If librarian's push hung or "
             "skipped, wrapper catches up. Belt-and-suspenders against subagent variance."),
        ],
        body_size=13, sub_size=11,
    )
    add_notes(
        s,
        "The librarian is the system's most important subagent for safety. Concentrating all side effects in one place is a load-bearing pattern: it means there's exactly one piece of code (the librarian's LEGAL-POLICY scan step) that decides whether content goes public. "
        "Other subagents can write to disk (briefer to threats/briefs/, grader to threats/findings/), but nothing leaves the local filesystem without librarian's blessing. "
        "The mode catalog is just how we organize the librarian's procedures — could be one giant procedure, but the modes make the doctrine easier to reason about."
    )
    add_footer(s)
    add_pagenum(s, 13, TOTAL)

    # =================================================================
    # Slide 14 — Orchestrator
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 2, "Subagents")
    add_title(s, "The orchestrator — what it does and never does",
              subtitle="Claude Code session itself. Schedules. Delegates. Never analyzes or writes.")

    # Two columns: does / never does
    left_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.6), Inches(6.0), Inches(0.4))
    left_box.line.fill.background()
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = C_KNOWN
    lt = left_box.text_frame
    lt.margin_top = Inches(0.06)
    lp = lt.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    lr = lp.add_run()
    lr.text = "DOES"
    lr.font.size = Pt(14)
    lr.font.bold = True
    lr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    lr.font.name = "Calibri"

    right_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.83), Inches(1.6), Inches(6.0), Inches(0.4))
    right_box.line.fill.background()
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = C_REFUSAL
    rt = right_box.text_frame
    rt.margin_top = Inches(0.06)
    rp = rt.paragraphs[0]
    rp.alignment = PP_ALIGN.CENTER
    rr = rp.add_run()
    rr.text = "NEVER DOES"
    rr.font.size = Pt(14)
    rr.font.bold = True
    rr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    rr.font.name = "Calibri"

    add_bullets(
        s,
        [
            "Receive trigger (scheduled / FLASH / on-demand)",
            "Read pipeline doctrine (which subagents in what order)",
            "Delegate to subagents via Task tool",
            "Pass minimal context between subagent invocations",
            "Track run state across handoffs",
            "Surface high-stakes decisions to operator",
            "Handle pipeline failures (retry / degrade / abort)",
        ],
        left=0.5, top=2.1, width=6.0, height=4.5, body_size=13, sub_size=11,
    )
    add_bullets(
        s,
        [
            "Read raw articles (collector's job)",
            "Grade findings (grader's job)",
            "Compose brief content (briefer's job)",
            "Write to git, Splunk, or Discord (librarian's job)",
            "Make attribution claims (no one's job — Hard Rule 2)",
            "Approve HIGH scoring (operator's job — Hard Rule 5)",
            "Read .env or anything credential-bearing",
        ],
        left=6.83, top=2.1, width=6.0, height=4.5, body_size=13, sub_size=11,
    )

    add_notes(
        s,
        "The orchestrator is what you start with when you open Claude Code in this repo. It reads CLAUDE.md, knows the agent's identity, and waits for triggers. "
        "When a trigger fires (you typed /brief morning, or the scheduled task invoked claude -p), the orchestrator looks at the pipeline doctrine in CLAUDE.md and delegates. "
        "Two non-obvious things in the 'never does' column: it never composes brief content (even if it knows what the brief should say) and it never makes attribution claims. Both of those would violate the role separation — and Hard Rule 2 specifically."
    )
    add_footer(s)
    add_pagenum(s, 14, TOTAL)

    # =================================================================
    # Slide 15 — MCP architecture
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 3, "MCPs + tool integration")
    add_title(s, "MCP architecture",
              subtitle="Model Context Protocol — stdio + JSON-RPC between Claude Code and tool servers.")

    add_bullets(
        s,
        [
            ("MCP = Model Context Protocol",
             "Open spec from Anthropic. Defines stdio+JSON-RPC contract between an LLM runtime (Claude Code) and a tool "
             "server. Tools are discovered at runtime via `tools/list`; called via `tools/call`."),
            ("Each MCP is its own process",
             "Claude Code spawns the MCP server as a subprocess on session start. .mcp.json declares 'this server is "
             "named X, run command Y in directory Z.' Server lives as long as the Claude Code session."),
            ("FastMCP — our framework choice",
             "Python wrapper around the MCP spec. Decorator-style tool definition. Handles protocol plumbing; we write "
             "the tool functions."),
            ("Per-MCP isolated venv",
             "uv workspace member per MCP. Each has its own pyproject.toml + dependencies + tests. Why: dependency "
             "conflicts between MCPs (e.g. SpiderFoot wants CherryPy 18; theHarvester wants dnspython 2.3) — isolation prevents them."),
            ("Tools surface as mcp__<server>__<tool>",
             "From the agent's perspective: mcp__shodan__lookup_host, mcp__virustotal__lookup_file, "
             "mcp__theharvester__enumerate, mcp__spiderfoot__passive_scan. Namespace per server."),
        ],
        body_size=13, sub_size=11,
    )
    add_notes(
        s,
        "If your audience knows MCP, this is review. If they don't, this is foundational. Three things to land: "
        "(1) MCP is a protocol, not a library — your rebuild can use any MCP client/server library (Python, TypeScript, Rust). "
        "(2) Tool servers are subprocesses — keeps them isolated, lets you crash one without taking down the agent. "
        "(3) The naming convention `mcp__<server>__<tool>` is how the agent thinks about tool calls — every tool reference in this codebase uses that exact pattern."
    )
    add_footer(s)
    add_pagenum(s, 15, TOTAL)

    # =================================================================
    # Slide 16 — MCP wrapper anatomy
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 3, "MCPs + tool integration")
    add_title(s, "Anatomy of an MCP wrapper",
              subtitle="Real directory structure — mcps/spiderfoot/")

    tree = """mcps/spiderfoot/
├── pyproject.toml                  # spiderfoot-mcp; httpx + mcp[cli] + pydantic
├── README.md                       # public-facing doc — what tools, what auth
└── src/spiderfoot_mcp/
    ├── __init__.py                 # version + main() entry point
    ├── config.py                   # pydantic-settings; reads .env
    │                               # SPIDERFOOT_URL, _USERNAME, _PASSWORD,
    │                               # _SCAN_TIMEOUT_SECONDS, _POLL_INTERVAL_SECONDS
    ├── exceptions.py               # SpiderFootError + 6 subclasses
    │                               # SpiderFootConnectionError, AuthError,
    │                               # RequestError, PolicyError, ScanError,
    │                               # TimeoutError(scan_id)
    ├── models.py                   # pydantic models + PASSIVE_MODULE_ALLOWLIST
    │                               # the load-bearing safety primitive
    ├── policy.py                   # validate_modules() — refuse non-passive
    │                               # before any HTTP call
    ├── spiderfoot_client.py        # synchronous httpx client
    │                               # start_scan, get_scan_status, export_results,
    │                               # run_passive_scan (orchestrates poll loop)
    └── server.py                   # @mcp.tool() decorated entry points
                                    # passive_scan, list_modules, health

tests/
├── test_config.py                  # config-loader unit tests
├── test_policy.py                  # allowlist enforcement — 14 cases
└── test_client.py                  # respx-mocked HTTP behavior — 23 cases"""
    add_code(s, tree, top=1.55, height=5.4, font_size=10, dark=False)
    add_notes(
        s,
        "Same 8-file structure for every MCP. Pattern stabilized after the third MCP — once you see it twice you can predict where everything lives. "
        "Three files are mandatory: config.py (env loading), <name>_client.py (the HTTP wrapper), server.py (the @mcp.tool() decorated endpoints). "
        "The other three (exceptions.py, models.py, policy.py) appear when the MCP has policy enforcement to do (like SpiderFoot's passive allowlist) or a complex domain model (like theHarvester's source list). Simpler MCPs (like rss-bridge) can skip policy.py."
    )
    add_footer(s)
    add_pagenum(s, 16, TOTAL)

    # =================================================================
    # Slide 17 — MCP wrapper code: policy primitive
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 3, "MCPs + tool integration")
    add_title(s, "The load-bearing safety primitive — policy.py",
              subtitle="Real code from mcps/spiderfoot/src/spiderfoot_mcp/ — Hard Rule 4 enforced before any HTTP call.")

    code = """# models.py — the allowlist + prohibited list (data)

PASSIVE_MODULE_ALLOWLIST: tuple[str, ...] = (
    "sfp_dnsresolve", "sfp_whois", "sfp_crt", "sfp_certspotter",
    "sfp_archiveorg", "sfp_threatfox", "sfp_virustotal",
    "sfp_bingsearch", "sfp_duckduckgo", "sfp_hibp_pastes",
    # ... 45 entries total
)

PROHIBITED_MODULES: tuple[str, ...] = (
    "sfp_tool_nmap",     "sfp_tool_nuclei",   "sfp_spider",
    "sfp_dnsbrute",      "sfp_screenshot",    "sfp_portscan_tcp",
    "sfp_bingsharedip",
)

# policy.py — the enforcement (logic)

def validate_modules(modules: list[str] | None) -> list[str]:
    \"\"\"Refuse non-passive modules BEFORE any HTTP call. Hard Rule 4.\"\"\"
    if not modules:
        return list(DEFAULT_MODULES)

    cleaned, rejected, prohibited_hits = [], [], []
    for raw in modules:
        stripped = raw.strip().removeprefix("module_")
        lower = stripped.lower()
        if lower not in _ALLOWLIST_LOWERCASE:
            if lower in (m.lower() for m in PROHIBITED_MODULES):
                prohibited_hits.append(stripped)
            else:
                rejected.append(stripped)
            continue
        cleaned.append(canonical_case(lower))

    if prohibited_hits or rejected:
        raise SpiderFootPolicyError(
            f"Module(s) refused: {prohibited_hits + rejected}. "
            "See LEGAL-POLICY.md SpiderFoot section."
        )
    return cleaned

# server.py — wired in BEFORE the HTTP call
@mcp.tool()
def passive_scan(target: str, modules: list[str] | None = None) -> dict:
    cleaned = validate_modules(modules)   # ← raises if non-passive
    return run_scan(target, cleaned)"""
    add_code(s, code, top=1.55, height=5.5, font_size=10, dark=True)
    add_notes(
        s,
        "This is the single most important code pattern in the system. Hard Rule 4 ('never scan third parties') is enforced as a function call BEFORE the HTTP request. The agent literally cannot make SpiderFoot scan an unauthorized target with an active module — the wrapper refuses with a PolicyError at the input boundary. "
        "The same pattern shows up in theharvester (PASSIVE_SOURCE_ALLOWLIST), in WebFetch (LEGAL-POLICY check on URLs), in WebSearch (prohibited query patterns). "
        "Rebuild lesson: encode your policy at the integration layer, not in prose-only doctrine. Doctrine tells the agent what to do; allowlists make it impossible to do otherwise."
    )
    add_footer(s)
    add_pagenum(s, 17, TOTAL)

    # =================================================================
    # Slide 18 — MCP inventory
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 3, "MCPs + tool integration")
    add_title(s, "MCP inventory", subtitle="8 wrappers — what each surfaces, auth model, policy enforcement.")

    rows = [
        ("splunk-query",     "first-party telemetry · SPL search · health", "Basic auth (REST 8089)",      "no allowlist · query rate limit"),
        ("virustotal",       "domain / IP / file / URL reputation",         "API key header",             "rate-limit aware"),
        ("shodan-mcp",       "host lookup · search_hosts · count · CDN",    "API key query param",        "no scan; index lookup only"),
        ("rss-bridge",       "direct RSS/Atom feed fetch + ETag cache",     "none (public feeds)",        "feed list curated"),
        ("urlscan",          "search · lookup · submit_scan",               "API-Key header",             "submit gated; search unmetered"),
        ("censys",           "host search · cert search · host lookup",     "Basic (API_ID:API_SECRET)",  "v2 only; index lookup"),
        ("theharvester",     "subdomain / host / IP enum via 45 sources",   "none + per-source keys",     "PASSIVE_SOURCE_ALLOWLIST"),
        ("spiderfoot",       "passive_scan against self-hosted SF daemon",  "optional HTTP Basic",        "PASSIVE_MODULE_ALLOWLIST"),
    ]
    add_table(
        s, rows,
        left=0.5, top=1.5,
        col_widths=(2.0, 4.6, 2.8, 3.13),
        row_h=0.5,
        header=("MCP", "What it surfaces", "Auth model", "Policy enforcement"),
        body_size=10,
    )

    # Footer
    fb = s.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12.33), Inches(0.8))
    ftf = fb.text_frame
    ftf.word_wrap = True
    fp = ftf.paragraphs[0]
    fr = fp.add_run()
    fr.text = (
        "Every MCP follows the same pattern: pydantic-settings config, httpx client, FastMCP server, "
        "respx-mocked unit tests, live-validation tests behind credentials. "
        "Two MCPs (theHarvester, SpiderFoot) carry explicit module-level allowlists because they wrap "
        "tools that could otherwise perform active recon."
    )
    fr.font.size = Pt(12)
    fr.font.italic = True
    fr.font.color.rgb = C_SUBTLE
    fr.font.name = "Calibri"

    add_notes(
        s,
        "Eight wrappers. Six just index-lookup or fetch passive data — those don't need allowlists, they need rate-limit and auth discipline. "
        "Two (theharvester, spiderfoot) wrap tools that CAN perform active recon — those carry mandatory allowlists at the wrapper level. "
        "Rebuild: start with the index-lookup wrappers (VT, Shodan, Censys, urlscan, RSS). Add the active-recon-capable ones (theHarvester, SpiderFoot) only when you've validated the allowlist pattern. Splunk-query is special — that's first-party, not OSINT, treated differently per Hard Rule 8."
    )
    add_footer(s)
    add_pagenum(s, 18, TOTAL)

    # =================================================================
    # Slide 19 — Scheduled brief pipeline
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 4, "Pipelines")
    add_title(s, "Scheduled brief pipeline",
              subtitle="08:00 morning, 16:00 afternoon. Five phases. ~25-40 min wall-clock end-to-end.")

    # Sequence diagram — vertical phases with timing
    phases = [
        ("07:30 EDT", "pre-brief collection", "collector subagent · WebFetch + RSS + MCP tools · writes to threats/raw-signal/ · returns count + source-health deltas"),
        ("~07:55 EDT", "grading (clustering + Admiralty)", "grader subagent · reads all un-promoted raw signal · clusters by topic/actor/vuln · applies INTEL-GRADING · writes to threats/findings/"),
        ("~08:00 EDT", "red-team review (conditional)", "red-team-analyst · ONLY if any finding has WEP ≥ very_likely · reads 1 finding at a time · writes red_team_review section"),
        ("~08:05 EDT", "brief composition", "briefer · reads findings + coverage-log + watch-config · composes Layer 1 + Layer 2 · 13-item preflight · regenerates failed sections"),
        ("~08:08 EDT", "ship", "librarian · LEGAL-POLICY scan · extract Layer 2 to temp file · discord-post.sh #intel-briefs · splunk-log.sh brief_published · git commit + push"),
    ]

    phase_h = Inches(1.0)
    phase_top = Inches(1.55)
    gap = Inches(0.05)
    for i, (time, name, body) in enumerate(phases):
        y = phase_top + (phase_h + gap) * i

        # Time column (left)
        tb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(1.5), phase_h)
        tb.line.fill.background()
        tb.fill.solid()
        tb.fill.fore_color.rgb = C_TITLE
        ttf = tb.text_frame
        ttf.margin_top = Inches(0.3)
        tp = ttf.paragraphs[0]
        tp.alignment = PP_ALIGN.CENTER
        tr = tp.add_run()
        tr.text = time
        tr.font.size = Pt(12)
        tr.font.bold = True
        tr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tr.font.name = "Calibri"

        # Body column (right)
        bb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.0), y, Inches(10.83), phase_h)
        bb.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        bb.line.width = Pt(0.5)
        bb.fill.solid()
        bb.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        btf = bb.text_frame
        btf.margin_top = Inches(0.12)
        btf.margin_left = Inches(0.2)
        btf.word_wrap = True
        bp = btf.paragraphs[0]
        br = bp.add_run()
        br.text = name.upper()
        br.font.size = Pt(13)
        br.font.bold = True
        br.font.color.rgb = C_ACCENT
        br.font.name = "Calibri"

        bp2 = btf.add_paragraph()
        bp2.space_before = Pt(3)
        br2 = bp2.add_run()
        br2.text = body
        br2.font.size = Pt(10)
        br2.font.color.rgb = C_BODY
        br2.font.name = "Calibri"

    add_notes(
        s,
        "Five phases for a scheduled brief. The 07:30 phase is split out so the agent has time to collect signal before composing — collector runs ahead of brief, grader+briefer run on brief time. "
        "Red-team-analyst is conditional: only invoked if any finding has WEP ≥ very_likely. Most briefs don't hit that threshold (5 of 5 actor scorings did NOT escalate to red-team because Intent capped below the trigger). "
        "Wall-clock end-to-end is ~25-40 min. Failure handling per phase: collector retry → degrade with caveat; grader hard-fail → halt pipeline; red-team hard-fail → ship brief with red_team_review=NOT_PERFORMED flag; briefer fail → BRIEF GENERATION FAILED message to Discord; librarian fail → commit but no Discord post."
    )
    add_footer(s)
    add_pagenum(s, 19, TOTAL)

    # =================================================================
    # Slide 20 — FLASH pipeline
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 4, "Pipelines")
    add_title(s, "FLASH pipeline — async, faster",
              subtitle="Every 6h. Same subagents, narrower scope. Quiet hours + critical override.")

    # Three columns: triggers / quiet hours / 72h clock
    columns = [
        ("7 TRIGGERS", C_ACCENT,
         "1.  Critical CVE actively exploited\n"
         "2.  New attribution to a tracked actor\n"
         "3.  Major breach in scope\n"
         "4.  Tracked actor TTP change\n"
         "5.  Sector-specific zero-day\n"
         "6.  Zero-day with no patch\n"
         "7.  Supply-chain compromise"),
        ("QUIET HOURS LOGIC", C_TITLE,
         "Active hours: 09:00 – 21:00 EDT\n\n"
         "Inside active hours → POST\n"
         "Outside + critical override → POST\n"
         "Outside + no override → QUEUE\n\n"
         "Critical override = ALL of:\n"
         "  · CVSS 10\n"
         "  · Confirmed active exploitation\n"
         "  · Tracked actor named\n"
         "  · A&D watchlist hit"),
        ("72h AUTO-DOWNGRADE", C_REFUSAL,
         "Single-source FLASHes ship with a\n"
         "72h clock. If no independent A/B-grade\n"
         "corroboration arrives within window:\n\n"
         "  attribution leg drops A2 → C3\n"
         "  campaign forensics hold\n"
         "  retraction additive in coverage log\n\n"
         "Live-tested 2026-05-09 on the\n"
         "MuddyWater / Rapid7 FLASH-0002.\n"
         "Fired clean."),
    ]
    col_w = Inches(4.1)
    col_h = Inches(5.3)
    col_top = Inches(1.5)
    col_gap = Inches(0.1)
    total_w = col_w * 3 + col_gap * 2
    start_x = (Inches(13.333) - total_w) / 2

    for i, (title, color, body) in enumerate(columns):
        x = start_x + (col_w + col_gap) * i
        # header
        hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, col_top, col_w, Inches(0.5))
        hdr.line.fill.background()
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        htf = hdr.text_frame
        htf.margin_top = Inches(0.08)
        hp = htf.paragraphs[0]
        hp.alignment = PP_ALIGN.CENTER
        hr = hp.add_run()
        hr.text = title
        hr.font.size = Pt(13)
        hr.font.bold = True
        hr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        hr.font.name = "Calibri"

        # body
        bb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, col_top + Inches(0.55), col_w, Inches(col_h.inches - 0.55))
        bb.line.color.rgb = color
        bb.line.width = Pt(1.5)
        bb.fill.solid()
        bb.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        btf = bb.text_frame
        btf.margin_top = Inches(0.2)
        btf.margin_left = Inches(0.2)
        btf.word_wrap = True
        bp = btf.paragraphs[0]
        br = bp.add_run()
        br.text = body
        br.font.size = Pt(11)
        br.font.color.rgb = C_BODY
        br.font.name = "Consolas"

    add_notes(
        s,
        "FLASH is the answer to 'what if something breaks between scheduled briefs?' Same subagent chain, but narrower scope: collector runs only against FLASH-eligible triggers, grader fast-paths single findings, red-team conditional, briefer composes FLASH format (shorter than scheduled brief), librarian posts to #flash-alerts. "
        "Quiet hours protect operators from being woken at 3 AM unless the threat genuinely warrants it. The critical-override threshold has fired exactly once in production — that's the calibration target. False FLASHes erode operator trust faster than missed FLASHes. "
        "72h auto-downgrade is how we handle single-source attribution responsibly: ship the FLASH because it's load-bearing, but commit to retracting/downgrading if independent corroboration doesn't arrive. Test pass on MuddyWater 2026-05-09 was the validation."
    )
    add_footer(s)
    add_pagenum(s, 20, TOTAL)

    # =================================================================
    # Slide 21 — On-demand commands
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 4, "Pipelines")
    add_title(s, "On-demand commands — Discord bridge",
              subtitle="scripts/discord_listener.py listens for slash commands; routes to orchestrator.")

    rows = [
        ("/investigate <target>",  "Deep dive on actor / domain / hash / CVE / campaign", "collector → grader → analyst → briefer"),
        ("/ioc-hunt <indicator>",  "Check IOC against repo + Splunk + external sources", "collector → librarian (Splunk search)"),
        ("/cve <cve-id>",          "Vuln research deep-dive against NVD + KEV + actors",  "vuln-tracker → briefer"),
        ("/new-actor <name>",      "Scaffold new actor dossier from scratch",              "actor-profiler → librarian"),
        ("/update-tracking",       "Refresh actor whose 90-day review is nearest due",     "actor-profiler → librarian"),
        ("/approve-scoring <id>",  "Operator sign-off on a HIGH actor threat-box scoring", "actor-profiler → librarian"),
        ("/brief <type>",          "Manually trigger a scheduled-type brief",              "full scheduled-brief pipeline"),
        ("/flash",                 "Manually run a FLASH sweep",                            "full FLASH pipeline"),
        ("/help",                  "List commands",                                         "discord_listener.py local"),
    ]
    add_table(
        s, rows,
        left=0.5, top=1.5,
        col_widths=(2.5, 5.8, 4.03),
        row_h=0.5,
        header=("Command", "Purpose", "Pipeline invoked"),
        body_size=10,
    )

    # Architecture note
    nb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.2), Inches(12.33), Inches(0.7))
    nb.line.color.rgb = C_SUBTLE
    nb.line.width = Pt(1)
    nb.fill.solid()
    nb.fill.fore_color.rgb = C_BG_LIGHT
    ntf = nb.text_frame
    ntf.margin_top = Inches(0.1)
    ntf.margin_left = Inches(0.2)
    ntf.margin_right = Inches(0.2)
    ntf.word_wrap = True
    np_ = ntf.paragraphs[0]
    nr = np_.add_run()
    nr.text = (
        "discord_listener.py runs as a Windows Task Scheduler entry — polls Discord for messages from "
        "DISCORD_OPERATOR_USER_ID only · validates the slash command · spawns `claude -p` with the "
        "command + args as the prompt · streams the response back to the channel · logs the exchange."
    )
    nr.font.size = Pt(11)
    nr.font.color.rgb = C_BODY
    nr.font.italic = True
    nr.font.name = "Calibri"

    add_notes(
        s,
        "On-demand is where the agent becomes interactive. Without the Discord listener, Archimedes is autonomous-only — runs on a schedule, can't be asked questions. The listener adds operator-in-the-loop. "
        "Security note: only DISCORD_OPERATOR_USER_ID can invoke commands. Messages from any other author are ignored. Single-operator design today; could extend to a team via roster YAML if multi-operator support is needed. "
        "Each slash command invokes a subset of the standard pipeline — /investigate is a partial pipeline (collect + grade + analyze the target, compose a Layer-1 brief), /new-actor is actor-profiler-only, /approve-scoring is the human-side of the HIGH gate."
    )
    add_footer(s)
    add_pagenum(s, 21, TOTAL)

    # =================================================================
    # Slide 22 — Layer 1 + Layer 2 brief composition
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 4, "Pipelines")
    add_title(s, "Layer 1 + Layer 2 — one file, two audiences",
              subtitle="Briefer writes both sections in one .md file. Librarian extracts Layer 2 to post to Discord.")

    code = """threats/briefs/2026-05-11-morning.md  (Layer 1 — the canonical record)

---
brief_id: 2026-05-11-morning
published_at: 2026-05-11T08:00:00-04:00
authored_by: archimedes-briefer
findings_referenced: [finding-2026-05-11-0001, finding-2026-05-11-0002]
word_count: 731
tlp: CLEAR
---

# Morning Brief — 2026-05-11

**[lead-with-impact sentence linked to source]**...
Per Hard Rule 2, the TeamPCP attribution to Checkmarx Jenkins AST is a
restatement of prior reporting...

## Active Threats
**[Headline link](https://...)** ... Digraph: B2 · WEP: likely
(procedural facts only — relational/attribution layer carries lower WEP)
· finding-2026-05-11-0001.

🔗 Connects to: Actor #001 TeamPCP (HIGH, roster) — Jenkins AST plugin
distribution channel is the first appearance of CI/CD-pipeline-poisoning
in the TeamPCP TTP register; dossier update queued via actor-profiler.

[... 700 more words ...]

## 📣 Discord Summary                                   ← Layer 2, last section

Good morning. Here's your 0800 brief — Monday, May 11.

🚨 **Active Threats**

• **[Checkmarx Jenkins AST plugin compromised...](https://www.securityweek.com/...)**
  Checkmarx warned Friday May 9 of a malicious Jenkins AST plugin...
  **DIB CI/CD:** pin to fix version *now*, capture plugin hashes for IOC backfill.

[... ≤1900 chars, natural-language dates, source-linked headlines ...]"""
    add_code(s, code, top=1.55, height=5.5, font_size=9, dark=False)
    add_notes(
        s,
        "Layer 1 is the analyst-grade record — full doctrine framings, Admiralty + WEP citations, Hard Rule 2 annotations, single-source-veto notes, Splunk first-party caveats. ~700 words. Lives in git forever. "
        "Layer 2 is the Discord summary — Smart Brevity, natural-language dates, source-linked headlines, ~250 words, ≤1900 chars. The librarian extracts only Layer 2 (heading-to-EOF read) and posts that to #intel-briefs. "
        "ONE canonical file, TWO renderings, two audiences. The intel team red-team-reviews Layer 1; leadership reads Layer 2 on their phone. "
        "Rebuild lesson: per-audience rendering is a doctrine concern (INTEL-BRIEF-STANDARDS specifies the format) and a librarian responsibility (extraction + delivery), but it's authored by the briefer in one pass. Don't split it across two subagents."
    )
    add_footer(s)
    add_pagenum(s, 22, TOTAL)

    # =================================================================
    # Slide 23 — Retraction handling
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 4, "Pipelines")
    add_title(s, "Retraction handling — additive, never silent",
              subtitle="When something ships wrong, the record of being wrong becomes part of the record.")

    add_bullets(
        s,
        [
            ("Retraction vs. correction",
             "RETRACTION = factually wrong / materially misleading / doctrine-violating. "
             "CORRECTION = typo / wrong link / wording change without meaning shift."),
            ("Retraction is ALWAYS additive",
             "Original brief file stays in threats/briefs/, unmodified. New retraction note APPENDED to the brief. "
             "Coverage-log entry gets retraction block. Discord retraction post in #intel-briefs. Never delete; never edit-to-remove."),
            ("72h auto-downgrade clock",
             "Single-source FLASHes ship with an auto-downgrade clock built into the finding frontmatter. "
             "Conditions for downgrade are SPECIFIED on the finding itself (no second A/B-grade source AND no Splunk hit "
             "AND no CISA pickup → drop attribution leg to C3). When the clock fires unbroken, librarian does in-place "
             "frontmatter downgrade — NOT a full retraction (campaign forensics hold; only the attribution leg moves)."),
            ("Tested live on MuddyWater 2026-05-09",
             "FLASH-0002 (Rapid7 single-source MuddyWater attribution) shipped 2026-05-06 12:18 EDT with a 72h clock. "
             "Clock fired unbroken at 2026-05-09 12:18 EDT. Librarian updated finding frontmatter: digraph A2 → C3, "
             "digraph_split.attribution_leg → C3, campaign forensics held at A2. Coverage log got the supersession entry. "
             "Morning brief 2026-05-10 carried the supersession in narrative."),
            ("Pattern check after every retraction",
             "Was grading process followed? Did red-team catch anything? What would have prevented this? "
             "If 3+ retractions trace to the same source in 90 days, actor-profiler proposes a source-grade downgrade."),
        ],
        body_size=12, sub_size=10,
    )
    add_notes(
        s,
        "Retraction handling is the single most under-built feature in most agent systems. The instinct is 'fix the error and move on.' That destroys auditability. "
        "Archimedes' pattern: every retraction is additive. Every retraction is logged. Every retraction triggers a pattern-check. After 3 retractions traced to the same source, the source itself gets graded down — that's the system learning. "
        "The 72h auto-downgrade clock is the sub-pattern for single-source FLASHes: ship now because it's load-bearing, but write the falsification conditions INTO the finding so the system can self-downgrade if corroboration doesn't arrive. Lived through MuddyWater 2026-05-09 — fired clean."
    )
    add_footer(s)
    add_pagenum(s, 23, TOTAL)

    # =================================================================
    # Slide 24 — Scheduler
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 5, "Infrastructure + ops")
    add_title(s, "Scheduler — Windows Task Scheduler + PowerShell wrapper",
              subtitle="8 scheduled tasks. infrastructure/scheduler/ holds the templates.")

    rows = [
        ("flash-sweep-0000", "00:00 EDT", "FLASH sweep · queue if quiet hours"),
        ("flash-sweep-0600", "06:00 EDT", "FLASH sweep · queue if quiet hours"),
        ("pre-brief-morning", "07:30 EDT", "Collector pre-pass · raw signal → disk · no brief"),
        ("morning-brief", "08:00 EDT", "Full scheduled-brief pipeline · 5 phases"),
        ("flash-sweep-1200", "12:00 EDT", "FLASH sweep"),
        ("pre-brief-afternoon", "15:30 EDT", "Collector pre-pass · raw signal → disk"),
        ("afternoon-brief", "16:00 EDT", "Full scheduled-brief pipeline · 5 phases"),
        ("flash-sweep-1800", "18:00 EDT", "FLASH sweep"),
    ]
    add_table(
        s, rows,
        left=0.5, top=1.5,
        col_widths=(3.0, 1.5, 8.33),
        row_h=0.45,
        header=("Task name", "Time", "What it does"),
        body_size=10,
    )

    code = """# Wrapper pattern — scripts/run_phase.ps1 (excerpt)

$ErrorActionPreference = 'Stop'
$phaseId = "morning-brief-$(Get-Date -Format yyyyMMdd-HHmmss)"

# Splunk: started event
& uv run python scripts/splunk_log.py --event-file $eventFile

# Invoke claude -p with the phase prompt
& claude -p "Execute the $phase pipeline. Mode 1 if scheduled brief."

# Catchup-push (mitigates librarian push variance)
$exitCode = & git push origin main; if ($?) { 0 } else { $LASTEXITCODE }

# Splunk: completed event with catchup_push_exit
& uv run python scripts/splunk_log.py --event-file $completedEvent"""
    add_code(s, code, top=5.4, height=1.7, font_size=10, dark=False)
    add_notes(
        s,
        "8 Windows Task Scheduler entries. XML templates in infrastructure/scheduler/. Each task invokes scripts/run_phase.ps1 with the phase name as an argument. "
        "The wrapper does three things: (1) logs Splunk 'started' event with run_id, (2) invokes `claude -p` with the phase prompt, (3) logs Splunk 'completed' event with exit code + catchup-push status. "
        "The catchup-push pattern (Session 11 discovery) compensates for the librarian's intermittent push behavior — if librarian skipped `git push origin main`, the wrapper does it. catchup_push_exit field in Splunk events tracks the outcome for dashboard visibility."
    )
    add_footer(s)
    add_pagenum(s, 24, TOTAL)

    # =================================================================
    # Slide 25 — Hooks
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 5, "Infrastructure + ops")
    add_title(s, "Hooks — bash scripts for deterministic side effects",
              subtitle=".claude/hooks/  — why hooks not subagents (no LLM cost, deterministic, easy to debug)")

    code = """.claude/hooks/discord-post.sh

#!/usr/bin/env bash
# Post a message file to a Discord channel via webhook.
# Invoked by librarian: bash discord-post.sh --channel intel-briefs --message-file <path>

set -euo pipefail

CHANNEL=""; MESSAGE_FILE=""; FLASH=false
while [[ $# -gt 0 ]]; do
    case "$1" in
        --channel)      CHANNEL="$2"; shift 2 ;;
        --message-file) MESSAGE_FILE="$2"; shift 2 ;;
        --flash)        FLASH=true; shift ;;
    esac
done

# Resolve webhook URL by channel name (from .env)
case "$CHANNEL" in
    intel-briefs)   WEBHOOK="$DISCORD_WEBHOOK_INTEL_BRIEFS" ;;
    flash-alerts)   WEBHOOK="$DISCORD_WEBHOOK_FLASH_ALERTS" ;;
    actor-review)   WEBHOOK="$DISCORD_WEBHOOK_ACTOR_REVIEW" ;;
    *) echo "Unknown channel: $CHANNEL"; exit 1 ;;
esac

# Read message body; truncate at 2000 chars if needed (Discord limit)
BODY=$(cat "$MESSAGE_FILE")
if [[ ${#BODY} -gt 2000 ]]; then
    echo "WARN: body >2000 chars; truncating"
    BODY="${BODY:0:1997}..."
fi

# POST
curl -X POST -H "Content-Type: application/json" \\
     -d "$(jq -Rs '{content: .}' <<< "$BODY")" \\
     "$WEBHOOK"

echo "POSTED to #$CHANNEL"   # parsed by librarian for status"""
    add_code(s, code, top=1.55, height=5.5, font_size=10, dark=True)
    add_notes(
        s,
        "Three hooks total: discord-post.sh, splunk-log.sh, and a couple others for index regen. Why hooks instead of subagents? "
        "(1) No LLM cost — bash scripts run free, deterministic, fast. (2) Deterministic behavior — Discord post either succeeds or fails based on HTTP status, not on model reasoning. (3) Easy to debug — you can run the hook by hand from the command line. "
        "Rebuild lesson: distinguish between 'I need the model to decide' (subagent) and 'I need the model to invoke a thing that always does the same operation' (hook). Hooks handle the second case. Don't make Claude write JSON to a webhook on every brief; have Claude invoke a hook that does that."
    )
    add_footer(s)
    add_pagenum(s, 25, TOTAL)

    # =================================================================
    # Slide 26 — Audit trail
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 5, "Infrastructure + ops")
    add_title(s, "Audit trail — three logs that reconstruct any action",
              subtitle="Together: git + Splunk + Discord = full provenance for everything the agent did.")

    add_bullets(
        s,
        [
            ("git history — the corpus changes",
             "Every brief committed. Every finding committed. Every dossier update. Every doctrine revision. "
             "60+ commits/week typical. Provenance is queryable: `git log --since`, `git blame`, `git log -p <file>`. "
             "Co-Authored-By: Claude trailer on every commit makes agent-vs-human authorship explicit."),
            ("Splunk events — operational telemetry",
             "Sourcetype: archimedes:operation. Events: run_started, run_completed, brief_published, "
             "git_committed, flash_queued, threat_box_scoring_completed, discord_format_preview_posted, "
             "retraction_logged, policy_violation. Indexed for dashboarding + alerting on agent behavior."),
            ("Discord channels — delivery log",
             "#intel-briefs (scheduled briefs + Layer 2 summaries). "
             "#flash-alerts (FLASH posts). "
             "#actor-review (HIGH-scoring proposals + retraction proposals). "
             "#commands (operator slash commands + responses). Every channel post has a message ID; ID recorded in Splunk."),
            ("Cross-reference via run_id",
             "Every scheduled phase has a run_id (e.g., morning-brief-20260511-080000). The run_id appears in the git "
             "commit message, the Splunk event, and the Discord post. Three logs joined on one key."),
            ("Refusal logging — policy-violations.yaml",
             "Every Hard Rule refusal logged to infrastructure/policy-violations.yaml with timestamp, subagent, "
             "section, attempted_action, triggering_prompt (sanitized), reason. Auditable record of what the agent "
             "WAS asked to do but refused."),
        ],
        body_size=12, sub_size=10,
    )
    add_notes(
        s,
        "Auditability is what makes the agent defensible in a regulated environment. Five years from now, if someone asks 'how did the agent arrive at this attribution on 2026-05-06,' the answer is: git log + Splunk events + Discord posts joined on run_id. "
        "That same auditability is also how you debug agent behavior. When the briefer makes a weird choice, you can pull the run_id, find the Splunk events, see what input it received, check the doctrine version in git at that commit. "
        "Rebuild lesson: design for audit BEFORE you have anything to audit. Adding it later is exponentially harder."
    )
    add_footer(s)
    add_pagenum(s, 26, TOTAL)

    # =================================================================
    # Slide 27 — Configuration
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 5, "Infrastructure + ops")
    add_title(s, "Configuration — what's where, what's gitignored",
              subtitle="Three classes: secrets (.env), agent-readable YAML (infrastructure/), runtime state (gitignored YAML).")

    rows = [
        (".env", "gitignored", "API keys · Discord webhooks · Splunk HEC token · ANTHROPIC_API_KEY (deliberately unset; use Claude Max)"),
        (".env.example", "committed", "Schema template · variable names + comments · NO secret values · current with the real .env"),
        ("infrastructure/source-grades.yaml", "committed", "Source → grade + category mapping · 45 entries · ratification log in source-grade-log.md"),
        ("infrastructure/source-health.yaml", "gitignored", "Runtime: status · last_successful_fetch · failure_count · stale_since · last_error · operator notes preserved"),
        ("infrastructure/watch-config.yaml", "committed", "Standing brief sections · sector watchlists · FLASH triggers · silent_day_template"),
        ("infrastructure/authorized-targets.yaml", "committed", "Hard Rule 4 enforcement list · only these targets may be actively scanned"),
        ("infrastructure/watchlists/*.yaml", "committed", "Per-watchlist YAML: A&D, ransomware, etc. Used by collector to filter raw signal"),
        ("infrastructure/flash-queue.yaml", "gitignored", "Quiet-hours FLASH queue · 09:00 catchup sweep processes it"),
        ("threats/threat-actors/_roster.yaml", "committed", "23-actor roster · per-actor primary_name, aliases, attribution, dossier path, scoring counters"),
        ("threats/vulnerabilities/_index.yaml", "committed", "Tracked CVE index"),
        ("threats/iocs/_master-index.yaml", "committed", "Regenerated by librarian on IOC changes · scripts/regenerate_ioc_index.py"),
    ]
    add_table(
        s, rows,
        left=0.5, top=1.5,
        col_widths=(3.8, 1.6, 6.93),
        row_h=0.42,
        header=("File", "Git status", "Contents"),
        body_size=9,
    )
    add_notes(
        s,
        "Configuration discipline: secrets gitignored, schemas committed, runtime state gitignored, operator-set context preserved. "
        "Two notable patterns: (1) .env.example mirrors the actual .env structure — anyone cloning the repo can copy-rename to get a working schema; (2) source-health.yaml has a mix of runtime fields (collector overwrites) and operator-set fields (notes preserved verbatim) — collector subagent doctrine specifies the field ownership boundary. "
        "If your rebuild has separate operator-set context that should survive runtime overwrites, codify which fields the agent owns vs. the operator owns — same pattern."
    )
    add_footer(s)
    add_pagenum(s, 27, TOTAL)

    # =================================================================
    # Slide 28 — Testing
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 5, "Infrastructure + ops")
    add_title(s, "Testing — per-MCP suites + the 13-item preflight",
              subtitle="231 unit tests across 8 MCPs. Live validation patterns where mocks aren't enough.")

    add_bullets(
        s,
        [
            ("Per-MCP pytest suites",
             "Each MCP has its own tests/ directory. test_config.py (env loading), test_<name>_client.py "
             "(respx-mocked HTTP), test_policy.py (where applicable), test_integration.py (live tests behind "
             "credentials, skipped by default). 231 unit tests passing across 8 MCPs."),
            ("respx for HTTP mocking",
             "All HTTP wrappers tested with respx — register expected requests + canned responses, run the client, "
             "assert. No real network in unit tests."),
            ("Live validation for third-party tools",
             "Mocks based on documentation are confidence-without-verification. theHarvester 4.10.1 + SpiderFoot "
             "4.0.0 each surfaced 3 bugs that unit tests couldn't have caught: API shape drift, off-by-one in "
             "scanstatus parsing, wrong JSON endpoint. Fix: live-test against real services."),
            ("Preflight checklist — 13 items",
             ".claude/skills/smart-brevity/references/preflight-checklist.md. Briefer runs this checklist before "
             "every brief publication. ALL must pass or regenerate. Items: source URL present, Admiralty digraph, "
             "coverage-log respected, lead-with-impact, banned phrases zero, standing sections present, WEP "
             "vocabulary, word count in range, actor/vuln links, TLP marked, single-source veto, quote discipline, "
             "Layer 2 section compliant."),
            ("Doctrine review as a test phase",
             "When doctrine changes, walk through every subagent that reads it and ask: does the change break "
             "anything? When INTEL-BRIEF-STANDARDS got Layer 2, briefer.md + librarian.md + preflight-checklist.md "
             "all needed coordinated updates."),
        ],
        body_size=12, sub_size=10,
    )
    add_notes(
        s,
        "Testing across three layers: unit tests for MCP wrappers (231 passing), preflight checklist for brief composition (13 items), live validation against real services (catches the 6 real bugs mocks missed). "
        "Doctrine changes need their own test phase — when you change a .md file, walk through every subagent that reads it. This is similar to schema changes in a database: a schema change isn't just a SQL migration, it's a coordinated update across all the code that reads the schema."
    )
    add_footer(s)
    add_pagenum(s, 28, TOTAL)

    # =================================================================
    # Slide 29 — Minimum viable Archimedes
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 6, "Rebuild guide")
    add_title(s, "Minimum viable Archimedes",
              subtitle="The smallest version that delivers value. Then iterate.")

    tree = """minimum-viable/
├── doctrine/
│   └── INTEL-GRADING.md            # ONE doctrine file. Admiralty + WEP.
│
├── .claude/
│   ├── agents/
│   │   ├── grader.md               # ONE subagent. Promotes raw → graded.
│   │   └── briefer.md              # ONE subagent. Composes the brief.
│   └── hooks/
│       └── discord-post.sh         # ONE delivery channel.
│
├── mcps/
│   └── rss-bridge/                 # ONE MCP. RSS feed fetcher.
│
├── threats/
│   ├── raw-signal/                 # collector output (manual to start)
│   ├── findings/                   # grader output
│   └── briefs/                     # briefer output
│
├── .env                            # DISCORD_WEBHOOK_URL only
├── pyproject.toml
└── CLAUDE.md                       # the agent's identity"""
    add_code(s, tree, top=1.55, height=4.5, font_size=10, dark=False)

    bullets_box = s.shapes.add_textbox(Inches(0.5), Inches(6.1), Inches(12.33), Inches(1.0))
    btf = bullets_box.text_frame
    btf.word_wrap = True
    bp = btf.paragraphs[0]
    br = bp.add_run()
    br.text = (
        "Start here. Run it manually for a week — operator pulls articles from feeds, drops them into raw-signal/, invokes the grader, "
        "invokes the briefer, posts to Discord via the hook. As friction surfaces, add the next piece: "
        "collector subagent → more MCPs → librarian subagent → scheduler → more doctrine files. "
        "Each addition replaces one source of manual work. The minimum viable version proves the value before you commit to building the rest."
    )
    br.font.size = Pt(12)
    br.font.italic = True
    br.font.color.rgb = C_BODY
    br.font.name = "Calibri"

    add_notes(
        s,
        "The trap is building Archimedes top-down — full doctrine, all 9 subagents, all 8 MCPs, scheduler, hooks, audit trail — before any of it has produced a single brief. That'll fail. "
        "The MVP version: one doctrine file, two subagents, one MCP, one delivery channel. Run it manually for a week. The friction tells you what to build next. "
        "This pattern (build the load-bearing thin slice; add to it from real friction) is also what Archimedes itself did — Session 1 had ONE doctrine file (LEGAL-POLICY) and one subagent (the grader was first). Everything else accreted over 13 sessions."
    )
    add_footer(s)
    add_pagenum(s, 29, TOTAL)

    # =================================================================
    # Slide 30 — Build order
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 6, "Rebuild guide")
    add_title(s, "Recommended build order",
              subtitle="What to add when. Time estimates assume the team is familiar with Claude + MCP.")

    rows = [
        ("Week 1",  "Doctrine + identity",       "Write CLAUDE.md + 1-2 doctrine files. No code yet. Just the agent's job description."),
        ("Week 1",  "Grader subagent + RSS MCP", "Build the grading path end-to-end against canned input. Run manually."),
        ("Week 2",  "Briefer subagent",          "Add brief composition. Hand-craft 3-5 sample briefs first; reverse-engineer the prompt from the samples."),
        ("Week 2",  "Discord delivery (hook)",   "discord-post.sh + webhook URL in .env. Ship the first brief to a channel."),
        ("Week 3",  "Librarian subagent + git",  "Move side effects (commit, post, log) into the librarian. Test the LEGAL-POLICY content scan gate."),
        ("Week 3",  "Collector subagent",        "Replace manual feed-pulling with the collector. Add 5+ MCPs (VT, Shodan, urlscan, RSS, web search)."),
        ("Week 4",  "Scheduler",                 "Task Scheduler XMLs · run_phase.ps1 wrapper · Splunk telemetry. First fully autonomous brief."),
        ("Week 4",  "Analyst + red-team",        "Add SAT skills + the analyst chain. Trigger red-team conditional on WEP threshold."),
        ("Week 5",  "Actor profiles + vulns",    "Add actor-profiler + vuln-tracker. Hard Rule 5 gate to #actor-review."),
        ("Week 5",  "Retraction handling",       "RETRACTION-POLICY + 72h auto-downgrade clock. Test on a hypothetical bad finding."),
        ("Week 6",  "On-demand / Discord listener", "scripts/discord_listener.py + the slash command surface. Now interactive."),
        ("Week 6",  "Hardening + observability", "Splunk dashboards. Refusal logging. Audit-trail end-to-end tests. Production cutover."),
    ]
    add_table(
        s, rows,
        left=0.5, top=1.5,
        col_widths=(1.2, 3.5, 7.63),
        row_h=0.42,
        header=("When", "Add", "Why now"),
        body_size=10,
    )

    add_notes(
        s,
        "Six weeks for a working version. Each week is two additions. The order matters: doctrine + identity before any code; grader + briefer before collector (so you can test on canned data first); manual operation in weeks 1-3, autonomous in week 4 once scheduler is in. "
        "The order I listed is the order Archimedes evolved through Sessions 1-7 — there's a reason; do it the way it grew."
    )
    add_footer(s)
    add_pagenum(s, 30, TOTAL)

    # =================================================================
    # Slide 31 — What to skip / defer
    # =================================================================
    s = blank_slide(prs)
    add_section_band(s, 6, "Rebuild guide")
    add_title(s, "What to skip / defer",
              subtitle="Things Archimedes has that you probably don't need on day 1.")

    add_bullets(
        s,
        [
            ("The web dashboard",
             "interfaces/dashboard/ is an empty placeholder. Originally planned as a Flask UI for browsing the corpus. "
             "Deferred indefinitely because the Discord listener (/cve, /investigate, /ioc-hunt) covers the interactive-query need "
             "and the markdown corpus is browsable in GitHub or any editor. Skip unless you need a non-technical-stakeholder browser."),
            ("8 MCP wrappers on day 1",
             "Start with 1-3. RSS bridge covers most feeds. VirusTotal covers IOC enrichment. Splunk-query covers first-party. "
             "Add more when you hit a 'this would be easier with X' moment. theHarvester + SpiderFoot can wait until you need passive recon."),
            ("All 9 subagents",
             "Start with 2-3. Grader + briefer covers the load-bearing path. Add collector (week 3), librarian (week 3), "
             "scheduler (week 4), analyst (week 4), red-team (week 4), actor-profiler (week 5), vuln-tracker (week 5). "
             "Building all 9 before anything ships is the failure mode."),
            ("Per-actor red-team variant",
             "We have red-team-analyst for findings. We also drafted a per-actor red-team for scoring runs — never built "
             "because the regular red-team chain was sufficient. Skip until proven necessary."),
            ("On-demand investigation slash commands beyond /investigate + /ioc-hunt",
             "The full list (8 slash commands) accreted over time. Start with 2-3 (/investigate + /ioc-hunt + /help). "
             "/new-actor, /update-tracking, /approve-scoring matter once you have an actor corpus to maintain."),
            ("FLASH if your cadence permits",
             "Scheduled briefs cover most operational needs. FLASH is the answer to 'CVE-10 broke at 3 AM, page me now.' "
             "If your operators don't need that pager-style escalation, scheduled briefs alone are sufficient."),
        ],
        body_size=12, sub_size=10,
    )
    add_notes(
        s,
        "Half the system isn't needed for V1. The dashboard is the most expensive thing we never built — Flask + auth + per-team filters + drilldown into CMDB would be 2-3 weeks of build, and the Discord listener + markdown corpus together cover the interactive-query need for one operator. Defer it. "
        "FLASH is the most expensive thing we DID build that might not justify the cost for every team. If your operators don't need a pager-style escalation channel, scheduled briefs cover the cadence. "
        "Take what makes the agent shape work for your team; defer what doesn't until you feel the pain."
    )
    add_footer(s)
    add_pagenum(s, 31, TOTAL)

    # =================================================================
    # Slide 32 — Q&A
    # =================================================================
    s = blank_slide(prs)
    box = s.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(12.33), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Questions?"
    r.font.size = Pt(54)
    r.font.bold = True
    r.font.color.rgb = C_TITLE
    r.font.name = "Calibri"

    box2 = s.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.33), Inches(0.5))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Where would you start?"
    r2.font.size = Pt(20)
    r2.font.italic = True
    r2.font.color.rgb = C_SUBTLE
    r2.font.name = "Calibri"

    box3 = s.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.33), Inches(0.5))
    tf3 = box3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = "[Presenter name]   ·   [Email / Slack / Teams handle]"
    r3.font.size = Pt(14)
    r3.font.color.rgb = C_SUBTLE
    r3.font.name = "Calibri"

    add_notes(
        s,
        "Likely questions to be ready for: "
        "(1) Token cost at this volume — order of magnitude per month, not exact. "
        "(2) Why Claude over Anthropic API direct — Claude Code's subagent runtime is the value-add; building the same with bare API would be 3-4 months of agent-runtime work before the CTI work starts. "
        "(3) What about open-source LLMs — possible but harder; the long-context + tool-use + structured-output discipline is where Claude is strongest. "
        "(4) How long to a rebuild — 6 weeks for a working version with the team familiar; 12 weeks if learning Claude as you go. "
        "(5) What broke that you wish you'd known earlier — schtasks UTF-16 LE BOM, PYTHONHOME poisoning under uv run, librarian push variance. All documented in CLAUDE.md operational notes."
    )
    add_footer(s)
    add_pagenum(s, 32, TOTAL)

    # ---------- save ----------
    out_path = Path(__file__).parent / "archimedes-architecture.pptx"
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    path = build()
    print(f"Wrote: {path}")
    print(f"Size:  {path.stat().st_size:,} bytes")
